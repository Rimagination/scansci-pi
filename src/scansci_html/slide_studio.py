"""Source-first PDF and document to editable PPTX production for ScanSci."""

from __future__ import annotations

import base64
import binascii
from datetime import datetime
import json
import re
import shutil
from pathlib import Path
from typing import Any, Callable, Protocol
from uuid import uuid4
import zipfile

from bs4 import BeautifulSoup
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from .builtin_skills import SCIENTIFIC_SLIDES_CONTRACT
from .easyslides_classic import render_classic_deck
from .easyslides_plugin import render_semantic_deck
from .slide_plan import build_slide_plan, safe_presentation_name, write_slide_plan
from .slides_templates import get_slide_template


class _ChatJsonClient(Protocol):
    def complete_json(self, messages: list[dict[str, str]], *, schema_name: str) -> Any:
        ...


_SUPPORTED_SUFFIXES = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".txt": "文本",
    ".md": "Markdown",
    ".markdown": "Markdown",
    ".html": "HTML",
    ".htm": "HTML",
}
_MAX_SOURCE_FILES = 3
_MAX_SOURCE_BYTES = 30 * 1024 * 1024
_MAX_EXTRACTED_CHARS = 90_000
_MAX_MODEL_CONTEXT_CHARS = 24_000
_MAX_PDF_PAGES = 80
_SAFE_NAME = re.compile(r"[^\w.-]+", flags=re.UNICODE)


def supported_slide_source_suffixes() -> tuple[str, ...]:
    return tuple(_SUPPORTED_SUFFIXES)


def persist_slide_sources(workspace: str | Path, values: object) -> list[dict[str, Any]]:
    """Copy user-selected presentation sources into a private durable area.

    Browser uploads may provide a data URL; desktop selections provide a path.
    The persisted record never exposes the original path back to the web UI.
    """

    sources = list(values or []) if isinstance(values, list) else []
    if not sources:
        raise ValueError("请先添加至少一个 PDF、Word、Markdown、TXT 或 HTML 文件")
    if len(sources) > _MAX_SOURCE_FILES:
        raise ValueError(f"一次最多可制作 {_MAX_SOURCE_FILES} 份材料的合并幻灯片")

    root = Path(workspace).resolve().parent / ".scansci-slide-sources"
    root.mkdir(parents=True, exist_ok=True)
    records: list[dict[str, Any]] = []
    total_bytes = 0
    for raw in sources:
        if not isinstance(raw, dict):
            raise ValueError("幻灯片来源格式无效")
        name = _safe_source_name(str(raw.get("name", "")))
        source_path = str(raw.get("path", "")).strip()
        data_url = str(raw.get("data_url", "")).strip()
        if source_path:
            candidate = Path(source_path).expanduser()
            try:
                candidate = candidate.resolve(strict=True)
            except OSError as error:
                raise ValueError("选择的演示材料已不存在") from error
            if not candidate.is_file():
                raise ValueError("选择的演示材料不是文件")
            suffix = candidate.suffix.casefold()
            name = name or _safe_source_name(candidate.name)
            payload_size = candidate.stat().st_size
            target = root / f"source-{uuid4().hex}{suffix}"
            _validate_source(name, suffix, payload_size, total_bytes)
            shutil.copy2(candidate, target)
        elif data_url:
            suffix = Path(name).suffix.casefold()
            if not suffix:
                raise ValueError("上传文件缺少可识别的扩展名")
            payload = _decode_data_url(data_url)
            payload_size = len(payload)
            _validate_source(name, suffix, payload_size, total_bytes)
            target = root / f"source-{uuid4().hex}{suffix}"
            target.write_bytes(payload)
        else:
            raise ValueError("请从本机选择演示材料")

        total_bytes += payload_size
        records.append(
            {
                "id": target.stem,
                "name": name,
                "path": str(target),
                "kind": _SUPPORTED_SUFFIXES[suffix],
                "size": payload_size,
            }
        )
    return records


def create_source_slide_deck(
    *,
    workspace: str | Path,
    sources: list[dict[str, Any]],
    topic: str = "",
    template_id: str = "",
    chat_client: _ChatJsonClient | None = None,
    slides_root: str | Path | None = None,
    on_progress: Callable[[float, str], None] | None = None,
) -> dict[str, Any]:
    """Extract supplied material, plan a template-led deck, and export PPTX."""

    _report_progress(on_progress, 0.06, "正在解析演示材料")
    extracted = [_extract_source(source) for source in sources]
    template = get_slide_template(template_id, slides_root) if str(template_id).strip() else None
    theme = _template_theme(template)
    _report_progress(on_progress, 0.20, "正在提炼论点与证据结构")
    outline, planning = _plan_deck(extracted, topic=topic, chat_client=chat_client, template=template)
    deck_path = _presentation_path(workspace, str(outline["title"]))

    if template and str(template.get("generation_mode", "")) == "easyslides-semantic":
        native = render_semantic_deck(
            workspace=workspace,
            deck_path=deck_path,
            outline=outline,
            sources=extracted,
            template=template,
            root=slides_root,
            on_progress=on_progress,
        )
        _report_progress(on_progress, 1.0, "EasySlides 原生成品已完成")
        return {
            "ok": True,
            "message": "已用 EasySlides 生成并质检可编辑 PPTX",
            "file_path": str(deck_path),
            "pptx_path": str(deck_path),
            "download_name": deck_path.name,
            "outline": outline,
            "slide_plan": native["plan"],
            "slide_plan_path": native["plan_path"],
            "execution_lock_path": native["execution_lock_path"],
            "project_path": native["project_path"],
            "renderer": "easyslides-semantic-native",
            "renderer_label": f"EasySlides {native['plugin'].get('version', '')} · 原生可编辑".strip(),
            "render_mode": "native",
            "quality_gate": native["quality_gate"],
            "academic_qa": native.get("academic_qa", {}),
            "visual_measure": native.get("visual_measure", {}),
            "visual_review": native.get("visual_review", {}),
            "execution_lock_validation": native.get("execution_lock_validation", {}),
            "easyslides": native["plugin"],
            "render_manifest": native["render_manifest"],
            "slide_count": len(list(native["plan"].get("slides", []) or [])),
            "source_count": len(extracted),
            "sources": _source_summaries(extracted),
            "template": template,
            "template_id": str(template.get("id", "")),
            "planning": planning,
        }

    if template and str(template.get("generation_mode", "")) == "easyslides-classic":
        native = render_classic_deck(
            workspace=workspace,
            deck_path=deck_path,
            outline=outline,
            sources=extracted,
            template=template,
            root=slides_root,
            on_progress=on_progress,
        )
        _report_progress(on_progress, 1.0, "EasySlides 正式模板成品已完成")
        return {
            "ok": True,
            "message": "已用 EasySlides 正式模板生成并质检可编辑 PPTX",
            "file_path": str(deck_path),
            "pptx_path": str(deck_path),
            "download_name": deck_path.name,
            "outline": outline,
            "slide_plan": native["plan"],
            "slide_plan_path": native["plan_path"],
            "project_path": native["project_path"],
            "renderer": "easyslides-classic-native",
            "renderer_label": f"EasySlides {native['plugin'].get('version', '')} · 正式模板原生可编辑".strip(),
            "render_mode": "native",
            "quality_gate": native["quality_gate"],
            "classic_text_fit": native.get("classic_text_fit", {}),
            "visual_measure": native.get("visual_measure", {}),
            "visual_review": native.get("visual_review", {}),
            "easyslides": native["plugin"],
            "render_manifest": native["render_manifest"],
            "slide_count": len(list(native["plan"].get("slides", []) or [])),
            "source_count": len(extracted),
            "sources": _source_summaries(extracted),
            "template": template,
            "template_id": str(template.get("id", "")),
            "planning": planning,
        }

    _report_progress(on_progress, 0.54, "正在使用 ScanSci 兼容排版器生成 PPTX")
    slide_plan = build_slide_plan(outline, extracted)
    if template:
        slide_plan["template"] = template
        slide_plan.setdefault("theme", {})["accent"] = str(template.get("primary_color") or "#14897D")
    slide_plan_path = deck_path.with_name(f"{deck_path.stem}.slide-plan.json")
    write_slide_plan(slide_plan_path, slide_plan)
    _render_pptx(deck_path, outline, extracted, theme=theme)
    _report_progress(on_progress, 1.0, "兼容版 PPTX 已完成")
    return {
        "ok": True,
        "message": "已从上传材料生成可编辑 PPTX",
        "file_path": str(deck_path),
        "pptx_path": str(deck_path),
        "download_name": deck_path.name,
        "outline": outline,
        "slide_plan": slide_plan,
        "slide_plan_path": str(slide_plan_path),
        "renderer": "scansci-python-pptx-compat",
        "renderer_label": "ScanSci 兼容排版器",
        "render_mode": "compatibility",
        "quality_gate": {"status": "not_run", "reason": "该模板尚未接入 EasySlides 原生生产门禁"},
        "enhanced_renderer": "pptxgenjs-browser",
        "slide_count": len(list(outline.get("slides", []) or [])) + 1,
        "source_count": len(extracted),
        "sources": _source_summaries(extracted),
        "template": template,
        "template_id": str(template.get("id", "")) if template else "",
        "planning": planning,
    }


def _source_summaries(extracted: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return [
        {
            "name": item["name"],
            "kind": item["kind"],
            "page_count": item["page_count"],
            "warnings": item["warnings"],
        }
        for item in extracted
    ]


def _report_progress(callback: Callable[[float, str], None] | None, fraction: float, summary: str) -> None:
    if callback is not None:
        callback(max(0.0, min(1.0, float(fraction))), summary)


def save_browser_rendered_deck(
    workspace: str | Path,
    *,
    file_name: str,
    base64_data: str,
) -> dict[str, Any]:
    """Validate and save a browser-rendered PptxGenJS presentation."""

    encoded = str(base64_data or "").strip()
    if encoded.startswith("data:"):
        encoded = encoded.partition(",")[2]
    try:
        payload = base64.b64decode(encoded, validate=True)
    except (ValueError, binascii.Error) as error:
        raise ValueError("PptxGenJS 导出数据无效") from error
    if not payload or len(payload) > 80 * 1024 * 1024:
        raise ValueError("PptxGenJS 导出文件为空或过大")

    try:
        from io import BytesIO

        with zipfile.ZipFile(BytesIO(payload)) as archive:
            names = set(archive.namelist())
            if "[Content_Types].xml" not in names or "ppt/presentation.xml" not in names:
                raise ValueError("导出内容不是有效的 PPTX")
    except zipfile.BadZipFile as error:
        raise ValueError("导出内容不是有效的 PPTX") from error

    root = Path(workspace).resolve().parent / "presentations"
    root.mkdir(parents=True, exist_ok=True)
    requested = Path(str(file_name or "")).stem
    stem = safe_presentation_name(requested, fallback="scansci_slides_enhanced")
    target = root / f"{stem}.pptx"
    counter = 2
    while target.exists():
        target = root / f"{stem}-{counter}.pptx"
        counter += 1
    temporary = target.with_name(f".{target.name}.{uuid4().hex}.tmp")
    try:
        temporary.write_bytes(payload)
        temporary.replace(target)
    finally:
        temporary.unlink(missing_ok=True)
    return {
        "ok": True,
        "file_path": str(target),
        "download_name": target.name,
        "download_url": f"/api/presentations/{target.name}",
        "renderer": "pptxgenjs-browser",
        "bytes": len(payload),
    }


def _validate_source(name: str, suffix: str, size: int, total_bytes: int) -> None:
    if suffix not in _SUPPORTED_SUFFIXES:
        types = "、".join(item.removeprefix(".").upper() for item in _SUPPORTED_SUFFIXES)
        raise ValueError(f"暂只支持 {types} 作为幻灯片来源")
    if not name:
        raise ValueError("上传文件缺少名称")
    if size <= 0:
        raise ValueError("上传文件为空")
    if size > _MAX_SOURCE_BYTES or total_bytes + size > _MAX_SOURCE_BYTES:
        raise ValueError("制作幻灯片的源文件总大小不能超过 30 MB")


def _decode_data_url(value: str) -> bytes:
    header, separator, encoded = value.partition(",")
    if not separator or ";base64" not in header.casefold():
        raise ValueError("上传数据格式无效")
    try:
        return base64.b64decode(encoded, validate=True)
    except ValueError as error:
        raise ValueError("上传文件无法解码") from error


def _safe_source_name(value: str) -> str:
    name = Path(value).name.strip()
    return name[:180]


def _extract_source(record: dict[str, Any]) -> dict[str, Any]:
    source_path = Path(str(record["path"])).resolve()
    suffix = source_path.suffix.casefold()
    if suffix not in _SUPPORTED_SUFFIXES or not source_path.is_file():
        raise ValueError("幻灯片来源已失效，请重新添加")
    warnings: list[str] = []
    page_count = 1
    if suffix == ".pdf":
        reader = PdfReader(str(source_path))
        if reader.is_encrypted:
            try:
                decrypted = reader.decrypt("")
            except Exception as error:  # noqa: BLE001 - third-party parser boundary
                raise ValueError("受密码保护的 PDF 暂不能直接制作幻灯片") from error
            if not decrypted:
                raise ValueError("受密码保护的 PDF 暂不能直接制作幻灯片")
        page_count = len(reader.pages)
        pieces: list[str] = []
        for index, page in enumerate(reader.pages[:_MAX_PDF_PAGES]):
            text = (page.extract_text() or "").strip()
            if text:
                pieces.append(f"[第 {index + 1} 页]\n{text}")
            if sum(len(item) for item in pieces) >= _MAX_EXTRACTED_CHARS:
                warnings.append("材料较长，已截取前部文本用于生成第一版幻灯片")
                break
        if page_count > _MAX_PDF_PAGES:
            warnings.append(f"PDF 共 {page_count} 页，已先解析前 {_MAX_PDF_PAGES} 页")
        text = "\n\n".join(pieces)
    elif suffix == ".docx":
        document = Document(str(source_path))
        text = "\n".join(paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip())
    elif suffix in {".html", ".htm"}:
        text = BeautifulSoup(source_path.read_text(encoding="utf-8", errors="replace"), "html.parser").get_text("\n", strip=True)
    else:
        text = source_path.read_text(encoding="utf-8", errors="replace")
    text = _compact_text(text, limit=_MAX_EXTRACTED_CHARS)
    if not text:
        raise ValueError(f"未能从「{record.get('name', source_path.name)}」提取可用文本；扫描件请先在文档处理设置中启用 OCR 后重试")
    return {
        "id": str(record.get("id", source_path.stem)),
        "name": str(record.get("name", source_path.name)),
        "path": str(source_path),
        "kind": _SUPPORTED_SUFFIXES[suffix],
        "page_count": page_count,
        "text": text,
        "warnings": warnings,
    }


def _plan_deck(
    sources: list[dict[str, Any]],
    *,
    topic: str,
    chat_client: _ChatJsonClient | None,
    template: dict[str, Any] | None = None,
) -> tuple[dict[str, Any], dict[str, str]]:
    fallback = _fallback_outline(sources, topic)
    if chat_client is None:
        return fallback, {"mode": "source-grounded", "reason": "未配置演示模型，已使用材料原文生成结构化初稿"}
    template_brief = ""
    if template:
        template_brief = (
            f"\nSelected presentation template: {template.get('name', template.get('id', ''))}. "
            f"Visual tone: {template.get('tone', '')}. "
            "Plan audience-facing presentation pages (context, question, approach, findings, implication, close). "
            "Do not make evidence chains, source verification, or generation process the visible story; put source notes in footers only.\n"
        )
    source_brief = "\n\n".join(
        f"来源：{item['name']}（{item['kind']}，{item['page_count']} 页）\n{item['text'][:_MAX_MODEL_CONTEXT_CHARS // max(1, len(sources))]}"
        for item in sources
    )
    source_brief = template_brief + source_brief
    prompt = f"""演示主题（可为空）：{topic.strip() or '请根据材料命名'}

以下是用户提供的材料。材料可能不完整或含有局限；不得补充材料外的事实。

{source_brief}
"""
    try:
        candidate = chat_client.complete_json(
            [{"role": "system", "content": SCIENTIFIC_SLIDES_CONTRACT}, {"role": "user", "content": prompt}],
            schema_name="scansci_scientific_slides",
        )
        outline = _normalise_model_outline(candidate, fallback=fallback, sources=sources)
        return outline, {"mode": "skill-aware-model", "reason": "已应用好问题、好故事与科研幻灯片内置技能"}
    except Exception as error:  # noqa: BLE001 - deck creation remains available offline
        return fallback, {"mode": "source-grounded", "reason": f"演示规划模型暂不可用，已改用材料原文生成初稿：{str(error)[:180]}"}


def _normalise_model_outline(candidate: Any, *, fallback: dict[str, Any], sources: list[dict[str, Any]]) -> dict[str, Any]:
    if not isinstance(candidate, dict):
        raise ValueError("演示模型没有返回结构化大纲")
    raw_slides = candidate.get("slides")
    if not isinstance(raw_slides, list) or len(raw_slides) < 3:
        raise ValueError("演示模型返回的大纲页数不足")
    max_pages = max(item["page_count"] for item in sources)
    source_corpus = " ".join(str(item.get("text", "")) for item in sources).casefold()
    slides: list[dict[str, Any]] = []
    for raw in raw_slides[:8]:
        if not isinstance(raw, dict):
            continue
        title = _short_text(raw.get("title"), 72)
        takeaway = _ground_model_adaptation_claim(
            _slide_excerpt(raw.get("takeaway"), limit=118),
            source_corpus,
        )
        layout = str(raw.get("layout", "")).strip().lower()
        if layout not in {"cards", "comparison", "process", "branches", "text", "text_focus"}:
            layout = "text"
        bullets = [
            _ground_model_adaptation_claim(_slide_excerpt(item, limit=108), source_corpus)
            for item in list(raw.get("bullets", []) or [])
            if _slide_excerpt(item, limit=108)
        ][:5]
        pages = [int(item) for item in list(raw.get("source_pages", []) or []) if str(item).isdigit() and 1 <= int(item) <= max_pages][:6]
        if title and (takeaway or bullets):
            slides.append({"title": title, "takeaway": takeaway, "layout": layout, "bullets": bullets, "source_pages": pages})
    if len(slides) < 3:
        raise ValueError("演示模型返回的页面内容不完整")
    slides = _rebalance_repetitive_layouts(slides)
    return {
        "title": _short_text(candidate.get("title"), 96) or str(fallback["title"]),
        "central_question": _short_text(candidate.get("central_question"), 300) or str(fallback["central_question"]),
        "story": _short_text(candidate.get("story"), 260) or str(fallback["story"]),
        "slides": slides,
        "evidence_linked": True,
        "source_first": True,
    }


def _ground_model_adaptation_claim(text: str, source_corpus: str) -> str:
    """Correct two common adaptation conflations when the supplied papers disprove them."""

    value = str(text or "").strip()
    folded = value.casefold()
    corpus = str(source_corpus or "").casefold()
    bert_output_only = (
        "bert" in folded
        and any(cue in value for cue in ("仅需微调一层输出层", "只需微调输出层", "仅微调输出层"))
        and "all parameters are fine-tuned" in corpus
    )
    if bert_output_only:
        return "BERT 在下游任务中新增任务输出层，并联合微调全部预训练参数。"
    gpt3_finetuning = (
        "gpt-3" in folded
        and "微调" in value
        and not any(cue in value for cue in ("无需微调", "不需微调", "不进行微调", "没有微调"))
        and "without any gradient updates or fine-tuning" in corpus
    )
    if gpt3_finetuning:
        return "GPT-3 论文的零样本、单样本与少样本评估通过上下文示例适配，不进行梯度更新或微调。"
    return value


def _rebalance_repetitive_layouts(slides: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Turn a one-layout model outline into a small visual story."""

    layouts = {str(item.get("layout", "")) for item in slides}
    if len(slides) < 4 or len(layouts) >= 3 or layouts.intersection({"text", "text_focus"}):
        return slides
    balanced: list[dict[str, Any]] = []
    last_index = len(slides) - 1
    for index, source in enumerate(slides):
        item = dict(source)
        title = str(item.get("title", "")).casefold()
        if index == last_index or any(cue in title for cue in ("总结", "结论", "路线", "分支", "展望")):
            layout = "branches"
        elif any(cue in title for cue in ("对比", "比较", "差异", " vs ", "versus")):
            layout = "comparison"
        elif any(cue in title for cue in ("过程", "演变", "规模", "效率", "阶段", "路径", "从 ")):
            layout = "process"
        elif index in {0, 2, 5} and 3 <= len(list(item.get("bullets", []) or [])) <= 4:
            layout = "process"
        else:
            layout = "cards"
        item["layout"] = layout
        balanced.append(item)
    return balanced


def _fallback_outline(sources: list[dict[str, Any]], topic: str) -> dict[str, Any]:
    title = _short_text(topic, 96) or Path(str(sources[0]["name"])).stem.replace("_", " ") or "研究材料汇报"
    source_slides: list[dict[str, Any]] = []
    comparison_bullets: list[str] = []
    for source in sources[:5]:
        key_sentences = _source_key_sentences(source, topic=title)
        if not key_sentences:
            key_sentences = [f"{_display_source_name(source['name'])} 已导入，但未提取到适合直接展示的完整句子。"]
        source_title = _display_source_name(source["name"])
        source_slides.append(
            {
                "title": source_title,
                "takeaway": _slide_excerpt(key_sentences[0], limit=118),
                "layout": "text",
                "bullets": [_slide_excerpt(item, limit=104) for item in (key_sentences[1:4] or key_sentences[:1])],
                "source_pages": list(range(1, min(3, int(source.get("page_count", 1))) + 1)),
            }
        )
        comparison_bullets.append(_labelled_slide_summary(source_title, key_sentences[0], limit=78))
    slides = [
        {
            "title": "研究范围与材料构成",
            "takeaway": f"本次演示围绕“{title}”整理 {len(sources)} 份上传材料。",
            "layout": "text",
            "bullets": [
                _labelled_slide_summary(item["title"], item["takeaway"], limit=68)
                for item in source_slides[:5]
            ],
            "source_pages": [],
        },
        *source_slides,
        *(
            [
                {
                    "title": "材料之间可以直接比较什么",
                    "takeaway": "以下对比只复述各材料中已经出现的陈述，不把并列路线改写为单线替代关系。",
                    "layout": "text",
                    "bullets": comparison_bullets[:5],
                    "source_pages": [],
                }
            ]
            if len(sources) > 1
            else [
                {
                    "title": "从材料到可核验结论",
                    "takeaway": "先保留材料原意，再组织要点、来源定位与证据边界。",
                    "layout": "process",
                    "bullets": [
                        "定位材料中的完整陈述与关键定义。",
                        "区分原文事实、解释和仍待验证的判断。",
                        "让每个演示要点保留材料名称与页面位置。",
                        "输出前复核数字、图表和结论适用范围。",
                    ],
                    "source_pages": [],
                }
            ]
        ),
        {
            "title": "证据边界与下一步",
            "takeaway": "离线初稿保留来源边界；正式汇报前应复核关键数字与结论范围。",
            "layout": "text",
            "bullets": [
                "所有要点均来自本次上传材料的可提取正文。",
                "图表、公式和被截断的长句需要回到原 PDF 复核。",
                "如需更强的跨文献叙事，可在模型服务恢复后重新生成大纲。",
            ],
            "source_pages": [],
        },
    ]
    return {
        "title": title,
        "central_question": f"围绕“{title}”，这些材料分别提供了哪些可核验的架构、训练目标、实验与适配证据？",
        "story": "先逐篇呈现可核验要点，再比较材料之间的共同基础与差异，最后明确证据边界。",
        "slides": slides,
        "evidence_linked": True,
        "source_first": True,
    }


def _source_key_sentences(source: dict[str, Any], *, topic: str) -> list[str]:
    name = _display_source_name(source.get("name", ""))
    keywords = {
        token.casefold()
        for token in re.findall(r"[A-Za-z][A-Za-z0-9-]{2,}", f"{name} {topic}")
        if token.casefold() not in {"the", "and", "from", "with", "pdf"}
    }
    lower_name = name.casefold()
    if "attention" in lower_name or "transformer" in lower_name:
        keywords.update({"attention", "transformer", "encoder", "decoder", "translation", "bleu"})
    if "bert" in lower_name:
        keywords.update({"bert", "bidirectional", "masked", "fine-tun", "glue", "squad"})
    if "gpt" in lower_name:
        keywords.update({"gpt-3", "few-shot", "zero-shot", "one-shot", "in-context", "language model"})
    candidates: list[tuple[float, int, str]] = []
    raw_sentences = re.split(r"(?<=[。！？.!?])\s+|\n+", str(source.get("text", "")))
    for index, sentence in enumerate(raw_sentences):
        clean = _clean_slide_source_line(sentence)
        had_page_marker = bool(re.match(r"^\[第\s*\d+\s*页\]", clean))
        clean = re.sub(r"^\[第\s*\d+\s*页\]\s*", "", clean)
        if len(clean) > 260:
            boundaries = [
                position + 1
                for position, character in enumerate(clean[:260])
                if position >= 70 and character in "，,；;。.!?！？"
            ]
            if not boundaries:
                continue
            clean = clean[: boundaries[-1]].strip()
        folded = clean.casefold()
        if not 35 <= len(clean) <= 260:
            continue
        if "@" in clean or _looks_like_operational_noise(clean) or any(
            noise in folded
            for noise in (
                "provided proper attribution",
                "listing order is random",
                "google brain",
                "google research",
                "university of toronto",
                "proposed replacing",
                "designed and implemented",
                "was responsible for",
                "we organize the appendix",
                "figure 1:",
            )
        ):
            continue
        score = sum(keyword in folded for keyword in keywords)
        score += 0.8 if any(cue in folded for cue in ("we propose", "we introduce", "we show", "achieves", "outperform", "uses", "model")) else 0
        # In unstructured sources such as README files, the opening definition
        # normally carries more presentation value than later build notes.
        score += max(0.0, 1.25 - (index * 0.08))
        if had_page_marker and "question:" not in folded and "研究问题" not in folded:
            score -= 4
        if score <= 0:
            continue
        candidates.append((float(score), index, clean))
    candidates.sort(key=lambda item: (-item[0], item[1]))
    if not candidates:
        for index, sentence in enumerate(raw_sentences):
            clean = _clean_slide_source_line(sentence)
            if 35 <= len(clean) <= 260 and "@" not in clean and not _looks_like_operational_noise(clean):
                excerpt = _slide_excerpt(clean, limit=155)
                if excerpt:
                    candidates.append((0.1, index, excerpt))
            if len(candidates) >= 5:
                break
    output: list[str] = []
    seen: set[str] = set()
    for _, _, sentence in candidates:
        key = sentence.casefold()
        if key in seen:
            continue
        seen.add(key)
        output.append(_slide_excerpt(sentence, limit=155))
        if len(output) >= 5:
            break
    return output


def _clean_slide_source_line(value: Any) -> str:
    """Remove markup that is useful in a source file but noisy on a slide."""

    text = str(value or "")
    text = re.sub(r"```.*?```", " ", text, flags=re.DOTALL)
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"!\[([^\]]*)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"^\s{0,3}#{1,6}\s*", "", text)
    text = re.sub(r"^\s*(?:[-*+] |\d+[.)]\s+)", "", text)
    text = " ".join(text.split()).strip(" |")
    return re.sub(r"^([A-Za-z][\w.-]{2,})\s+\1\b", r"\1", text, flags=re.IGNORECASE)


def _labelled_slide_summary(title: Any, summary: Any, *, limit: int) -> str:
    label = _display_source_name(title)
    excerpt = _slide_excerpt(summary, limit=limit)
    first_token = label.split(maxsplit=1)[0] if label else ""
    if first_token and excerpt.casefold().startswith(first_token.casefold() + " "):
        excerpt = excerpt[len(first_token) :].lstrip(" ：:-—")
    return f"{label}：{excerpt}" if label else excerpt


def _looks_like_operational_noise(value: str) -> bool:
    """Reject install, preview and shell instructions from offline slide prose."""

    folded = str(value or "").casefold()
    return any(
        cue in folded
        for cue in (
            "127.0.0.1",
            "localhost",
            "start-process",
            "browser preview",
            "output preview",
            "always serves this checkout",
            "npm install",
            "npm run",
            "pip install",
            "python -m",
            "powershell",
            "http://",
            "https://",
            "windows 构建",
            "构建会将",
            "用户无需单独安装",
            "正式发布由",
            "发布门禁",
            "安装包",
            "代码目录",
            "在浏览器打开",
            "兼容入口",
            "scansci-html 命令",
        )
    )


def _render_pptx(path: Path, outline: dict[str, Any], sources: list[dict[str, Any]], *, theme: dict[str, Any]) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slides = _rebalance_repetitive_layouts([dict(item) for item in list(outline.get("slides", []) or [])])
    if not slides:
        raise ValueError("幻灯片大纲为空")
    cover = presentation.slides.add_slide(presentation.slide_layouts[6])
    _render_cover(cover, outline, sources, theme=theme)
    total = len(slides) + 1
    for index, slide_data in enumerate(slides, start=2):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _render_content_slide(slide, slide_data, index=index, total=total, sources=sources, theme=theme)
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_name(f".{path.name}.{uuid4().hex}.tmp.pptx")
    try:
        presentation.save(str(temporary))
        temporary.replace(path)
    finally:
        temporary.unlink(missing_ok=True)


def _render_cover(slide: Any, outline: dict[str, Any], sources: list[dict[str, Any]], *, theme: dict[str, Any]) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*theme["cover_bg"])
    title = str(outline.get("title", "研究汇报"))
    title_size, title_height, question_top = _cover_title_metrics(title)
    _add_text(
        slide,
        _balanced_cover_lines(title),
        0.9,
        1.35,
        11.4,
        title_height,
        size=title_size,
        color=theme["cover_text"],
        bold=True,
        preserve_newlines=True,
    )
    central_question = str(outline.get("central_question", ""))
    _add_text(
        slide,
        _balanced_cover_lines(central_question),
        0.94,
        question_top,
        10.8,
        1.55,
        size=16 if len(central_question) > 110 else 18 if len(central_question) > 72 else 21,
        color=theme["cover_muted"],
        preserve_newlines=True,
    )
    _add_text(slide, f"基于 {len(sources)} 份上传材料生成", 0.94, 5.95, 10.7, 0.35, size=12, color=(167, 190, 199))
    _add_text(slide, " · ".join(_display_source_name(item["name"]) for item in sources), 0.94, 6.35, 10.8, 0.38, size=12, color=(207, 232, 230))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(0.88), Inches(1.65), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*theme["accent"])
    bar.line.fill.background()


def _render_content_slide(
    slide: Any,
    data: dict[str, Any],
    *,
    index: int,
    total: int,
    sources: list[dict[str, Any]],
    theme: dict[str, Any],
) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*theme["background"])
    _render_template_chrome(slide, theme=theme, index=index, total=total)
    content_left = float(theme.get("content_left", 1.02))
    content_width = 12.15 - content_left
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(content_left - 0.28), Inches(0.55), Inches(0.11), Inches(0.78))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(*theme["accent"])
    accent.line.fill.background()
    title = str(data.get("title", "研究要点"))
    title_size, title_height = _content_title_metrics(title)
    _add_text(slide, title, content_left, 0.62, content_width, title_height, size=title_size, color=theme["heading"], bold=True)
    takeaway = str(data.get("takeaway", "")).strip()
    takeaway_top = 0.72 + title_height
    takeaway_height = 0.76 if len(takeaway) > 100 else 0.52
    if takeaway:
        _add_text(slide, takeaway, content_left, takeaway_top, content_width - 0.2, takeaway_height, size=17, color=theme["body"])
    bullets = list(data.get("bullets", []) or [])[:5]
    module_top = takeaway_top + takeaway_height + 0.25
    if str(theme.get("chrome", "header")) == "header":
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0),
            Inches(module_top - 0.12),
            Inches(13.333),
            Inches(0.05),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = RGBColor(*theme["accent"])
        divider.line.fill.background()
        module_top += 0.08
    _render_content_modules(
        slide,
        bullets,
        layout=str(data.get("layout", "")),
        index=index,
        theme=theme,
        left=content_left,
        width=content_width - 0.2,
        top=module_top,
    )
    pages = [str(item) for item in list(data.get("source_pages", []) or []) if str(item).strip()]
    source_names = " · ".join(_display_source_name(item.get("name", "")) for item in sources[:3])
    source_label = source_names or "本次上传材料"
    if pages:
        source_label += f" | 涉及原文页码 {', '.join(pages)}"
    _add_text(slide, source_label, content_left, 6.57, content_width - 1.0, 0.38, size=11, color=theme["muted"])


def _template_theme(template: dict[str, Any] | None) -> dict[str, Any]:
    """Translate an EasySlides template into editable PPTX design tokens."""

    template_id = str((template or {}).get("id", "academic_general"))
    themes: dict[str, dict[str, Any]] = {
        "academic_general": {
            "cover_bg": (0, 51, 102), "cover_text": (255, 255, 255), "cover_muted": (232, 244, 252),
            "background": (255, 255, 255), "heading": (51, 51, 51), "body": (51, 51, 51),
            "muted": (102, 102, 102), "primary": (0, 51, 102), "accent": (0, 102, 204),
            "emphasis": (204, 0, 0), "surface": (245, 247, 250), "border": (208, 215, 224),
            "chrome": "header", "content_left": 1.02,
        },
        "academic_scqa": {
            "cover_bg": (11, 47, 107), "cover_text": (255, 255, 255), "cover_muted": (207, 240, 250),
            "background": (247, 250, 255), "heading": (11, 47, 107), "body": (42, 68, 105),
            "muted": (103, 124, 150), "primary": (11, 47, 107), "accent": (0, 166, 214),
            "emphasis": (0, 137, 123), "surface": (241, 246, 255), "border": (188, 213, 239),
            "chrome": "header", "content_left": 1.02,
        },
        "defense_leftnav": {
            "cover_bg": (139, 0, 18), "cover_text": (255, 255, 255), "cover_muted": (250, 228, 232),
            "background": (255, 255, 255), "heading": (61, 23, 31), "body": (51, 51, 51),
            "muted": (112, 112, 112), "primary": (139, 0, 18), "accent": (139, 0, 18),
            "emphasis": (104, 0, 13), "surface": (253, 245, 246), "border": (228, 205, 209),
            "chrome": "leftnav", "content_left": 3.15,
        },
        "defense_topnav": {
            "cover_bg": (24, 58, 106), "cover_text": (255, 255, 255), "cover_muted": (226, 234, 245),
            "background": (255, 255, 255), "heading": (24, 58, 106), "body": (45, 55, 72),
            "muted": (110, 117, 128), "primary": (24, 58, 106), "accent": (24, 58, 106),
            "emphasis": (92, 115, 145), "surface": (231, 230, 230), "border": (210, 214, 220),
            "chrome": "topnav", "content_left": 0.92,
        },
        "literature_minimal": {
            "cover_bg": (255, 255, 255), "cover_text": (14, 40, 65), "cover_muted": (83, 111, 139),
            "background": (255, 255, 255), "heading": (14, 40, 65), "body": (35, 48, 61),
            "muted": (113, 126, 139), "primary": (13, 93, 190), "accent": (13, 93, 190),
            "emphasis": (13, 93, 190), "surface": (246, 250, 255), "border": (211, 224, 239),
            "chrome": "minimal", "content_left": 0.92,
        },
    }
    return dict(themes.get(template_id, themes["academic_general"]))


def _render_template_chrome(slide: Any, *, theme: dict[str, Any], index: int, total: int) -> None:
    primary = RGBColor(*theme["primary"])
    accent = RGBColor(*theme["accent"])
    chrome = str(theme.get("chrome", "header"))
    if chrome == "leftnav":
        rail = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(2.55), Inches(7.5))
        rail.fill.solid(); rail.fill.fore_color.rgb = RGBColor(241, 241, 241); rail.line.fill.background()
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.35 + ((index - 2) % 5) * 0.66), Inches(2.72), Inches(0.54))
        band.fill.solid(); band.fill.fore_color.rgb = primary; band.line.fill.background()
        labels = ["研究背景", "研究目标", "材料与方法", "结果与分析", "结论与讨论"]
        for position, label in enumerate(labels):
            _add_text(slide, label, 0.36, 1.49 + position * 0.66, 1.9, 0.24, size=12, color=(255, 255, 255) if position == (index - 2) % 5 else (99, 99, 99), bold=position == (index - 2) % 5)
    elif chrome == "topnav":
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.48))
        bar.fill.solid(); bar.fill.fore_color.rgb = primary; bar.line.fill.background()
        labels = ["背景", "目标", "方法", "结果", "结论"]
        active = (index - 2) % len(labels)
        for position, label in enumerate(labels):
            x = 3.1 + position * 1.45
            if position == active:
                tab = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x - 0.08), Inches(0.08), Inches(1.22), Inches(0.34))
                tab.fill.solid(); tab.fill.fore_color.rgb = RGBColor(255, 255, 255); tab.line.fill.background()
            _add_text(slide, label, x, 0.14, 1.05, 0.16, size=10, color=theme["primary"] if position == active else (231, 237, 245), alignment=PP_ALIGN.CENTER)
    elif chrome == "minimal":
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.48), Inches(0.38), Inches(12.36), Inches(0.04))
        line.fill.solid(); line.fill.fore_color.rgb = accent; line.line.fill.background()
        footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.24), Inches(13.333), Inches(0.26))
        footer.fill.solid(); footer.fill.fore_color.rgb = accent; footer.line.fill.background()
    else:
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.46))
        header.fill.solid(); header.fill.fore_color.rgb = primary; header.line.fill.background()
    _add_text(slide, f"{index}/{total}", 11.86, 6.76, 0.55, 0.3, size=10, color=theme["muted"], alignment=PP_ALIGN.RIGHT)


def _render_content_modules(
    slide: Any,
    bullets: list[str],
    *,
    layout: str,
    index: int,
    theme: dict[str, Any],
    left: float,
    width: float,
    top: float,
) -> None:
    items = bullets or ["根据材料提炼本页要点"]
    layout_kinds = {"cards": 0, "comparison": 1, "process": 2, "branches": 3}
    kind = layout_kinds.get(str(layout).strip().lower(), (index - 2) % 3)
    gap = 0.22
    surface = RGBColor(*theme["surface"])
    border = RGBColor(*theme["border"])
    accent = RGBColor(*theme["accent"])
    if kind == 0:
        count = min(3, len(items))
        card_width = (width - gap * (count - 1)) / count
        groups = _balanced_item_groups(items, count)
        for position, group in enumerate(groups):
            x = left + position * (card_width + gap)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(card_width), Inches(3.28))
            card.fill.solid(); card.fill.fore_color.rgb = surface; card.line.color.rgb = border
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.22), Inches(top + 0.24), Inches(0.42), Inches(0.42))
            marker.fill.solid(); marker.fill.fore_color.rgb = accent; marker.line.fill.background()
            _add_text(slide, str(position + 1), x + 0.22, top + 0.29, 0.42, 0.14, size=10, color=(255, 255, 255), bold=True, alignment=PP_ALIGN.CENTER)
            text = _group_text(group)
            font_size = 14 if len(text) > 150 else 15 if len(text) > 100 else 16
            _add_group_text(slide, group, x + 0.22, top + 0.82, card_width - 0.44, 2.12, size=font_size, color=theme["body"])
    elif kind == 1:
        midpoint = left + (width - gap) / 2
        groups = _balanced_item_groups(items, min(2, len(items)))
        for position, group in enumerate(groups):
            x = left if position == 0 else midpoint + gap
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches((width - gap) / 2), Inches(3.28))
            card.fill.solid(); card.fill.fore_color.rgb = surface; card.line.color.rgb = border
            _add_text(slide, "证据" if position == 0 else "比较", x + 0.28, top + 0.28, 1.0, 0.28, size=13, color=theme["accent"], bold=True)
            text = _group_text(group)
            font_size = 14 if len(text) > 190 else 15 if len(text) > 125 else 17
            _add_group_text(slide, group, x + 0.28, top + 0.82, (width - gap) / 2 - 0.55, 2.18, size=font_size, color=theme["body"])
    elif kind == 2:
        count = min(4, len(items))
        node_width = (width - 0.4) / count
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + 0.55), Inches(top + 1.18), Inches(max(0.5, node_width * (count - 1))), Inches(0.05))
        line.fill.solid(); line.fill.fore_color.rgb = accent; line.line.fill.background()
        for position, item in enumerate(items[:count]):
            x = left + position * node_width
            node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.38), Inches(top + 0.88), Inches(0.62), Inches(0.62))
            node.fill.solid(); node.fill.fore_color.rgb = accent; node.line.fill.background()
            _add_text(slide, str(position + 1), x + 0.38, top + 1.08, 0.62, 0.16, size=11, color=(255, 255, 255), bold=True, alignment=PP_ALIGN.CENTER)
            _add_text(slide, item, x, top + 1.78, node_width - 0.12, 1.15, size=15, color=theme["body"], alignment=PP_ALIGN.CENTER)
    else:
        groups = _balanced_item_groups(items, 2)
        root_width = min(3.1, width * 0.34)
        root_left = left + (width - root_width) / 2
        root = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(root_left), Inches(top + 0.05), Inches(root_width), Inches(0.78))
        root.fill.solid(); root.fill.fore_color.rgb = accent; root.line.fill.background()
        _add_text(slide, "共享基础", root_left + 0.18, top + 0.28, root_width - 0.36, 0.24, size=16, color=(255, 255, 255), bold=True, alignment=PP_ALIGN.CENTER)
        branch_width = (width - gap) / 2
        junction_left = left + branch_width / 2
        junction_width = branch_width + gap
        trunk = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(root_left + root_width / 2 - 0.025), Inches(top + 0.83), Inches(0.05), Inches(0.31))
        trunk.fill.solid(); trunk.fill.fore_color.rgb = accent; trunk.line.fill.background()
        crossbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(junction_left), Inches(top + 1.11), Inches(junction_width), Inches(0.05))
        crossbar.fill.solid(); crossbar.fill.fore_color.rgb = accent; crossbar.line.fill.background()
        for position, group in enumerate(groups):
            x = left + position * (branch_width + gap)
            connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + branch_width / 2 - 0.025), Inches(top + 1.11), Inches(0.05), Inches(0.16))
            connector.fill.solid(); connector.fill.fore_color.rgb = accent; connector.line.fill.background()
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(top + 1.22), Inches(branch_width), Inches(2.15))
            card.fill.solid(); card.fill.fore_color.rgb = surface; card.line.color.rgb = border
            _add_text(slide, f"分支 {position + 1}", x + 0.25, top + 1.47, branch_width - 0.5, 0.28, size=13, color=theme["accent"], bold=True)
            _add_group_text(slide, group, x + 0.25, top + 1.92, branch_width - 0.5, 1.15, size=16, color=theme["body"])


def _balanced_item_groups(items: list[str], group_count: int) -> list[list[str]]:
    count = max(1, min(int(group_count), len(items)))
    base_size, remainder = divmod(len(items), count)
    groups: list[list[str]] = []
    cursor = 0
    for index in range(count):
        size = base_size + (1 if index < remainder else 0)
        groups.append([str(item) for item in items[cursor : cursor + size]])
        cursor += size
    return groups


def _group_text(items: list[str]) -> str:
    if len(items) <= 1:
        return str(items[0]) if items else ""
    return "\n".join(f"• {item}" for item in items)


def _display_source_name(value: object) -> str:
    name = Path(str(value or "")).stem
    return " ".join(name.replace("_", " ").replace("-", "-").split()) or "上传材料"


def _balanced_cover_lines(value: object) -> str:
    """Insert one deliberate line break instead of leaving a one-word orphan."""

    text = " ".join(str(value or "").split()).strip()
    if len(text) < 34:
        return text
    midpoint = len(text) / 2
    candidates = [
        index + 1
        for index, character in enumerate(text)
        if character in "：:；;，,？?" or (character.isspace() and index > 8)
    ]
    if not candidates:
        return text
    split = min(candidates, key=lambda index: abs(index - midpoint))
    if split < len(text) * 0.28 or split > len(text) * 0.72:
        return text
    return f"{text[:split].rstrip()}\n{text[split:].lstrip()}"


def _add_group_text(
    slide: Any,
    items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int,
    color: tuple[int, int, int],
) -> None:
    """Render each item as its own paragraph so bullets never collapse inline."""

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    rows = [str(item).strip() for item in items if str(item).strip()] or ["根据材料提炼本页要点"]
    for index, value in enumerate(rows):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"• {_short_text(value, 180)}" if len(rows) > 1 else _short_text(value, 180)
        paragraph.space_after = Pt(9)
        paragraph.font.name = "Microsoft YaHei"
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = RGBColor(*color)


def _add_text(
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: int,
    color: tuple[int, int, int],
    bold: bool = False,
    alignment: Any | None = None,
    preserve_newlines: bool = False,
) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    paragraph = frame.paragraphs[0]
    if preserve_newlines:
        paragraph.text = "\n".join(" ".join(line.split()) for line in str(text).splitlines())[:900]
    else:
        paragraph.text = _short_text(text, 900)
    if alignment is not None:
        paragraph.alignment = alignment
    for run in paragraph.runs:
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*color)
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0


def _add_bullets(slide: Any, bullets: list[str], left: float, top: float, width: float, height: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    for index, value in enumerate(bullets or ["材料中尚未提炼出足够的要点。"]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = _short_text(value, 155)
        paragraph.level = 0
        paragraph.space_after = Pt(14)
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = RGBColor(30, 54, 69)
        paragraph.font.bold = False


def _cover_title_metrics(title: str) -> tuple[int, float, float]:
    """Reserve vertical space for long scientific titles before rendering."""

    length = len(title.strip())
    if length > 78:
        return 26, 2.7, 4.3
    if length > 48:
        return 27, 2.35, 4.05
    if length > 36:
        return 30, 2.05, 3.85
    return 42, 1.35, 3.05


def _content_title_metrics(title: str) -> tuple[int, float]:
    """Keep claim titles distinct from the takeaway instead of letting them overlap."""

    length = len(title.strip())
    if length > 68:
        return 24, 1.22
    if length > 34:
        return 28, 1.02
    return 30, 0.72


def _presentation_path(workspace: str | Path, title: str) -> Path:
    root = Path(workspace).resolve().parent / "presentations"
    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    stem = _SAFE_NAME.sub("_", title).strip("._ ")[:70] or "scansci_slides"
    candidate = root / f"{stem}_{stamp}.pptx"
    counter = 2
    while candidate.exists():
        candidate = root / f"{stem}_{stamp}-{counter}.pptx"
        counter += 1
    return candidate


def _sentences(value: str) -> list[str]:
    chunks = re.split(r"(?<=[。！？.!?])\s+|\n+", value)
    rows: list[str] = []
    for chunk in chunks:
        candidate = _compact_text(chunk, limit=190)
        if len(candidate) >= 24 and candidate not in rows:
            rows.append(candidate)
    return rows


def _compact_text(value: str, *, limit: int) -> str:
    return " ".join(str(value or "").split())[:limit].strip()


def _short_text(value: Any, limit: int) -> str:
    return _compact_text(str(value or ""), limit=limit)


def _slide_excerpt(value: Any, *, limit: int) -> str:
    """Return a visibly bounded excerpt without silently cutting a word in half."""

    text = " ".join(str(value or "").split()).strip(" •")
    if len(text) <= limit:
        return text
    window = text[: limit + 1]
    boundaries = [
        position + 1
        for position, character in enumerate(window)
        if position >= max(36, int(limit * 0.55)) and (character.isspace() or character in "，,；;：:。.!?！？")
    ]
    if boundaries:
        clipped = window[: boundaries[-1]].rstrip(" ，,；;：:")
    else:
        clipped = window[:limit].rstrip(" ，,；;：:")
    return f"{clipped}…" if clipped else ""
