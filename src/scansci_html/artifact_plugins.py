"""Built-in office and LaTeX artifact tools for the ScanSci Pi runtime.

The tools in this module are deliberately narrow.  They only create artifacts
under the local ScanSci workspace and never expose an arbitrary shell or an
unbounded filesystem writer to the model.
"""

from __future__ import annotations

from importlib.util import find_spec
import json
import os
from pathlib import Path
import re
import shutil
import subprocess
from typing import Any
from zipfile import ZipFile


_SAFE_NAME = re.compile(r"[^A-Za-z0-9._\-\u4e00-\u9fff]+")


BUILTIN_PLUGINS: list[dict[str, Any]] = [
    {
        "id": "zotero",
        "name": "Zotero",
        "description": "连接本机 Zotero 文献库，检索条目与附件全文，并导出 BibTeX 或格式化引用。",
        "source": "ScanSci built-in · adapted from OpenAI Zotero plugin · MIT",
        "enabled": True,
        "builtin": True,
        "icon": "zotero",
        "skills": ["Search library", "Read attachments", "Export BibTeX", "Format citations"],
        "tool_names": [
            "zotero_search",
            "zotero_status",
            "zotero_fulltext",
            "zotero_attachment",
            "zotero_export_bibtex",
            "zotero_citations",
        ],
    },
    {
        "id": "documents",
        "name": "Documents",
        "description": "创建和读取 DOCX 文档，并在交付前验证文档结构。",
        "source": "ScanSci built-in · adapted from Codex Documents workflow",
        "enabled": True,
        "builtin": True,
        "icon": "document",
        "skills": ["创建 DOCX", "读取 DOCX", "结构校验"],
        "tool_names": ["create_document"],
    },
    {
        "id": "pdf",
        "name": "PDF",
        "description": "创建、读取并检查 PDF 页数、元数据和可提取文本。",
        "source": "ScanSci built-in · adapted from Codex PDF workflow",
        "enabled": True,
        "builtin": True,
        "icon": "pdf",
        "skills": ["创建 PDF", "读取 PDF", "版面前置校验"],
        "tool_names": ["create_pdf"],
    },
    {
        "id": "spreadsheets",
        "name": "Spreadsheets",
        "description": "创建可编辑 XLSX，保留数值类型、公式与基础表格样式。",
        "source": "ScanSci built-in · adapted from Codex Spreadsheets workflow",
        "enabled": True,
        "builtin": True,
        "icon": "spreadsheet",
        "skills": ["创建 XLSX", "公式写入", "工作簿校验"],
        "tool_names": ["create_spreadsheet"],
    },
    {
        "id": "presentations",
        "name": "Presentations",
        "description": "创建可编辑 PPTX，并与 ScanSci 的科研幻灯片工作流共用产物目录。",
        "source": "ScanSci built-in · adapted from Codex Presentations workflow",
        "enabled": True,
        "builtin": True,
        "icon": "presentation",
        "skills": ["创建 PPTX", "讲述结构", "文件校验"],
        "tool_names": ["create_presentation", "build_presentation_outline"],
    },
    {
        "id": "latex",
        "name": "LaTeX",
        "description": "优先使用内置 Tectonic 编译 TeX；复杂工程可回退到现有 TeX Live。",
        "source": "ScanSci built-in · adapted from Codex LaTeX workflow",
        "enabled": True,
        "builtin": True,
        "icon": "latex",
        "skills": ["LaTeX Doctor", "LaTeX Compile", "TeX Live 检测"],
        "tool_names": ["compile_latex"],
    },
]


def default_plugin_records() -> list[dict[str, Any]]:
    return [{**item, "skills": list(item["skills"]), "tool_names": list(item["tool_names"])} for item in BUILTIN_PLUGINS]


def plugin_runtime_statuses() -> dict[str, dict[str, Any]]:
    """Return inexpensive local readiness information for the extension UI."""

    modules = {
        "documents": ["docx"],
        "spreadsheets": ["openpyxl"],
        "presentations": ["pptx"],
    }
    statuses: dict[str, dict[str, Any]] = {
        "zotero": {
            "ready": True,
            "status": "ready",
            "detail": "内置连接器",
            "missing": [],
        }
    }
    for plugin_id, requirements in modules.items():
        missing = [name for name in requirements if find_spec(name) is None]
        statuses[plugin_id] = {
            "ready": not missing,
            "status": "ready" if not missing else "missing_dependency",
            "detail": "可用" if not missing else f"缺少依赖：{', '.join(missing)}",
            "missing": missing,
        }
    pdf_missing = [name for name in ("pypdf",) if find_spec(name) is None]
    if find_spec("reportlab") is None and find_spec("fitz") is None:
        pdf_missing.append("reportlab or PyMuPDF")
    statuses["pdf"] = {
        "ready": not pdf_missing,
        "status": "ready" if not pdf_missing else "missing_dependency",
        "detail": "可用" if not pdf_missing else f"缺少依赖：{', '.join(pdf_missing)}",
        "missing": pdf_missing,
    }
    tectonic = find_tectonic()
    texlive = shutil.which("latexmk") or shutil.which("pdflatex")
    statuses["latex"] = {
        "ready": bool(tectonic or texlive),
        "status": "ready" if tectonic or texlive else "missing_runtime",
        "detail": "Tectonic 可用" if tectonic else ("TeX Live 可用" if texlive else "未检测到 Tectonic 或 TeX Live"),
        "runtime": str(tectonic or texlive or ""),
    }
    return statuses


def enrich_builtin_plugins(records: list[dict[str, Any]]) -> list[dict[str, Any]]:
    statuses = plugin_runtime_statuses()
    return [{**record, "runtime": statuses.get(str(record.get("id", "")), {})} for record in records]


def execute_artifact_tool(name: str, arguments: dict[str, Any], *, workspace: str | Path) -> dict[str, Any]:
    dispatch = {
        "create_document": create_document_artifact,
        "create_pdf": create_pdf_artifact,
        "create_spreadsheet": create_spreadsheet_artifact,
        "create_presentation": create_presentation_artifact,
        "compile_latex": compile_latex_artifact,
    }
    handler = dispatch.get(name)
    if handler is None:
        raise ValueError(f"Unsupported artifact tool: {name}")
    return handler(arguments, workspace=workspace)


def create_document_artifact(arguments: dict[str, Any], *, workspace: str | Path) -> dict[str, Any]:
    from docx import Document
    from docx.shared import Pt

    title = _text(arguments.get("title"), "ScanSci 文档")
    content = _text(arguments.get("content"), "")
    if not content:
        raise ValueError("content must not be empty")
    path = _artifact_path(workspace, "documents", arguments.get("output_name"), title, ".docx")
    document = Document()
    styles = document.styles
    styles["Normal"].font.name = "Microsoft YaHei"
    styles["Normal"].font.size = Pt(10.5)
    document.add_heading(title, level=0)
    for raw_line in content.splitlines():
        line = raw_line.strip()
        if not line:
            document.add_paragraph()
        elif line.startswith("### "):
            document.add_heading(line[4:], level=3)
        elif line.startswith("## "):
            document.add_heading(line[3:], level=2)
        elif line.startswith("# "):
            document.add_heading(line[2:], level=1)
        elif line.startswith(("- ", "* ")):
            document.add_paragraph(line[2:], style="List Bullet")
        else:
            document.add_paragraph(line)
    document.save(path)
    with ZipFile(path) as archive:
        required = {"[Content_Types].xml", "word/document.xml"}
        if not required.issubset(archive.namelist()):
            raise RuntimeError("DOCX structure verification failed")
    return _result("document", path, title, paragraphs=len(document.paragraphs))


def create_pdf_artifact(arguments: dict[str, Any], *, workspace: str | Path) -> dict[str, Any]:
    from pypdf import PdfReader

    title = _text(arguments.get("title"), "ScanSci PDF")
    content = _text(arguments.get("content"), "")
    if not content:
        raise ValueError("content must not be empty")
    path = _artifact_path(workspace, "pdf", arguments.get("output_name"), title, ".pdf")
    if find_spec("reportlab") is not None:
        _write_pdf_with_reportlab(path, title, content)
    elif find_spec("fitz") is not None:
        _write_pdf_with_pymupdf(path, title, content)
    else:
        raise ModuleNotFoundError("PDF creation requires reportlab or PyMuPDF")
    reader = PdfReader(str(path))
    if not reader.pages:
        raise RuntimeError("PDF verification failed")
    return _result("pdf", path, title, pages=len(reader.pages))


def _write_pdf_with_reportlab(path: Path, title: str, content: str) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.pdfgen.canvas import Canvas

    font_name = "Helvetica"
    for font_path in (
        Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts" / "msyh.ttc",
        Path(os.environ.get("WINDIR", "C:/Windows")) / "Fonts" / "simhei.ttf",
    ):
        if font_path.is_file():
            try:
                pdfmetrics.registerFont(TTFont("ScanSciCJK", str(font_path), subfontIndex=0))
                font_name = "ScanSciCJK"
                break
            except Exception:
                continue
    canvas = Canvas(str(path), pagesize=A4)
    _, height = A4
    y = height - 58
    canvas.setFont(font_name, 18)
    canvas.drawString(54, y, title[:70])
    y -= 34
    canvas.setFont(font_name, 10.5)
    for paragraph in content.splitlines() or [content]:
        for line in _wrap_text(paragraph, 82) or [""]:
            if y < 54:
                canvas.showPage()
                canvas.setFont(font_name, 10.5)
                y = height - 54
            canvas.drawString(54, y, line)
            y -= 16
        y -= 5
    canvas.save()


def _write_pdf_with_pymupdf(path: Path, title: str, content: str) -> None:
    import fitz

    document = fitz.open()
    page = document.new_page(width=595, height=842)
    page.insert_textbox(fitz.Rect(54, 50, 541, 92), title, fontname="china-s", fontsize=18)
    body = "\n".join(content.splitlines())
    page.insert_textbox(fitz.Rect(54, 108, 541, 788), body, fontname="china-s", fontsize=10.5, lineheight=1.45)
    document.save(path)
    document.close()


def create_spreadsheet_artifact(arguments: dict[str, Any], *, workspace: str | Path) -> dict[str, Any]:
    from openpyxl import Workbook, load_workbook
    from openpyxl.styles import Font, PatternFill

    title = _text(arguments.get("title"), "ScanSci 表格")
    columns = [str(value) for value in list(arguments.get("columns") or [])]
    rows = list(arguments.get("rows") or [])
    if not columns:
        raise ValueError("columns must contain at least one header")
    path = _artifact_path(workspace, "spreadsheets", arguments.get("output_name"), title, ".xlsx")
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = _safe_sheet_name(title)
    sheet.append(columns)
    for row in rows[:10_000]:
        if isinstance(row, dict):
            values = [row.get(column, "") for column in columns]
        elif isinstance(row, (list, tuple)):
            values = list(row)[: len(columns)]
        else:
            values = [row]
        sheet.append(values)
    for cell in sheet[1]:
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = PatternFill("solid", fgColor="267454")
    sheet.freeze_panes = "A2"
    sheet.auto_filter.ref = sheet.dimensions
    for column_cells in sheet.columns:
        width = min(42, max(10, max(len(str(cell.value or "")) for cell in column_cells) + 2))
        sheet.column_dimensions[column_cells[0].column_letter].width = width
    workbook.save(path)
    verified = load_workbook(path, data_only=False, read_only=True)
    row_count = verified.active.max_row
    verified.close()
    return _result("spreadsheet", path, title, rows=max(0, row_count - 1), columns=len(columns))


def create_presentation_artifact(arguments: dict[str, Any], *, workspace: str | Path) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    title = _text(arguments.get("title"), "ScanSci 演示文稿")
    slides = list(arguments.get("slides") or [])
    if not slides:
        raise ValueError("slides must contain at least one slide")
    path = _artifact_path(workspace, "presentations", arguments.get("output_name"), title, ".pptx")
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = title
    if len(cover.placeholders) > 1:
        cover.placeholders[1].text = _text(arguments.get("subtitle"), "ScanSci")
    for item in slides[:60]:
        record = dict(item or {}) if isinstance(item, dict) else {"title": str(item), "bullets": []}
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = _text(record.get("title"), "未命名页面")
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = list(record.get("bullets") or [])
        if not bullets:
            bullets = [_text(record.get("content"), "")]
        for index, bullet in enumerate(value for value in bullets if str(value).strip()):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.font.size = Pt(22)
    presentation.save(path)
    with ZipFile(path) as archive:
        if "ppt/presentation.xml" not in archive.namelist():
            raise RuntimeError("PPTX structure verification failed")
    return _result("presentation", path, title, slides=len(presentation.slides))


def compile_latex_artifact(arguments: dict[str, Any], *, workspace: str | Path) -> dict[str, Any]:
    source = _text(arguments.get("source"), "")
    if not source:
        raise ValueError("source must not be empty")
    title = _text(arguments.get("title"), "ScanSci LaTeX")
    root = _artifact_root(workspace, "latex")
    stem = _safe_stem(arguments.get("output_name"), title)
    tex_path = root / f"{stem}.tex"
    tex_path.write_text(source, encoding="utf-8")
    tectonic = find_tectonic()
    if tectonic:
        command = [str(tectonic), "-X", "compile", "--outdir", str(root), "--outfmt", "pdf", "--print", "--untrusted", tex_path.name]
        runtime = "tectonic"
    else:
        latexmk = shutil.which("latexmk")
        if not latexmk:
            raise FileNotFoundError("未检测到 Tectonic 或 TeX Live；请先在插件页运行 LaTeX Doctor")
        command = [latexmk, "-pdf", "-interaction=nonstopmode", "-halt-on-error", f"-outdir={root}", str(tex_path)]
        runtime = "texlive"
    completed = subprocess.run(command, cwd=root, capture_output=True, text=True, timeout=180, check=False)
    pdf_path = root / f"{stem}.pdf"
    if completed.returncode != 0 or not pdf_path.is_file():
        message = (completed.stderr or completed.stdout or "LaTeX compilation failed")[-4000:]
        raise RuntimeError(message)
    return _result("latex", pdf_path, title, source_path=str(tex_path), runtime=runtime, log=(completed.stdout or "")[-2000:])


def find_tectonic() -> Path | None:
    candidates = []
    env_path = os.environ.get("SCANSCI_TECTONIC_PATH", "").strip()
    if env_path:
        candidates.append(Path(env_path))
    package_root = Path(__file__).resolve().parent
    candidates.extend((package_root / "runtime" / "latex" / name for name in ("tectonic.exe", "tectonic")))
    path_command = shutil.which("tectonic")
    if path_command:
        candidates.append(Path(path_command))
    codex_plugins = Path.home() / ".codex" / "plugins" / "cache" / "openai-bundled" / "latex"
    if codex_plugins.is_dir():
        candidates.extend(sorted(codex_plugins.glob("*/bin/tectonic.exe"), reverse=True))
    for candidate in candidates:
        if candidate.is_file():
            return candidate.resolve()
    return None


def _artifact_root(workspace: str | Path, kind: str) -> Path:
    root = Path(workspace).resolve().parent / "artifacts" / kind
    root.mkdir(parents=True, exist_ok=True)
    return root


def _artifact_path(workspace: str | Path, kind: str, requested: object, title: str, suffix: str) -> Path:
    root = _artifact_root(workspace, kind)
    stem = _safe_stem(requested, title)
    candidate = root / f"{stem}{suffix}"
    index = 2
    while candidate.exists():
        candidate = root / f"{stem}-{index}{suffix}"
        index += 1
    return candidate


def _safe_stem(requested: object, fallback: str) -> str:
    raw = Path(str(requested or "")).stem or fallback
    return (_SAFE_NAME.sub("-", raw).strip("-._") or "scansci-artifact")[:80]


def _safe_sheet_name(value: str) -> str:
    return re.sub(r"[\\/*?:\[\]]", " ", value).strip()[:31] or "Sheet1"


def _text(value: object, fallback: str) -> str:
    text = str(value or "").strip()
    return text or fallback


def _wrap_text(value: str, width: int) -> list[str]:
    text = str(value or "")
    return [text[index : index + width] for index in range(0, len(text), width)]


def _result(kind: str, path: Path, title: str, **details: Any) -> dict[str, Any]:
    return {
        "ok": True,
        "artifact_type": kind,
        "title": title,
        "file_path": str(path),
        "size_bytes": path.stat().st_size,
        **details,
    }


def status_json() -> str:
    return json.dumps(plugin_runtime_statuses(), ensure_ascii=False, sort_keys=True)
