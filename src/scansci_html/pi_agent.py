"""Process-isolated bridge from ScanSci's Python runtime to the Pi Agent SDK.

The Node sidecar owns model orchestration.  Every actionable tool call crosses
this JSONL boundary and is executed by a narrow ScanSci dispatcher, so Pi never
receives shell or arbitrary filesystem tools.
"""

from __future__ import annotations

from collections import deque
from collections.abc import Iterator
from concurrent.futures import Future, ThreadPoolExecutor
from dataclasses import dataclass, field
import hashlib
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
import json
import os
from pathlib import Path
from queue import Empty, Full, Queue
import re
import requests
import shutil
import subprocess
import sys
import threading
import time
from typing import Any, Callable
from uuid import uuid4

from .academic_search import search_academic_papers
from .academic_planning import plan_academic_search
from .agent_capabilities import builtin_capability_descriptor
from .agent_contract import compile_task_contract as compile_host_task_contract
from .agent_reach import run_agent_reach
from .agent_skill_tools import ProgressiveSkillRuntime, SKILL_STATE_SCHEMA
from .app_settings import load_settings
from .artifact_plugins import execute_artifact_tool
from .checkpoints import CheckpointError, CheckpointStore
from .context_policy import build_token_envelope, prune_stale_tool_results
from .image_attachments import validate_pi_image_blocks
from .ingestion import SUPPORTED_INGESTION_SUFFIXES, extract_local_document
from .library_manager import import_library_files
from .model_metadata import ModelRuntimeDescriptor, descriptor_from_model_record
from .qa.agent import answer_question
from .research_runs import ResearchRunStore, StageSpec
from .research_tools import (
    _BatchCancelled,
    analyze_references,
    build_ppt_outline,
    capability_snapshot,
    download_papers,
    search_journals,
    verify_doi_metadata,
)
from .retrieval import search_evidence_store
from .runtime_components import default_node_component
from .web_search import search_public_web
from .web_access import browser_access_read, browser_access_status
from .workspace import initialize_notebook, load_workspace_summary
from .zotero_integration import (
    search_zotero_library,
    zotero_export_bibtex,
    zotero_file_url,
    zotero_formatted_citations,
    zotero_fulltext,
    zotero_status,
)
from .zotero_scope import (
    filter_zotero_result,
    resolve_zotero_tag_scope,
)
from .obsidian_integration import obsidian_backlinks, obsidian_status, read_obsidian_note, search_obsidian_vault
from .prefix_diagnostics import build_prefix_shape
from .run_manifest import RunManifest
from .task_contract import TaskContract
from .tool_authorization import (
    ApprovalToken,
    approval_token_from_response,
    authorize_tool_call,
)


class PiRuntimeUnavailable(RuntimeError):
    """Raised when the bundled Pi sidecar or Node runtime cannot be found."""


class PiMultimodalUnavailable(PiRuntimeUnavailable):
    """Raised when a validated image turn targets a text-only descriptor."""


class PiAgentRunError(RuntimeError):
    """Structured Pi failure with machine-actionable recovery metadata."""

    def __init__(self, message: str, *, failure: dict[str, Any] | None = None) -> None:
        super().__init__(message)
        self.failure = dict(failure or {})


_MANAGED_GATEWAY_HOST = "scansci-glm-gateway.932196440.workers.dev"
_MAX_MODEL_TOOL_RESULT_BYTES = 16_000
_MAX_GATEWAY_REQUEST_BYTES = 2_000_000
_MAX_GATEWAY_RESPONSE_BYTES = 8_000_000
_GATEWAY_TRANSPORT_ATTEMPTS = 3
_GATEWAY_MAX_RETRY_AFTER_SECONDS = 60.0
_SESSION_REGISTRY_LOCK = threading.Lock()
_PI_PROTOCOL_VERSION = 6
_PI_REQUIRED_FEATURES = (
    "task_contract_v2",
    "explicit_empty_leases",
    "host_tool_authorization",
    "structured_mcp_effects",
    "current_request_context",
    "dynamic_tools",
    "ephemeral_sessions",
    "progressive_skills",
    "parallel_tool_dispatch",
    "lifecycle_hooks_v1",
    "acked_session_commands",
    "model_runtime_descriptor",
    "token_envelope",
    "multimodal_turns",
    "scientific_subagents_v1",
    "session_controls_v2",
    "thinking_max",
)
_MAX_JSONL_LINE_BYTES = 20 * 1024 * 1024
_TOOL_TAG_PATTERN = re.compile(r"<SCANSCI_TOOL_CALL>\s*(?P<body>.*?)\s*</SCANSCI_TOOL_CALL>", re.DOTALL)
_TOOL_CALL_PATTERN = re.compile(
    r"tool_call\s*\(\s*name\s*=\s*[\"'](?P<name>[A-Za-z0-9_-]+)[\"']"
    r"(?:\s*,\s*arguments\s*=\s*(?P<arguments>\{.*?\}))?\s*\)",
    re.DOTALL,
)
_DOI_IN_TEXT_PATTERN = re.compile(
    r"(?<![A-Za-z0-9])10\.\d{4,9}/[A-Za-z0-9._;()/:+-]+",
    re.IGNORECASE,
)
_ARXIV_IN_TEXT_PATTERN = re.compile(
    r"(?<![A-Za-z0-9])(?:arxiv:)?(?:\d{4}\.\d{4,5}|[a-z-]+/\d{7})(?:v\d+)?",
    re.IGNORECASE,
)
_SKILL_STATE_ID_PATTERN = re.compile(r"^[a-z0-9][a-z0-9._-]{0,99}$")
_SKILL_STATE_HASH_PATTERN = re.compile(r"^sha256:[0-9a-f]{64}$")
_MAX_SESSION_STATE_ENTRIES = 200_000
_MAX_SESSION_STATE_LINE_CHARS = 16_000_000
_MAX_PERSISTED_SKILL_RESOURCES = 512

# A read-only risk label alone does not prove thread safety.  This host-owned
# allow-list intentionally excludes composite retrieval, mutable caches,
# browser automation, control/loaders, MCP, and every effectful tool.
_PARALLEL_SAFE_TOOL_NAMES = frozenset({
    "inspect_available_tools",
    "zotero_status",
    "zotero_fulltext",
    "zotero_attachment",
    "zotero_export_bibtex",
    "zotero_citations",
    "obsidian_status",
    "obsidian_search",
    "obsidian_read",
    "obsidian_backlinks",
    "verify_doi",
    "search_web",
    "agent_reach",
    "search_journal",
    "audit_references",
    "list_scientific_agents",
    "collect_scientific_agents",
})
_TOOL_EXECUTOR_WORKERS = 4
_TOOL_EXECUTOR_CAPACITY = 16
_TOOL_COMPLETION_CAPACITY = 32
_DISPATCH_AUDIT_CAPACITY = 256
_CANCEL_DRAIN_SECONDS = 1.5


@dataclass(frozen=True)
class _ToolCompletion:
    request_id: str
    generation: int
    call_id: str
    name: str
    arguments: dict[str, Any]
    parallel_safe: bool = False
    result: Any = None
    error_type: str = ""
    error: str = ""


@dataclass(frozen=True)
class _PendingToolCall:
    call_id: str
    name: str
    arguments: dict[str, Any]
    parallel_safe: bool


@dataclass(frozen=True)
class _PreparedToolResult:
    temporary_path: Path
    final_path: Path
    reference: str
    encoded_bytes: int


@dataclass
class _RunContext:
    request_id: str
    session_id: str
    generation: int
    supports_images: bool = False
    cancel_event: threading.Event = field(default_factory=threading.Event)
    history: list[dict[str, Any]] = field(default_factory=list)
    history_lock: threading.Lock = field(default_factory=threading.Lock)
    completions: Queue[_ToolCompletion] = field(
        default_factory=lambda: Queue(maxsize=_TOOL_COMPLETION_CAPACITY)
    )
    in_flight: dict[str, Future[_ToolCompletion]] = field(default_factory=dict)
    deferred: deque[_PendingToolCall] = field(default_factory=deque)
    parallel_in_flight: int = 0
    sequential_in_flight: bool = False
    authorized_call_count: int = 0
    active: bool = True
    pending_terminal: dict[str, Any] | None = None
    terminal_tool_completed: str = ""
    lock: threading.Lock = field(default_factory=threading.Lock)
    # Linearizes the final tool.result wire commit against cancellation.  File
    # preparation and rename stay outside this lock so cancellation remains
    # responsive during slow disk I/O.
    result_commit_lock: threading.Lock = field(default_factory=threading.Lock)

    def has_pending_tools(self) -> bool:
        with self.lock:
            return bool(self.in_flight or self.deferred)


class _ProtocolDispatcher:
    """Route one sidecar stdout stream without letting waiters steal events."""

    def __init__(
        self,
        raw_output: Queue[str | None],
        *,
        generation: int,
        audit: deque[dict[str, Any]],
    ) -> None:
        self.raw_output = raw_output
        self.generation = generation
        self.audit = audit
        self._lock = threading.Lock()
        self._requests: dict[str, Queue[dict[str, Any] | None]] = {}
        self._commands: dict[str, Queue[dict[str, Any] | None]] = {}
        self._request_tombstones: set[str] = set()
        self._command_tombstones: set[str] = set()
        self._thread: threading.Thread | None = None

    def start(self) -> None:
        with self._lock:
            if self._thread is not None and self._thread.is_alive():
                return
            self._thread = threading.Thread(target=self._run, daemon=True, name="scansci-pi-dispatch")
            self._thread.start()

    def register_request(self, request_id: str) -> Queue[dict[str, Any] | None]:
        channel: Queue[dict[str, Any] | None] = Queue()
        with self._lock:
            if request_id in self._requests:
                raise RuntimeError(f"Pi request {request_id} is already registered")
            self._request_tombstones.discard(request_id)
            self._requests[request_id] = channel
        return channel

    def unregister_request(self, request_id: str) -> None:
        with self._lock:
            self._requests.pop(request_id, None)
            self._request_tombstones.add(request_id)
            if len(self._request_tombstones) > _DISPATCH_AUDIT_CAPACITY:
                self._request_tombstones.pop()

    def register_command(self, command_id: str) -> Queue[dict[str, Any] | None]:
        channel: Queue[dict[str, Any] | None] = Queue()
        with self._lock:
            if command_id in self._commands:
                raise RuntimeError(f"Pi command {command_id} is already registered")
            self._command_tombstones.discard(command_id)
            self._commands[command_id] = channel
        return channel

    def unregister_command(self, command_id: str) -> None:
        with self._lock:
            self._commands.pop(command_id, None)
            self._command_tombstones.add(command_id)
            if len(self._command_tombstones) > _DISPATCH_AUDIT_CAPACITY:
                self._command_tombstones.pop()

    def _record(self, kind: str, event: dict[str, Any] | None = None) -> None:
        record = {
            "kind": kind,
            "generation": self.generation,
            "type": str((event or {}).get("type", ""))[:80],
            "request_id": str((event or {}).get("request_id", ""))[:80],
            "command_id": str((event or {}).get("command_id", ""))[:80],
        }
        with self._lock:
            self.audit.append(record)

    def _broadcast_eof(self) -> None:
        with self._lock:
            channels = [*self._requests.values(), *self._commands.values()]
        for channel in channels:
            channel.put(None)

    def _run(self) -> None:
        while True:
            raw = self.raw_output.get()
            if raw is None:
                self._record("eof")
                self._broadcast_eof()
                return
            try:
                event = json.loads(raw)
            except (json.JSONDecodeError, TypeError):
                self._record("invalid_json")
                with self._lock:
                    channels = [*self._requests.values(), *self._commands.values()]
                failure = {"type": "protocol.error", "error": "Pi Agent returned invalid JSON"}
                for channel in channels:
                    channel.put(dict(failure))
                continue
            if not isinstance(event, dict):
                self._record("invalid_message")
                continue
            request_id = str(event.get("request_id", ""))
            command_id = str(event.get("command_id", ""))
            event_generation = event.get("generation")
            if command_id and (
                event_generation is None
                or not str(event_generation).isdigit()
                or int(event_generation) != self.generation
            ):
                self._record("cross_generation", event)
                continue
            delivered = False
            with self._lock:
                request_channel = self._requests.get(request_id) if request_id else None
                command_channel = self._commands.get(command_id) if command_id else None
                sole_request_channel = (
                    next(iter(self._requests.values()))
                    if len(self._requests) == 1
                    else None
                )
                late_request = bool(request_id and request_id in self._request_tombstones)
                late_command = bool(command_id and command_id in self._command_tombstones)
            # Correlated command acknowledgements take their command route.  If
            # an event intentionally carries both identities, fan it out.
            if command_channel is not None:
                command_channel.put(event)
                delivered = True
            if request_channel is not None:
                request_channel.put(event)
                delivered = True
            elif not request_id and not command_id and str(event.get("type", "")) == "protocol.error":
                # An uncorrelated protocol error is generation-fatal.  Fan it
                # out so no request or command silently waits until timeout.
                with self._lock:
                    channels = [*self._requests.values(), *self._commands.values()]
                for channel in channels:
                    channel.put(dict(event))
                delivered = bool(channels)
            elif (
                sole_request_channel is not None
                and str(event.get("type", "")) in {"tool.call", "skill.call"}
                and not late_request
            ):
                # Executable messages with a missing/unknown request identity
                # must reach the sole active host gate so it can return an
                # explicit denial instead of leaving the sidecar promise hung.
                sole_request_channel.put(event)
                self._record("unauthorized_executable", event)
                delivered = True
            if not delivered:
                self._record("late" if late_request or late_command else "orphan", event)


def _persisted_skill_state(session_file: str | Path) -> dict[str, Any]:
    """Read only bounded Skill metadata from the active Pi JSONL branch."""

    path = Path(session_file) if str(session_file or "").strip() else None
    if path is None or path.suffix.lower() != ".jsonl" or not path.is_file():
        return {"schema": SKILL_STATE_SCHEMA, "loaded": []}
    entries: dict[str, tuple[str | None, dict[str, Any] | None]] = {}
    leaf_id = ""
    try:
        with path.open("r", encoding="utf-8") as handle:
            for index, line in enumerate(handle):
                if index >= _MAX_SESSION_STATE_ENTRIES or len(line) > _MAX_SESSION_STATE_LINE_CHARS:
                    return {"schema": SKILL_STATE_SCHEMA, "loaded": []}
                raw = json.loads(line)
                if not isinstance(raw, dict) or raw.get("type") == "session":
                    continue
                entry_id = str(raw.get("id", "") or "")
                parent_value = raw.get("parentId")
                parent_id = None if parent_value is None else str(parent_value)
                if (
                    not entry_id
                    or len(entry_id) > 128
                    or (parent_id is not None and len(parent_id) > 128)
                ):
                    return {"schema": SKILL_STATE_SCHEMA, "loaded": []}
                custom = None
                if raw.get("type") == "custom" and raw.get("customType") == SKILL_STATE_SCHEMA:
                    data = raw.get("data")
                    custom = dict(data) if isinstance(data, dict) else None
                entries[entry_id] = (parent_id, custom)
                leaf_id = entry_id
    except (OSError, UnicodeError, ValueError, TypeError):
        return {"schema": SKILL_STATE_SCHEMA, "loaded": []}

    loaded_by_key: dict[str, dict[str, Any]] = {}
    visited: set[str] = set()
    current = leaf_id
    while current:
        if current in visited or current not in entries:
            return {"schema": SKILL_STATE_SCHEMA, "loaded": []}
        visited.add(current)
        parent_id, custom = entries[current]
        if custom is not None:
            skill_id = str(custom.get("skill_id", "") or "").strip().lower()
            resource = str(custom.get("resource", "") or "")
            package_hash = str(custom.get("package_hash", "") or "")
            content_hash = str(custom.get("content_hash", "") or "")
            provenance = str(custom.get("provenance", "resume") or "resume")[:32]
            try:
                byte_count = int(custom.get("bytes", -1))
            except (TypeError, ValueError):
                byte_count = -1
            if (
                _SKILL_STATE_ID_PATTERN.fullmatch(skill_id)
                and 0 < len(resource) <= 500
                and _SKILL_STATE_HASH_PATTERN.fullmatch(package_hash)
                and _SKILL_STATE_HASH_PATTERN.fullmatch(content_hash)
                and 0 <= byte_count <= 64 * 1024
            ):
                key = f"{skill_id}:{resource}"
                if key not in loaded_by_key and len(loaded_by_key) < _MAX_PERSISTED_SKILL_RESOURCES:
                    loaded_by_key[key] = {
                        "skill_id": skill_id,
                        "resource": resource,
                        "package_hash": package_hash,
                        "content_hash": content_hash,
                        "provenance": provenance,
                        "bytes": byte_count,
                    }
        current = parent_id or ""
    return {
        "schema": SKILL_STATE_SCHEMA,
        "loaded": list(loaded_by_key.values()),
    }


def _explicit_paper_identifiers(text: str) -> list[str]:
    """Extract exact identifiers without turning a DOI into a fuzzy query."""

    identifiers: list[str] = []
    for pattern in (_DOI_IN_TEXT_PATTERN, _ARXIV_IN_TEXT_PATTERN):
        for match in pattern.finditer(str(text or "")):
            value = match.group(0).rstrip(".,;:!?，。；：！？、")
            value = re.sub(r"^https?://(?:dx\.)?doi\.org/", "", value, flags=re.IGNORECASE)
            value = re.sub(r"^doi:\s*", "", value, flags=re.IGNORECASE)
            if value:
                identifiers.append(value)
    return list(dict.fromkeys(identifiers))


def _presentation_arguments_from_result(result: dict[str, Any], user_request: str) -> dict[str, Any]:
    """Turn a retrieved/document-map result into a small verified PPTX input."""

    query = re.sub(r"\s+", " ", str(result.get("query", "") or "")).strip()
    title = f"{query[:64]} 文献综述" if query else "ScanSci 文献综述"

    candidates: list[dict[str, Any]] = []
    for key in ("documents", "hits", "items", "results"):
        candidates.extend(
            dict(item)
            for item in list(result.get(key, []) or [])
            if isinstance(item, dict)
        )
    for container_key in ("libraries", "zotero", "obsidian"):
        for library in list(result.get(container_key, []) or []):
            if not isinstance(library, dict):
                continue
            for key in ("documents", "hits", "items", "results", "citations"):
                candidates.extend(
                    dict(item)
                    for item in list(library.get(key, []) or [])
                    if isinstance(item, dict)
                )
    for library in list(result.get("libraries", []) or []):
        if not isinstance(library, dict):
            continue
        reader_answer = dict(library.get("reader_answer", {}) or {})
        candidates.extend(
            dict(item)
            for item in list(reader_answer.get("citations", []) or [])
            if isinstance(item, dict)
        )

    slides: list[dict[str, Any]] = [
        {
            "title": "研究任务",
            "bullets": [
                re.sub(r"\s+", " ", user_request).strip()[:320],
                f"已读取 {len(candidates)} 条可用证据或文档映射",
            ],
        }
    ]
    for index, item in enumerate(candidates[:10], start=1):
        item_title = str(
            item.get("title")
            or item.get("name")
            or item.get("document_title")
            or f"来源 {index}"
        ).strip()
        bullets: list[str] = []
        for label, key in (
            ("研究问题", "research_question"),
            ("方法", "methods"),
            ("主要发现", "findings"),
            ("局限", "limitations"),
            ("摘要", "abstract"),
            ("证据", "excerpt"),
            ("全文证据", "fulltext_excerpt"),
            ("内容", "text"),
            ("总结", "summary"),
        ):
            value = re.sub(r"\s+", " ", str(item.get(key, "") or "")).strip()
            if value:
                bullets.append(f"{label}：{value[:300]}")
            if len(bullets) >= 4:
                break
        if not bullets:
            bullets.append("该来源已纳入检索结果，但未返回可用于幻灯片的全文摘录。")
        slides.append({"title": item_title[:120], "bullets": bullets})
    if len(slides) == 1:
        slides.append(
            {
                "title": "证据状态",
                "bullets": ["当前工具结果没有返回可展开的来源条目；请在交付时明确这一限制。"],
            }
        )
    slides.append(
        {
            "title": "综合与证据边界",
            "bullets": [
                "仅总结工具实际返回的内容，不把题录或搜索片段冒充全文证据。",
                "未读取或失败的来源应在后续研究中单独补齐。",
            ],
        }
    )
    return {"title": title, "slides": slides}


class _ManagedGatewayAdapter:
    """Normalize ScanSci's text-only managed gateway to OpenAI tool-call SSE."""

    def __init__(self, *, upstream_base_url: str, api_key: str, minimum_interval_seconds: float = 0.0) -> None:
        self.upstream_base_url = upstream_base_url.rstrip("/")
        self.api_key = api_key
        self.minimum_interval_seconds = max(0.0, float(minimum_interval_seconds))
        self._rate_lock = threading.Lock()
        self._last_request_started = 0.0
        self.records: list[dict[str, Any]] = []
        self._records_lock = threading.Lock()
        adapter = self

        class Handler(BaseHTTPRequestHandler):
            def do_POST(self) -> None:  # noqa: N802 - stdlib handler contract
                adapter._handle(self)

            def log_message(self, _format: str, *_args: object) -> None:
                return

        self.server = ThreadingHTTPServer(("127.0.0.1", 0), Handler)
        self.thread = threading.Thread(target=self.server.serve_forever, daemon=True)
        self.thread.start()

    @property
    def base_url(self) -> str:
        return f"http://127.0.0.1:{self.server.server_port}/v1"

    def close(self) -> None:
        self.server.shutdown()
        self.thread.join(timeout=3)
        self.server.server_close()

    def _handle(self, handler: BaseHTTPRequestHandler) -> None:
        started = time.perf_counter()
        request_summary: dict[str, Any] = {}
        try:
            length = int(handler.headers.get("Content-Length", "0"))
            if length < 0 or length > _MAX_GATEWAY_REQUEST_BYTES:
                self._write_json(
                    handler,
                    413,
                    json.dumps(
                        {"error": {"message": "Managed gateway request exceeded the local input byte limit."}},
                        ensure_ascii=False,
                    ),
                )
                return
            original_payload = json.loads(handler.rfile.read(length))
            expects_stream = bool(original_payload.get("stream", False))
            followup_tool_call = self._managed_followup_tool_call(original_payload)
            payload = self._normalize_payload(original_payload)
            request_summary = {
                "model": str(payload.get("model", "")),
                "message_count": len(list(payload.get("messages", []) or [])),
                "tool_count": len(list(payload.get("tools", []) or [])),
                "stream_requested": expects_stream,
                "request_bytes": length,
            }
            if followup_tool_call is not None:
                self._record_transport(
                    request_summary,
                    status=200,
                    started=started,
                    usage={},
                )
                self._write_openai_response(
                    handler,
                    {
                        "id": f"chatcmpl-scansci-orchestrator-{uuid4().hex}",
                        "model": str(payload.get("model", "")),
                        "choices": [
                            {
                                "index": 0,
                                "message": {
                                    "role": "assistant",
                                    "content": None,
                                    "tool_calls": [followup_tool_call],
                                },
                                "finish_reason": "tool_calls",
                            }
                        ],
                        "usage": {"prompt_tokens": 0, "completion_tokens": 0, "total_tokens": 0},
                    },
                    stream=expects_stream,
                )
                return
            self._await_rate_slot()
            response = self._post_with_transport_retry(payload)
            response_content = bytes(getattr(response, "content", b"") or b"")
            if len(response_content) > _MAX_GATEWAY_RESPONSE_BYTES:
                raise RuntimeError("Managed gateway response exceeded the local output byte limit")
            if response.status_code >= 400:
                rate_type = str(response.headers.get("x-ratelimit-type", ""))
                retry_after = str(response.headers.get("Retry-After", ""))
                self._record_transport(
                    request_summary,
                    status=response.status_code,
                    started=started,
                    error=response.text[:500],
                    response_headers={
                        key: value
                        for key, value in response.headers.items()
                        if "rate" in key.lower() or "retry" in key.lower()
                    },
                )
                error_body = response.text
                lowered_error = response.text.casefold()
                if (
                    response.status_code == 400
                    and length > 0
                    and "exceeds the managed service input limit" in lowered_error
                ):
                    error_body = json.dumps(
                        {
                            "error": {
                                "message": (
                                    "context_length_exceeded: The managed service input limit was exceeded. "
                                    "Compact the conversation or reduce tool output before retrying."
                                ),
                                "code": "context_length_exceeded",
                            }
                        },
                        ensure_ascii=False,
                    )
                elif rate_type or retry_after:
                    try:
                        error_payload = response.json()
                        message = str(dict(error_payload.get("error", {}) or {}).get("message", response.text))
                    except (ValueError, TypeError):
                        message = response.text.strip()
                    error_body = json.dumps(
                        {
                            "error": {
                                "message": f"{message} [rate_limit_type={rate_type}; retry_after={retry_after}]",
                                "code": "rate_limit_exceeded",
                            }
                        },
                        ensure_ascii=False,
                    )
                self._write_json(handler, response.status_code, error_body)
                return
            body = response.json()
            self._record_transport(
                request_summary,
                status=response.status_code,
                started=started,
                usage=dict(body.get("usage", {}) or {}),
            )
            self._write_openai_response(handler, body, stream=expects_stream)
        except Exception as error:  # noqa: BLE001 - protocol adapter must return an HTTP error
            self._record_transport(request_summary, status=502, started=started, error=f"{type(error).__name__}: {error}")
            self._write_json(
                handler,
                502,
                json.dumps({"error": {"message": f"Managed gateway adapter failed: {type(error).__name__}: {error}"}}),
            )

    def _post_with_transport_retry(self, payload: dict[str, Any]) -> requests.Response:
        """Retry transient gateway failures before Pi sees an empty turn.

        The managed gateway forwards the upstream provider's ``Retry-After``
        header.  Honouring a bounded version of that signal is important: an
        eight-second blind retry simply turns a brief upstream quota window
        into a visible failed turn.  The handler remains bounded so one user
        request can never occupy a worker indefinitely.
        """

        retryable_statuses = {429, 500, 502, 503, 504}
        last_error: Exception | None = None
        for attempt in range(_GATEWAY_TRANSPORT_ATTEMPTS):
            try:
                response = requests.post(
                    f"{self.upstream_base_url}/chat/completions",
                    headers={
                        "Authorization": f"Bearer {self.api_key}",
                        "Content-Type": "application/json",
                        "Accept": "application/json",
                        "User-Agent": "OpenAI/Python 2.8.1",
                    },
                    json=payload,
                    timeout=(15, 150),
                )
                if response.status_code not in retryable_statuses or attempt == _GATEWAY_TRANSPORT_ATTEMPTS - 1:
                    return response
                retry_after = response.headers.get("Retry-After", "")
                try:
                    delay = (
                        min(_GATEWAY_MAX_RETRY_AFTER_SECONDS, max(0.5, float(retry_after)))
                        if retry_after
                        else float(2**attempt)
                    )
                except ValueError:
                    delay = float(2**attempt)
            except requests.RequestException as error:
                last_error = error
                if attempt == _GATEWAY_TRANSPORT_ATTEMPTS - 1:
                    raise
                delay = float(2**attempt)
            time.sleep(delay)
        if last_error is not None:  # pragma: no cover - loop always returns or raises
            raise last_error
        raise RuntimeError("Managed gateway retry loop ended unexpectedly")

    def _await_rate_slot(self) -> None:
        if self.minimum_interval_seconds <= 0:
            return
        with self._rate_lock:
            delay = self.minimum_interval_seconds - (time.monotonic() - self._last_request_started)
            if delay > 0:
                time.sleep(delay)
            self._last_request_started = time.monotonic()

    def _record_transport(
        self,
        request_summary: dict[str, Any],
        *,
        status: int,
        started: float,
        usage: dict[str, Any] | None = None,
        error: str = "",
        response_headers: dict[str, str] | None = None,
    ) -> None:
        record = {
            **request_summary,
            "status": status,
            "latency_seconds": round(time.perf_counter() - started, 6),
            "usage": usage or {},
            "error": str(_redact_tool_value(error))[:1_000],
            "response_headers": response_headers or {},
        }
        with self._records_lock:
            self.records.append(record)

    @staticmethod
    def _managed_followup_tool_call(payload: dict[str, Any]) -> dict[str, Any] | None:
        """Deterministically advance explicit download/index/synthesis workflows."""

        messages = [dict(item) for item in list(payload.get("messages", []) or []) if isinstance(item, dict)]
        tools = {
            str(dict(item.get("function", {}) or {}).get("name", ""))
            for item in list(payload.get("tools", []) or [])
            if isinstance(item, dict)
        }
        tool_names_by_id: dict[str, str] = {}
        for message in messages:
            if message.get("role") != "assistant":
                continue
            for call in list(message.get("tool_calls", []) or []):
                if not isinstance(call, dict):
                    continue
                call_id = str(call.get("id", ""))
                name = str(dict(call.get("function", {}) or {}).get("name", ""))
                if call_id and name:
                    tool_names_by_id[call_id] = name

        completed: list[tuple[str, dict[str, Any]]] = []
        for message in messages:
            if message.get("role") != "tool":
                continue
            name = tool_names_by_id.get(str(message.get("tool_call_id", "")), "")
            if not name:
                continue
            content = message.get("content", "")
            if isinstance(content, list):
                content = "".join(
                    str(part.get("text", ""))
                    for part in content
                    if isinstance(part, dict) and part.get("type") == "text"
                )
            try:
                result = json.loads(str(content or "{}"))
            except (TypeError, ValueError, json.JSONDecodeError):
                # Successful ScanSci bridge tools always serialize a JSON
                # object.  Non-JSON tool content is an SDK error message and
                # must never advance the deterministic workflow.
                return None
            if bool(message.get("is_error", message.get("isError", False))) or (
                isinstance(result, dict) and result.get("ok") is False
            ):
                return None
            completed.append((name, dict(result) if isinstance(result, dict) else {}))
        wrapped_user_requests = []
        direct_user_requests = []
        for message in messages:
            if message.get("role") != "user":
                continue
            content = str(message.get("content", "")).strip()
            if "[USER]\n" in content:
                wrapped_user_requests.append(content.rsplit("[USER]\n", 1)[-1].strip())
            elif content:
                direct_user_requests.append(content)
        user_request = str(
            (wrapped_user_requests[-1:] or direct_user_requests[:1] or [""])[0]
        ).casefold()
        wants_download = any(
            marker in user_request
            for marker in ("download", "acquire", "full text", "下载", "全文", "获取论文", "获取文献", "索引")
        )
        wants_presentation = (
            "create_presentation" in tools
            and any(marker in user_request for marker in ("ppt", "powerpoint", "slide", "幻灯片", "演示文稿"))
            and any(
                marker in user_request
                for marker in ("create", "generate", "make", "build", "actual", "downloadable", "创建", "生成", "制作", "实际", "可下载")
            )
        )
        if not completed:
            explicit_identifiers = _explicit_paper_identifiers(user_request)
            if wants_download and explicit_identifiers and "download_and_index" in tools:
                return {
                    "id": f"call_scansci_orchestrated_{uuid4().hex[:16]}",
                    "type": "function",
                    "function": {
                        "name": "download_and_index",
                        "arguments": json.dumps(
                            {"identifiers": explicit_identifiers},
                            ensure_ascii=False,
                            separators=(",", ":"),
                        ),
                    },
                }
            task_document_reference = any(
                marker in user_request
                for marker in ("刚才", "这些文献", "这些论文", "当前任务", "下载的文献", "下载的论文", "task documents")
            )
            wants_document_summary = any(
                marker in user_request
                for marker in ("总结", "比较", "对比", "分析", "方法", "结论", "summar", "compar", "analy")
            )
            if (
                task_document_reference
                and wants_document_summary
                and "summarize_documents" in tools
                and "read_task_documents" in tools
                and "discover_papers" not in tools
            ):
                return {
                    "id": f"call_scansci_orchestrated_{uuid4().hex[:16]}",
                    "type": "function",
                    "function": {
                        "name": "summarize_documents",
                        "arguments": "{}",
                    },
                }
            return None

        wants_synthesis = any(
            marker in user_request
            for marker in (
                "summar",
                "synthesi",
                "review",
                "compar",
                "analy",
                "method",
                "conclusion",
                "总结",
                "综述",
                "比较",
                "对比",
                "分析",
                "归纳",
                "方法",
                "结论",
                "共同点",
            )
        ) or wants_presentation
        last_name, last_result = completed[-1]
        completed_names = {name for name, _ in completed}
        next_name = ""
        arguments: dict[str, Any] = {}
        if (
            last_name == "discover_papers"
            and wants_download
            and "download_and_index" in tools
            and "download_and_index" not in completed_names
        ):
            requested_count = re.search(r"(\d{1,2})\s*(?:篇|papers?)", user_request, re.IGNORECASE)
            download_limit = min(8, max(1, int(requested_count.group(1)))) if requested_count else 3
            identifiers = [
                str(value).strip()
                for value in list(last_result.get("download_identifiers", []) or [])
                if str(value).strip()
            ][:download_limit]
            if identifiers:
                next_name = "download_and_index"
                arguments = {"identifiers": identifiers}
        elif (
            last_name == "download_and_index"
            and wants_synthesis
            and "summarize_documents" in tools
            and "summarize_documents" not in completed_names
        ):
            next_name = "summarize_documents"
            run_id = str(last_result.get("run_id", "")).strip()
            arguments = {"run_id": run_id} if run_id else {}
        elif (
            last_name in {
                "summarize_documents",
                "kb_search",
                "search_local_evidence",
                "zotero_search",
                "zotero_fulltext",
                "obsidian_search",
                "obsidian_read",
                "build_verified_answer",
            }
            and wants_presentation
            and "create_presentation" in tools
            and "create_presentation" not in completed_names
        ):
            next_name = "create_presentation"
            arguments = _presentation_arguments_from_result(last_result, user_request)
        elif (
            last_name == "create_presentation"
            and "download_and_index" in completed_names
            and "check_task_completion" in tools
            and "check_task_completion" not in completed_names
        ):
            next_name = "check_task_completion"
            run_id = next(
                (
                    str(result.get("run_id", "")).strip()
                    for name, result in reversed(completed)
                    if name == "download_and_index" and str(result.get("run_id", "")).strip()
                ),
                "",
            )
            arguments = {"run_id": run_id} if run_id else {}
        elif (
            last_name in {"download_and_index", "summarize_documents"}
            and "check_task_completion" in tools
            and "check_task_completion" not in completed_names
            and not (last_name == "summarize_documents" and wants_presentation)
        ):
            next_name = "check_task_completion"
            run_id = str(last_result.get("run_id", "")).strip()
            if not run_id:
                run_id = next(
                    (
                        str(result.get("run_id", "")).strip()
                        for name, result in reversed(completed)
                        if name == "download_and_index" and str(result.get("run_id", "")).strip()
                    ),
                    "",
                )
            arguments = {"run_id": run_id} if run_id else {}
        if not next_name:
            return None
        return {
            "id": f"call_scansci_orchestrated_{uuid4().hex[:16]}",
            "type": "function",
            "function": {
                "name": next_name,
                "arguments": json.dumps(arguments, ensure_ascii=False, separators=(",", ":")),
            },
        }

    @staticmethod
    def _normalize_payload(payload: dict[str, Any]) -> dict[str, Any]:
        normalized = dict(payload)
        normalized["stream"] = False
        normalized.pop("stream_options", None)
        messages: list[dict[str, Any]] = []
        for item in list(normalized.get("messages", []) or []):
            message = dict(item)
            content = message.get("content", "")
            if isinstance(content, list):
                content = "".join(
                    str(part.get("text", ""))
                    for part in content
                    if isinstance(part, dict) and part.get("type") == "text"
                )
            message["content"] = "" if content is None else str(content)
            messages.append(message)
        functions = []
        for tool in list(normalized.get("tools", []) or []):
            function = dict(tool.get("function", {}) or {}) if isinstance(tool, dict) else {}
            if function.get("name"):
                functions.append(function)
        has_tool_result = any(message.get("role") == "tool" for message in messages)
        completed_tool_names = {
            str(call.get("function", {}).get("name", ""))
            for message in messages
            if message.get("role") == "assistant"
            for call in list(message.get("tool_calls", []) or [])
            if isinstance(call, dict) and isinstance(call.get("function"), dict)
        }
        if len(functions) == 1 and not has_tool_result:
            function = functions[0]
            name = str(function["name"])
            properties = dict(dict(function.get("parameters", {}) or {}).get("properties", {}) or {})
            required = list(dict(function.get("parameters", {}) or {}).get("required", []) or [])
            final_user_value = _final_user_value(messages)
            example_arguments = {
                str(key): final_user_value if str(key) in {"question", "query"} and final_user_value else "COPY THE USER VALUE VERBATIM"
                for key in required
                if key in properties
            }
            exact_call = json.dumps({"name": name, "arguments": example_arguments}, separators=(",", ":"))
            instruction = (
                "A tool call is mandatory for this request. Reply only with this exact JSON object: "
                f"{exact_call}. Do not add XML, Markdown, explanation, or any other characters. "
                "The key is name (never e) and arguments must remain a JSON object."
            )
            # The managed gateway exposes text completion rather than native
            # function calling.  A small deterministic generation window is
            # enough for the model to select and encode this single action and
            # prevents repetition loops from consuming the model's full limit.
            normalized["max_tokens"] = min(int(normalized.get("max_tokens", 192) or 192), 192)
            normalized["temperature"] = 0
            normalized["thinking"] = {"type": "disabled"}
        elif has_tool_result and "build_verified_answer" in completed_tool_names:
            # build_verified_answer is terminal for the evidence endpoint: its
            # verified payload is delivered by the bridge, while Pi still gets
            # a clean assistant turn to close and persist the SDK loop.
            instruction = "The tool completed successfully. Reply with DONE only."
            normalized["max_tokens"] = min(int(normalized.get("max_tokens", 16) or 16), 16)
            normalized["temperature"] = 0
            normalized["thinking"] = {"type": "disabled"}
            normalized["tools"] = []
            normalized.pop("tool_choice", None)
        elif has_tool_result:
            instruction = (
                "A ScanSci tool has returned real results. Use those results to answer the user's original request now. "
                "Include source URLs supplied by the tool, distinguish search snippets from verified page content, "
                "and do not reply with DONE, a plan, setup instructions, or a request to wait."
            )
            normalized["tools"] = []
            normalized.pop("tool_choice", None)
        else:
            instruction = (
                "When a supplied function is needed, emit exactly one tool intent as "
                '<SCANSCI_TOOL_CALL>{"name":"function_name","arguments":{}}</SCANSCI_TOOL_CALL>. '
                "Do not say that tools are unavailable. After a tool result, answer normally."
            )
        if messages and messages[0].get("role") == "system":
            messages[0]["content"] = f"{messages[0]['content']}\n\n{instruction}"
        else:
            messages.insert(0, {"role": "system", "content": instruction})
        normalized["messages"] = messages
        return normalized

    @staticmethod
    def _write_json(handler: BaseHTTPRequestHandler, status: int, body: str) -> None:
        encoded = body.encode("utf-8")
        handler.send_response(status)
        handler.send_header("Content-Type", "application/json; charset=utf-8")
        handler.send_header("Content-Length", str(len(encoded)))
        handler.end_headers()
        handler.wfile.write(encoded)

    @staticmethod
    def _write_openai_response(handler: BaseHTTPRequestHandler, body: dict[str, Any], *, stream: bool) -> None:
        choice = dict((body.get("choices") or [{}])[0] or {})
        message = dict(choice.get("message") or {})
        model = str(body.get("model", ""))
        completion_id = str(body.get("id", f"chatcmpl-scansci-{uuid4().hex}"))
        created = int(body.get("created", int(time.time())) or int(time.time()))
        tool_calls = list(message.get("tool_calls") or [])
        content = str(message.get("content") or "")
        if not tool_calls:
            parsed_calls: list[dict[str, Any]] = []
            for index, (name, arguments_value) in enumerate(_parse_text_tool_intents(content)):
                arguments = json.dumps(arguments_value, ensure_ascii=False, separators=(",", ":"))
                parsed_calls.append(
                    {
                        "index": index,
                        "id": f"call_scansci_{uuid4().hex[:16]}",
                        "type": "function",
                        "function": {"name": name, "arguments": arguments},
                    }
                )
            tool_calls = parsed_calls
        if tool_calls:
            delta = {"role": "assistant", "tool_calls": tool_calls}
            finish_reason = "tool_calls"
        else:
            cleaned = re.sub(r"\[(?:ASSISTANT|USER)\]\s*", "", content).strip()
            delta = {"role": "assistant", "content": cleaned}
            finish_reason = str(choice.get("finish_reason") or "stop")
        if not stream:
            response_body = dict(body)
            response_body["object"] = "chat.completion"
            response_body["choices"] = [
                {
                    "index": 0,
                    "message": delta,
                    "finish_reason": finish_reason,
                }
            ]
            _ManagedGatewayAdapter._write_json(handler, 200, json.dumps(response_body, ensure_ascii=False))
            return
        chunks = [
            {
                "id": completion_id,
                "object": "chat.completion.chunk",
                "created": created,
                "model": model,
                "choices": [{"index": 0, "delta": delta, "finish_reason": None}],
            },
            {
                "id": completion_id,
                "object": "chat.completion.chunk",
                "created": created,
                "model": model,
                "choices": [{"index": 0, "delta": {}, "finish_reason": finish_reason}],
                "usage": body.get("usage") or {},
            },
        ]
        encoded = ("".join(f"data: {json.dumps(chunk, ensure_ascii=False)}\n\n" for chunk in chunks) + "data: [DONE]\n\n").encode("utf-8")
        handler.send_response(200)
        handler.send_header("Content-Type", "text/event-stream; charset=utf-8")
        handler.send_header("Content-Length", str(len(encoded)))
        handler.end_headers()
        handler.wfile.write(encoded)


def _parse_text_tool_intents(content: str) -> list[tuple[str, dict[str, Any]]]:
    calls: list[tuple[str, dict[str, Any]]] = []
    stripped = content.strip()
    if stripped.startswith("```") and stripped.endswith("```"):
        stripped = re.sub(r"^```(?:json)?\s*|\s*```$", "", stripped, flags=re.IGNORECASE | re.DOTALL).strip()
    try:
        payload = json.loads(stripped)
    except json.JSONDecodeError:
        payload = None
    if isinstance(payload, dict) and re.fullmatch(r"[A-Za-z0-9_-]+", str(payload.get("name", ""))):
        arguments = _coerce_tool_arguments(payload.get("arguments", {}))
        if arguments is not None:
            return [(str(payload["name"]), arguments)]
    for match in _TOOL_TAG_PATTERN.finditer(content):
        try:
            payload = json.loads(match.group("body"))
        except json.JSONDecodeError:
            continue
        if not isinstance(payload, dict) or not re.fullmatch(r"[A-Za-z0-9_-]+", str(payload.get("name", ""))):
            continue
        arguments = _coerce_tool_arguments(payload.get("arguments", {}))
        if arguments is not None:
            calls.append((str(payload["name"]), arguments))
    if calls:
        return calls
    for match in _TOOL_CALL_PATTERN.finditer(content):
        try:
            arguments = json.loads(match.group("arguments") or "{}")
        except json.JSONDecodeError:
            arguments = {}
        calls.append((match.group("name"), dict(arguments) if isinstance(arguments, dict) else {}))
    return calls


def _coerce_tool_arguments(value: Any) -> dict[str, Any] | None:
    """Accept native argument objects and JSON-string encodings of objects."""

    if isinstance(value, dict):
        return dict(value)
    if isinstance(value, str):
        try:
            decoded = json.loads(value)
        except json.JSONDecodeError:
            return None
        if isinstance(decoded, dict):
            return dict(decoded)
    return None


def _final_user_value(messages: list[dict[str, Any]]) -> str:
    """Recover the final user value from Pi's text conversation wrapper."""

    for message in reversed(messages):
        if message.get("role") != "user":
            continue
        content = str(message.get("content", "")).strip()
        marker = "[USER]\n"
        if marker in content:
            content = content.rsplit(marker, 1)[-1].strip()
        return content
    return ""


class PiAgentClient:
    """Own a persistent Pi sidecar and bridge approved ScanSci tools."""

    def __init__(
        self,
        *,
        workspace: str | Path,
        evidence_db: str | Path,
        root_evidence_db: str | Path | None = None,
        additional_evidence_dbs: list[str | Path] | None = None,
        notebook_ids: list[str] | None = None,
        knowledge_scope: dict[str, Any] | None = None,
        embedding_provider: Any | None = None,
        reranker: Any | None = None,
        active_run_id: str = "",
        scientific_agent_control: Callable[[str, str, dict[str, Any]], dict[str, Any]] | None = None,
    ) -> None:
        self.workspace = Path(workspace).resolve()
        self.evidence_db = Path(evidence_db).resolve()
        # ``evidence_db`` can point at an isolated notebook index.  Keep the
        # workspace-level root as a separate value so high-level tools that
        # import new files derive the same notebook index path as the rest of
        # the desktop runtime instead of nesting a second ``.libraries`` tree.
        self.root_evidence_db = Path(root_evidence_db or evidence_db).resolve()
        self.evidence_dbs = list(dict.fromkeys([
            self.evidence_db,
            *(Path(path).resolve() for path in (additional_evidence_dbs or [])),
        ]))
        self.notebook_ids = list(dict.fromkeys(str(value).strip() for value in (notebook_ids or []) if str(value).strip()))
        self.knowledge_scope = dict(knowledge_scope or {})
        self._knowledge_scope_cache: dict[str, dict[str, Any]] = {}
        self.embedding_provider = embedding_provider
        self.reranker = reranker
        self.active_run_id = str(active_run_id or "").strip()
        self._scientific_agent_control = scientific_agent_control
        self.agent_dir = self.workspace.parent / ".scansci-pi-agent"
        self.agent_dir.mkdir(parents=True, exist_ok=True)
        self._process: subprocess.Popen[str] | None = None
        self._output: Queue[str | None] = Queue()
        self._process_generation = 0
        self._protocol_generation = 0
        self._dispatcher: _ProtocolDispatcher | None = None
        self._dispatch_audit: deque[dict[str, Any]] = deque(maxlen=_DISPATCH_AUDIT_CAPACITY)
        self._errors: list[str] = []
        self._stdin_lock = threading.Lock()
        # Legacy direct-tool history remains available outside a run.  Active
        # Pi turns use the request-scoped history stored in _RunContext.
        self._tool_history: list[dict[str, Any]] = []
        self._tool_history_lock = threading.Lock()
        self._worker_context = threading.local()
        self._tool_executor = ThreadPoolExecutor(
            max_workers=_TOOL_EXECUTOR_WORKERS,
            thread_name_prefix="scansci-pi-tool",
        )
        self._tool_slots = threading.BoundedSemaphore(_TOOL_EXECUTOR_CAPACITY)
        self._run_lock = threading.Lock()
        # Own process/session lifecycle transitions for the full duration of a
        # run or management command.  Control writes intentionally use only
        # _stdin_lock so another thread can still cancel an active run.
        self._lifecycle_lock = threading.RLock()
        self._cancel_requested = threading.Event()
        self._contexts_lock = threading.Lock()
        self._run_contexts: dict[str, _RunContext] = {}
        self._request_generation = 0
        self._session_command_locks_guard = threading.Lock()
        self._session_command_locks: dict[str, threading.Lock] = {}
        self._active_request_id = ""
        self._active_session_id = ""
        self._interaction_lock = threading.Lock()
        self._pending_interaction_kinds: dict[str, tuple[str, str]] = {}
        self._approval_tokens: dict[str, ApprovalToken] = {}
        self._skill_runtimes: dict[str, ProgressiveSkillRuntime] = {}
        self._provider_key_fingerprint = ""
        self._gateway_adapter: _ManagedGatewayAdapter | None = None
        self._gateway_adapter_signature = ""

    @staticmethod
    def runtime_paths() -> tuple[Path, Path]:
        if getattr(sys, "frozen", False):
            bundle_root = Path(getattr(sys, "_MEIPASS", Path(sys.executable).parent)).resolve()
            # A user-installed Node component takes precedence over the
            # bundled copy, so slim builds without an embedded Node still run
            # the Pi sidecar after one confirmed component install.
            managed_node = default_node_component().executable()
            node_candidates = [
                *( [managed_node] if managed_node is not None else [] ),
                bundle_root / "pi_runtime" / "node.exe",
                Path(sys.executable).parent / "pi_runtime" / "node.exe",
            ]
            script_candidates = [
                bundle_root / "pi_runtime" / "main.mjs",
                Path(sys.executable).parent / "pi_runtime" / "main.mjs",
            ]
        else:
            project_root = Path(__file__).resolve().parents[2]
            node_on_path = shutil.which("node")
            node_candidates = [Path(node_on_path)] if node_on_path else []
            script_candidates = [project_root / "pi-runtime" / "dist" / "main.mjs"]

        node_path = next((path for path in node_candidates if path.is_file()), None)
        script_path = next((path for path in script_candidates if path.is_file()), None)
        if node_path is None:
            raise PiRuntimeUnavailable("The ScanSci Pi Node runtime is unavailable")
        if script_path is None:
            raise PiRuntimeUnavailable("The ScanSci Pi sidecar bundle is unavailable")
        return node_path, script_path

    @staticmethod
    def _node_environment() -> dict[str, str]:
        """Return a clean environment for the Node sidecar and its children."""

        environment = dict(os.environ)
        # Python-only import paths can make Windows reject a nested stdio
        # child spawn when the release gate injects PYTHONPATH. The Node
        # runtime does not need this variable, so do not pass it downstream.
        for variable in ("PYTHONPATH", "PYTHONUTF8"):
            environment.pop(variable, None)
        return environment

    @classmethod
    def runtime_status(cls, *, timeout_seconds: float = 8.0) -> dict[str, Any]:
        node_path, script_path = cls.runtime_paths()
        creation_flags = getattr(subprocess, "CREATE_NO_WINDOW", 0)
        process = subprocess.Popen(
            [str(node_path), str(script_path)],
            env=cls._node_environment(),
            stdin=subprocess.PIPE,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            creationflags=creation_flags,
        )
        try:
            assert process.stdin is not None
            assert process.stdout is not None
            process.stdin.write(json.dumps({
                "type": "ping",
                "pi_protocol_version": _PI_PROTOCOL_VERSION,
                "required_features": list(_PI_REQUIRED_FEATURES),
            }) + "\n")
            process.stdin.flush()
            output: Queue[str | None] = Queue()

            def read_one() -> None:
                output.put(process.stdout.readline() or None)

            threading.Thread(target=read_one, daemon=True).start()
            try:
                line = output.get(timeout=timeout_seconds)
            except Empty as error:
                raise PiRuntimeUnavailable("The ScanSci Pi sidecar did not respond") from error
            if not line:
                raise PiRuntimeUnavailable("The ScanSci Pi sidecar exited before responding")
            response = json.loads(line)
            if response.get("type") != "pong":
                raise PiRuntimeUnavailable(f"Unexpected Pi sidecar response: {response.get('type', '')}")
            protocol = int(response.get("protocol", 0) or 0)
            capabilities = {str(value) for value in list(response.get("capabilities", []) or [])}
            missing_features = [feature for feature in _PI_REQUIRED_FEATURES if feature not in capabilities]
            if protocol != _PI_PROTOCOL_VERSION or missing_features:
                detail = f"protocol={protocol}, missing={','.join(missing_features) or 'none'}"
                raise PiRuntimeUnavailable(f"The ScanSci Pi sidecar protocol is incompatible ({detail})")
            return {
                "ready": True,
                "runtime": str(response.get("runtime", "pi")),
                "version": str(response.get("version", "")),
                "protocol": protocol,
                "capabilities": sorted(capabilities),
                "node": str(node_path),
                "sidecar": str(script_path),
            }
        finally:
            process.kill()
            process.wait(timeout=5)

    @classmethod
    def probe_mcp_server(
        cls,
        *,
        workspace: str | Path,
        server: dict[str, Any],
        timeout_seconds: float = 20.0,
    ) -> dict[str, Any]:
        """Connect one saved MCP server and return its discovered tool surface."""

        node_path, script_path = cls.runtime_paths()
        creation_flags = getattr(subprocess, "CREATE_NO_WINDOW", 0)
        process = subprocess.Popen(
            [str(node_path), str(script_path)],
            env=cls._node_environment(),
            cwd=Path(workspace).resolve().parent,
            stdin=subprocess.PIPE,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
            encoding="utf-8",
            errors="replace",
            creationflags=creation_flags,
        )
        output: Queue[str | None] = Queue()
        try:
            assert process.stdin is not None
            assert process.stdout is not None

            def drain_stdout() -> None:
                for line in process.stdout:
                    output.put(line)
                output.put(None)

            threading.Thread(target=drain_stdout, daemon=True).start()
            request_id = uuid4().hex
            process.stdin.write(
                json.dumps(
                    {
                        "type": "mcp.probe",
                        "request_id": request_id,
                        "cwd": str(Path(workspace).resolve().parent),
                        "mcp_servers": [{**dict(server), "enabled": True}],
                    },
                    ensure_ascii=False,
                )
                + "\n"
            )
            process.stdin.flush()
            deadline = time.monotonic() + max(1.0, float(timeout_seconds))
            diagnostics: list[dict[str, Any]] = []
            while time.monotonic() < deadline:
                try:
                    line = output.get(timeout=min(0.25, max(0.01, deadline - time.monotonic())))
                except Empty:
                    continue
                if not line:
                    break
                message = json.loads(line)
                if str(message.get("request_id", "")) != request_id:
                    continue
                if message.get("type") == "status.update":
                    diagnostics.append(dict(message))
                    continue
                if message.get("type") == "mcp.probe.failed":
                    raise PiRuntimeUnavailable(str(message.get("error", "MCP probe failed")))
                if message.get("type") == "mcp.probe.completed":
                    return {**dict(message), "diagnostics": diagnostics}
            stderr = process.stderr.read() if process.stderr is not None and process.poll() is not None else ""
            detail = diagnostics[-1].get("error", "") if diagnostics else stderr.strip()
            raise PiRuntimeUnavailable(f"MCP connection test timed out{f': {detail}' if detail else ''}")
        finally:
            if process.poll() is None:
                process.kill()
            process.wait(timeout=5)

    def stream_chat(
        self,
        *,
        provider_kind: str,
        base_url: str,
        api_key: str,
        model_id: str,
        api_surface: str = "chat_completions",
        responses_enabled: bool = False,
        messages: list[dict[str, Any]],
        thinking_level: str = "medium",
        task_mode: str = "general",
        task_contract: dict[str, Any] | None = None,
        selected_skills: list[dict[str, Any]] | None = None,
        images: list[dict[str, Any]] | None = None,
        model_runtime: dict[str, Any] | ModelRuntimeDescriptor | None = None,
        timeout_seconds: float = 900.0,
        session_id: str | None = None,
    ) -> Iterator[dict[str, Any]]:
        """Yield normalized Pi events, optionally continuing a durable session."""

        if isinstance(model_runtime, ModelRuntimeDescriptor):
            descriptor = model_runtime
        elif isinstance(model_runtime, dict):
            descriptor = ModelRuntimeDescriptor.from_payload(model_runtime)
        else:
            descriptor = descriptor_from_model_record(
                provider_id="runtime",
                provider_kind=provider_kind,
                model_id=model_id,
                model_record=None,
                api_surface=str(api_surface or "chat_completions"),
            )
        if (
            descriptor.provider_kind != str(provider_kind)
            or descriptor.model_id != str(model_id)
            or descriptor.api_surface != str(api_surface or "chat_completions")
        ):
            raise ValueError(
                "Model runtime descriptor does not match the final provider route"
            )
        validated_images = validate_pi_image_blocks(list(images or []))
        if validated_images and not descriptor.supports_images:
            raise PiMultimodalUnavailable(
                "The selected model runtime descriptor does not support image input"
            )

        messages, stale_context_report = prune_stale_tool_results(messages)
        system_parts: list[str] = []
        conversation: list[str] = []
        for item in messages:
            role = str(item.get("role", "user")).strip().lower()
            content = item.get("content", "")
            if not isinstance(content, str):
                if role in {"tool", "toolresult", "tool_result"}:
                    content = json.dumps(content, ensure_ascii=False, default=str)
                else:
                    raise ValueError("Pi text bridge does not accept image message blocks")
            if role == "system":
                system_parts.append(content)
            else:
                conversation.append(f"[{role.upper()}]\n{content}")
        if not conversation:
            raise ValueError("Pi requires at least one conversational message")

        request_id = uuid4().hex
        transient_session = not session_id
        logical_session_id = session_id or uuid4().hex
        registry = self._load_session_registry()
        session_file = str(registry.get(logical_session_id, ""))
        is_recovery = bool(session_file and Path(session_file).is_file())
        if is_recovery:
            final_user = next(
                (
                    str(item.get("content", ""))
                    for item in reversed(messages)
                    if str(item.get("role", "user")).strip().lower() == "user"
                ),
                "",
            )
            prompt = final_user or "Continue the persisted ScanSci session."
        else:
            prompt = (
                "Continue the following ScanSci conversation. Reply to the final USER message.\n\n"
                + "\n\n".join(conversation)
            )
        final_request = next(
            (
                str(item.get("content", ""))
                for item in reversed(messages)
                if str(item.get("role", "user")).strip().lower() == "user"
            ),
            prompt,
        )
        if task_contract is None:
            # A direct transport caller still goes through the trusted host
            # compiler.  In contrast, a supplied payload that omits
            # ``allowed_tools`` remains omitted and therefore fail-closed.
            effective_contract = compile_host_task_contract(
                task_mode=task_mode,
                user_text=final_request,
            )
        else:
            effective_contract = TaskContract.from_payload(
                task_contract,
                request=final_request,
                task_mode=task_mode,
            ).to_dict()
        messages, token_envelope_report = build_token_envelope(
            messages,
            descriptor=descriptor,
            host_contract=effective_contract,
        )
        system_parts = []
        conversation = []
        for item in messages:
            role = str(item.get("role", "user")).strip().lower()
            content = item.get("content", "")
            if not isinstance(content, str):
                if role in {"tool", "toolresult", "tool_result"}:
                    content = json.dumps(content, ensure_ascii=False, default=str)
                else:
                    raise ValueError("Pi conversation text must be separate from canonical image blocks")
            if role == "system":
                system_parts.append(content)
            else:
                conversation.append(f"[{role.upper()}]\n{content}")
        if not conversation:
            raise ValueError("Pi requires at least one conversational message")
        if is_recovery:
            prompt = final_request or "Continue the persisted ScanSci session."
        else:
            prompt = (
                "Continue the following ScanSci conversation. Reply to the final USER message.\n\n"
                + "\n\n".join(conversation)
            )
        context_report = {
            **stale_context_report.to_dict(),
            "token_envelope": token_envelope_report.to_dict(),
            "model_runtime": {
                "schema_version": descriptor.schema_version,
                "context_window_tokens": descriptor.context_window_tokens,
                "provider_input_tokens": descriptor.provider_input_tokens,
                "input_modalities": list(descriptor.input_modalities),
                "degraded": descriptor.degraded,
                "degradation_reasons": list(descriptor.degradation_reasons),
            },
        }
        effective_base_url = self._effective_base_url(base_url=base_url, api_key=api_key)
        skill_selection = [
            dict(item)
            for item in list(selected_skills or [])
            if isinstance(item, dict) and str(item.get("id", "")).strip()
        ]
        persisted_skill_state = _persisted_skill_state(session_file) if is_recovery else {
            "schema": SKILL_STATE_SCHEMA,
            "loaded": [],
        }
        restored_by_key = {
            f"{str(item.get('skill_id', '')).lower()}:{str(item.get('resource', 'SKILL.md') or 'SKILL.md')}": dict(item)
            for item in list(persisted_skill_state.get("loaded", []) or [])
            if isinstance(item, dict) and str(item.get("skill_id", "")).strip()
        }
        for item in skill_selection:
            if str(item.get("status", "")) != "loaded":
                continue
            metadata = {
                "skill_id": str(item.get("id", "")),
                "resource": str(item.get("resource", "SKILL.md") or "SKILL.md"),
                "package_hash": str(item.get("package_hash", "")),
                "content_hash": str(item.get("content_hash", "")),
                "provenance": str(item.get("provenance", "explicit") or "explicit"),
                "bytes": int(item.get("bytes", 0) or 0),
            }
            restored_by_key[f"{metadata['skill_id'].lower()}:{metadata['resource']}"] = metadata
        restored_skill_state = {
            "schema": SKILL_STATE_SCHEMA,
            "loaded": list(restored_by_key.values()),
        }
        priority_ids = [
            str(item.get("id", ""))
            for item in skill_selection
            if str(item.get("provenance", "")) in {"explicit", "inferred"}
        ]
        priority_ids.extend(
            str(item.get("skill_id", ""))
            for item in list(persisted_skill_state.get("loaded", []) or [])
            if isinstance(item, dict)
        )
        skill_runtime = ProgressiveSkillRuntime(
            self.workspace,
            restored_state=restored_skill_state,
            priority_ids=priority_ids,
        )
        self._skill_runtimes[request_id] = skill_runtime
        registered_tools = sorted({
            "ask_user",
            "search_tools",
            "submit_plan",
            *[str(name) for name in list(effective_contract.get("allowed_tools", []) or []) if str(name)],
        })
        required_initial = {
            str(name)
            for group in list(effective_contract.get("required_tool_groups", []) or [])
            if isinstance(group, (list, tuple, set, frozenset))
            for name in group
            if str(name)
        }
        active_tools = sorted({
            "ask_user",
            "search_tools",
            "submit_plan",
            *[str(name) for name in list(effective_contract.get("initial_tools", []) or []) if str(name)],
            *required_initial,
        } & set(registered_tools))
        tool_set = {
            "mcp_servers": [str(item.get("id", "")) for item in self._enabled_mcp_servers() if isinstance(item, dict)],
            "disabled_tools": self._disabled_artifact_tools(),
            "registered_tools": registered_tools,
            "active_tools": active_tools,
        }
        contract_shape = {
            key: effective_contract.get(key)
            for key in (
                "schema_version", "output_format", "constraints", "required_evidence",
                "allowed_tools", "pause_policy", "success_criteria",
            )
        }
        contract_shape["model_runtime"] = descriptor.to_dict()
        prefix_shape = build_prefix_shape(
            provider=provider_kind,
            model=model_id,
            api_surface=str(api_surface or "chat_completions"),
            system_prompt="\n\n".join(system_parts),
            tool_set=tool_set,
            selected_skills=re.findall(r'<selected_skill\s+id="([^"]+)"', "\n\n".join(system_parts), flags=re.IGNORECASE),
            task_contract=contract_shape,
        )
        start_message = {
            "type": "run.start",
            "pi_protocol_version": _PI_PROTOCOL_VERSION,
            "required_features": list(_PI_REQUIRED_FEATURES),
            "request_id": request_id,
            "session_id": logical_session_id,
            "ephemeral_session": transient_session,
            "session_file": session_file,
            "cwd": str(self.workspace.parent),
            "agent_dir": str(self.agent_dir),
            "provider_kind": provider_kind,
            "base_url": effective_base_url,
            "model_id": model_id,
            "api_surface": str(api_surface or "chat_completions"),
            "responses_enabled": bool(responses_enabled),
            "thinking_level": thinking_level,
            "system_prompt": "\n\n".join(system_parts),
            "prompt": prompt,
            "images": validated_images,
            "model_runtime": descriptor.to_dict(),
            "task_mode": task_mode,
            "mcp_servers": self._enabled_mcp_servers(),
            "disabled_tools": self._disabled_artifact_tools(),
            "task_contract": effective_contract,
            "tool_set": tool_set,
            "prefix_shape": prefix_shape,
            "context_policy": context_report,
            "skill_catalog": skill_runtime.catalog(),
            "skill_selection": skill_selection,
            "skill_state": skill_runtime.state(),
        }
        manifest: RunManifest | None = None
        try:
            manifest = RunManifest.start(
                self.workspace,
                harness="pi",
                provider=provider_kind,
                model=model_id,
                api_surface=str(api_surface or "chat_completions"),
                session_id=logical_session_id,
                prompt=prompt,
                tool_set=tool_set,
                timeout_seconds=timeout_seconds,
                prefix_shape=prefix_shape,
                task_contract=effective_contract,
                context_policy=context_report,
            )
        except OSError:
            # Diagnostics are valuable but must never prevent a research run.
            manifest = None
        started_at = time.monotonic()
        first_visible_at: float | None = None
        retry_count = 0
        tool_calls = 0
        tool_failures = 0
        try:
            for event in self._run_request(start_message, api_key=api_key, timeout_seconds=timeout_seconds):
                event_type = str(event.get("type", ""))
                if event_type == "delta" and str(event.get("content", "")) and first_visible_at is None:
                    first_visible_at = time.monotonic()
                elif event_type == "retry":
                    retry_count += 1
                elif event_type == "tool.completed":
                    tool_calls += 1
                elif event_type == "tool.failed":
                    tool_calls += 1
                    tool_failures += 1
                if manifest is not None:
                    manifest.record(
                        f"agent.{event_type or 'event'}",
                        name=event.get("name", ""),
                        status=event.get("status", ""),
                    )
                    if event_type in {"done", "session_stats"}:
                        stats = event.get("stats") or event.get("value")
                        if isinstance(stats, dict):
                            runtime_prefix = stats.get("prefixShape")
                            manifest.record_context_stats(
                                stats,
                                prefix_shape=runtime_prefix if isinstance(runtime_prefix, dict) else prefix_shape,
                            )
                yield event
            if manifest is not None:
                manifest.finish(
                    status="completed",
                    duration_ms=round((time.monotonic() - started_at) * 1000, 2),
                    ttft_ms=round((first_visible_at - started_at) * 1000, 2) if first_visible_at else None,
                    retry_count=retry_count,
                    tool_calls=tool_calls,
                    tool_failures=tool_failures,
                )
        except BaseException as error:
            if manifest is not None:
                manifest.metric("duration_ms", round((time.monotonic() - started_at) * 1000, 2))
                manifest.fail(error, retryable=isinstance(error, (TimeoutError, ConnectionError)))
            raise
        finally:
            self._skill_runtimes.pop(request_id, None)
            if transient_session:
                self.close()

    def _execute_skill_tool(
        self,
        request_id: str,
        name: str,
        arguments: dict[str, Any],
    ) -> dict[str, Any]:
        """Execute one bounded instruction-plane operation for the current request."""

        runtime = self._skill_runtimes.get(str(request_id or ""))
        if runtime is None:
            raise PermissionError("Skill call belongs to an inactive request")
        if name == "search_skills":
            return runtime.search_skills(arguments.get("query", ""), arguments.get("limit", 8))
        if name == "load_skill":
            return runtime.load_skill(
                arguments.get("skill_id", ""),
                resource=arguments.get("resource"),
                provenance=arguments.get("provenance", "model"),
            )
        if name == "restore_skill":
            return runtime.restore_skill(
                arguments.get("skill_id", ""),
                resource=arguments.get("resource"),
            )
        raise PermissionError(f"Unknown Skill instruction operation: {name}")

    def _ensure_process(self, *, api_key: str) -> subprocess.Popen[str]:
        with self._lifecycle_lock:
            fingerprint = hashlib.sha256(api_key.encode("utf-8")).hexdigest()
            if self._process is not None and self._process.poll() is None:
                if self._provider_key_fingerprint == fingerprint:
                    return self._process
                self.close()
            node_path, script_path = self.runtime_paths()
            environment = self._node_environment()
            environment["SCANSCIPI_PROVIDER_KEY"] = api_key
            creation_flags = getattr(subprocess, "CREATE_NO_WINDOW", 0)
            raw_output: Queue[str | None] = Queue()
            error_lines: list[str] = []
            self._output = raw_output
            self._process_generation += 1
            self._dispatcher = _ProtocolDispatcher(
                raw_output,
                generation=self._process_generation,
                audit=self._dispatch_audit,
            )
            self._errors = error_lines
            process = subprocess.Popen(
                [str(node_path), str(script_path)],
                stdin=subprocess.PIPE,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                errors="replace",
                bufsize=1,
                env=environment,
                creationflags=creation_flags,
            )
            assert process.stdin is not None
            assert process.stdout is not None
            assert process.stderr is not None
            self._process = process
            self._provider_key_fingerprint = fingerprint

            def drain_stdout() -> None:
                for line in process.stdout:
                    raw_output.put(line)
                raw_output.put(None)

            def drain_stderr() -> None:
                for line in process.stderr:
                    if len(error_lines) < 80:
                        error_lines.append(line.rstrip())

            threading.Thread(target=drain_stdout, daemon=True).start()
            threading.Thread(target=drain_stderr, daemon=True).start()
            return process

    def _ensure_dispatcher(self) -> _ProtocolDispatcher:
        dispatcher = self._dispatcher
        if dispatcher is None or dispatcher.raw_output is not self._output:
            if self._process_generation <= 0:
                self._process_generation = 1
            dispatcher = _ProtocolDispatcher(
                self._output,
                generation=self._process_generation,
                audit=self._dispatch_audit,
            )
            self._dispatcher = dispatcher
        return dispatcher

    def _ensure_protocol_negotiated(self, *, timeout_seconds: float) -> None:
        """Fail closed before a fresh sidecar may mutate durable sessions."""

        if self._protocol_generation == self._process_generation and self._process_generation > 0:
            return
        try:
            event = self._await_command(
                {
                    "type": "ping",
                    "command_id": uuid4().hex,
                    "pi_protocol_version": _PI_PROTOCOL_VERSION,
                    "required_features": list(_PI_REQUIRED_FEATURES),
                },
                terminal_types={"pong"},
                timeout_seconds=max(0.05, float(timeout_seconds)),
                session_lock_id=f"__protocol__:{self._process_generation}",
                timeout_message="Pi sidecar protocol handshake exceeded the timeout",
            )
        except Exception as error:
            raise PiRuntimeUnavailable(
                f"The ScanSci Pi sidecar protocol handshake failed: {error}"
            ) from error
        protocol = int(event.get("protocol", 0) or 0)
        capabilities = {str(value) for value in list(event.get("capabilities", []) or [])}
        missing_features = [feature for feature in _PI_REQUIRED_FEATURES if feature not in capabilities]
        if protocol != _PI_PROTOCOL_VERSION or missing_features:
            detail = f"protocol={protocol}, missing={','.join(missing_features) or 'none'}"
            raise PiRuntimeUnavailable(f"The ScanSci Pi sidecar protocol is incompatible ({detail})")
        self._protocol_generation = self._process_generation

    def _write(self, message: dict[str, Any]) -> None:
        process = self._process
        if process is None or process.poll() is not None or process.stdin is None:
            raise PiRuntimeUnavailable("The ScanSci Pi sidecar is not running")
        with self._stdin_lock:
            encoded = json.dumps(message, ensure_ascii=False, separators=(",", ":"))
            if len(encoded.encode("utf-8")) > _MAX_JSONL_LINE_BYTES:
                raise ValueError("Pi protocol message exceeds the bounded JSONL line limit")
            process.stdin.write(encoded + "\n")
            process.stdin.flush()

    def _session_command_lock(self, session_id: str) -> threading.Lock:
        key = str(session_id or "").strip() or "__runtime__"
        with self._session_command_locks_guard:
            return self._session_command_locks.setdefault(key, threading.Lock())

    def _current_cancel_event(self) -> threading.Event:
        context = getattr(self._worker_context, "run_context", None)
        if isinstance(context, _RunContext):
            return context.cancel_event
        return self._cancel_requested

    def _record_dispatch_audit(
        self,
        kind: str,
        *,
        request_id: str = "",
        command_id: str = "",
        event_type: str = "",
        generation: int = 0,
    ) -> None:
        self._dispatch_audit.append({
            "kind": str(kind)[:80],
            "generation": int(generation or self._process_generation),
            "type": str(event_type)[:80],
            "request_id": str(request_id)[:80],
            "command_id": str(command_id)[:80],
        })

    def _context_is_active(self, context: _RunContext) -> bool:
        with self._contexts_lock:
            return self._run_contexts.get(context.request_id) is context and context.active

    @staticmethod
    def _mark_context_cancelled(context: _RunContext) -> None:
        # Signal workers immediately, then cross the final result-commit lock.
        # If the commit already owns the lock it linearizes first and finishes
        # file/wire/history consistently; otherwise the commit observes the
        # event and is suppressed before writing to the sidecar.
        context.cancel_event.set()
        with context.result_commit_lock:
            pass

    def _execute_tool_worker(
        self,
        context: _RunContext,
        pending: _PendingToolCall,
    ) -> _ToolCompletion:
        self._worker_context.run_context = context
        try:
            if context.cancel_event.is_set():
                raise InterruptedError("Pi tool execution was cancelled before it started")
            result = self._execute_tool(pending.name, dict(pending.arguments))
            if context.cancel_event.is_set():
                raise InterruptedError("Pi tool execution was cancelled")
            retried = _retry_empty_search(
                pending.name,
                pending.arguments,
                result,
                executor=self._execute_tool,
            )
            if retried is not None:
                result = retried
            return _ToolCompletion(
                request_id=context.request_id,
                generation=context.generation,
                call_id=pending.call_id,
                name=pending.name,
                arguments=dict(pending.arguments),
                parallel_safe=pending.parallel_safe,
                result=result,
            )
        except Exception as error:  # noqa: BLE001 - serialized on the owning event thread
            return _ToolCompletion(
                request_id=context.request_id,
                generation=context.generation,
                call_id=pending.call_id,
                name=pending.name,
                arguments=dict(pending.arguments),
                parallel_safe=pending.parallel_safe,
                error_type=type(error).__name__,
                error=str(error),
            )
        finally:
            self._worker_context.run_context = None

    def _submit_tool(self, context: _RunContext, pending: _PendingToolCall) -> None:
        if not self._tool_slots.acquire(blocking=False):
            context.completions.put_nowait(_ToolCompletion(
                request_id=context.request_id,
                generation=context.generation,
                call_id=pending.call_id,
                name=pending.name,
                arguments=dict(pending.arguments),
                parallel_safe=pending.parallel_safe,
                error_type="RuntimeError",
                error="Pi tool executor capacity is exhausted",
            ))
            return
        try:
            future = self._tool_executor.submit(self._execute_tool_worker, context, pending)
        except BaseException:
            self._tool_slots.release()
            raise
        with context.lock:
            context.in_flight[pending.call_id] = future
            if pending.parallel_safe:
                context.parallel_in_flight += 1
            else:
                context.sequential_in_flight = True

        def completed(done: Future[_ToolCompletion]) -> None:
            self._tool_slots.release()
            try:
                completion = done.result()
            except BaseException as error:  # pragma: no cover - worker converts ordinary failures
                completion = _ToolCompletion(
                    request_id=context.request_id,
                    generation=context.generation,
                    call_id=pending.call_id,
                    name=pending.name,
                    arguments=dict(pending.arguments),
                    parallel_safe=pending.parallel_safe,
                    error_type=type(error).__name__,
                    error=str(error),
                )
            if not self._context_is_active(context):
                self._record_dispatch_audit(
                    "late_tool_completion",
                    request_id=context.request_id,
                    event_type="tool.result",
                    generation=context.generation,
                )
                return
            try:
                context.completions.put_nowait(completion)
            except Full:
                self._mark_context_cancelled(context)
                self._record_dispatch_audit(
                    "completion_queue_full",
                    request_id=context.request_id,
                    event_type="tool.result",
                    generation=context.generation,
                )

        future.add_done_callback(completed)

    def _enqueue_tool(self, context: _RunContext, pending: _PendingToolCall) -> None:
        with context.lock:
            if not context.active:
                return
            can_start = (
                not context.sequential_in_flight
                and not context.deferred
                and (pending.parallel_safe or not context.in_flight)
            )
            if not can_start:
                context.deferred.append(pending)
                return
        self._submit_tool(context, pending)

    def _release_tool_lane(self, context: _RunContext, completion: _ToolCompletion) -> None:
        to_start: list[_PendingToolCall] = []
        with context.lock:
            future = context.in_flight.pop(completion.call_id, None)
            if future is not None:
                if completion.parallel_safe:
                    context.parallel_in_flight = max(0, context.parallel_in_flight - 1)
                else:
                    context.sequential_in_flight = False
            if context.in_flight or not context.active:
                return
            while context.deferred:
                pending = context.deferred.popleft()
                to_start.append(pending)
                if not pending.parallel_safe:
                    break
                if context.deferred and not context.deferred[0].parallel_safe:
                    break
        for pending in to_start:
            self._submit_tool(context, pending)

    def _tool_result_commit_checkpoint(
        self,
        _stage: str,
        _context: _RunContext,
        _completion: _ToolCompletion,
    ) -> None:
        """Test seam for deterministic cancellation at result commit phases."""

    def _finish_tool_completion(
        self,
        context: _RunContext,
        completion: _ToolCompletion,
        *,
        task_mode: str,
    ) -> dict[str, Any] | None:
        self._release_tool_lane(context, completion)
        if not self._context_is_active(context) or context.cancel_event.is_set():
            self._record_dispatch_audit(
                "late_tool_completion",
                request_id=context.request_id,
                event_type="tool.result",
                generation=context.generation,
            )
            return None
        if completion.error_type:
            public_error = str(_redact_tool_value(completion.error))[:500]
            self._tool_result_commit_checkpoint("before_error_wire", context, completion)
            with context.result_commit_lock:
                if not self._context_is_active(context) or context.cancel_event.is_set():
                    self._record_dispatch_audit(
                        "late_tool_completion",
                        request_id=context.request_id,
                        event_type="tool.result",
                        generation=context.generation,
                    )
                    return None
                self._write({
                    "type": "tool.result",
                    "request_id": context.request_id,
                    "call_id": completion.call_id,
                    "ok": False,
                    "error": f"{completion.error_type}: {public_error}",
                })
                with context.history_lock:
                    context.history.append({
                        "name": completion.name,
                        "status": "failed",
                        "error": public_error[:200],
                    })
                self._tool_result_commit_checkpoint("after_error_wire", context, completion)
            return {"type": "tool.failed", "name": completion.name, "error": public_error}

        safe_result = _redact_tool_value(_json_safe(completion.result))
        prepared: _PreparedToolResult | None = None
        wire_committed = False
        try:
            if completion.name == "discover_papers" and isinstance(safe_result, dict):
                prepared = self._prepare_tool_result(
                    "discover-papers",
                    safe_result,
                    request_id=context.request_id,
                    generation=context.generation,
                    call_id=completion.call_id,
                )
                model_result = _compact_academic_search_result_for_model(
                    safe_result,
                    requested_limit=_bounded_limit(
                        completion.arguments.get("result_limit", completion.arguments.get("limit")),
                        default=8,
                    ),
                    result_reference=prepared.reference,
                    result_bytes=prepared.encoded_bytes,
                )
                result_meta: dict[str, int | bool] = {
                    "original_bytes": prepared.encoded_bytes,
                    "model_bytes": len(json.dumps(model_result, ensure_ascii=False).encode("utf-8")),
                    "truncated": True,
                    "persist_full": True,
                }
            else:
                model_result, result_meta = _bounded_tool_result_for_model(completion.name, safe_result)
                if result_meta["persist_full"]:
                    prepared = self._prepare_tool_result(
                        completion.name,
                        safe_result if isinstance(safe_result, dict) else {"result": safe_result},
                        request_id=context.request_id,
                        generation=context.generation,
                        call_id=completion.call_id,
                    )
                    model_result["_full_result_reference"] = prepared.reference
                    model_result["_persisted_bytes"] = prepared.encoded_bytes
            if prepared is not None:
                self._tool_result_commit_checkpoint("persist_prepared", context, completion)
                if not self._context_is_active(context) or context.cancel_event.is_set():
                    self._record_dispatch_audit(
                        "late_tool_completion",
                        request_id=context.request_id,
                        event_type="tool.result",
                        generation=context.generation,
                    )
                    return None
                os.replace(prepared.temporary_path, prepared.final_path)
                self._tool_result_commit_checkpoint("result_renamed", context, completion)
                if not self._context_is_active(context) or context.cancel_event.is_set():
                    self._record_dispatch_audit(
                        "late_tool_completion",
                        request_id=context.request_id,
                        event_type="tool.result",
                        generation=context.generation,
                    )
                    return None

            self._tool_result_commit_checkpoint("before_wire", context, completion)
            with context.result_commit_lock:
                if not self._context_is_active(context) or context.cancel_event.is_set():
                    self._record_dispatch_audit(
                        "late_tool_completion",
                        request_id=context.request_id,
                        event_type="tool.result",
                        generation=context.generation,
                    )
                    return None
                # This dedicated lock has no sidecar callback dependency.  It
                # only orders cancellation against the short stdin commit.
                self._write({
                    "type": "tool.result",
                    "request_id": context.request_id,
                    "call_id": completion.call_id,
                    "ok": True,
                    "result": model_result,
                })
                wire_committed = True
                with context.history_lock:
                    context.history.append({
                        "name": completion.name,
                        "status": "ok",
                        "result_summary": _summarize_tool_result(completion.name, model_result),
                    })
                self._tool_result_commit_checkpoint("after_wire", context, completion)
        finally:
            if prepared is not None and not wire_committed:
                self._discard_prepared_tool_result(prepared)
        if completion.name == "build_verified_answer" and task_mode == "verified-answer":
            context.terminal_tool_completed = completion.name
            self._write({
                "type": "run.cancel",
                "request_id": context.request_id,
                "command_id": uuid4().hex,
                "generation": self._process_generation,
            })
        return {
            "type": "tool.completed",
            "name": completion.name,
            "result": model_result,
            "result_bytes": result_meta["original_bytes"],
            "model_result_bytes": result_meta["model_bytes"],
            "result_truncated": result_meta["truncated"],
        }

    def _teardown_run_context(
        self,
        context: _RunContext,
        dispatcher: _ProtocolDispatcher,
        event_channel: Queue[dict[str, Any] | None],
        *,
        terminal_received: bool,
        request_started: bool,
    ) -> None:
        self._mark_context_cancelled(context)
        with context.lock:
            context.active = False
            pending_futures = list(context.in_flight.values())
            context.deferred.clear()
        for future in pending_futures:
            future.cancel()
        with self._contexts_lock:
            if self._run_contexts.get(context.request_id) is context:
                self._run_contexts.pop(context.request_id, None)

        if request_started and not terminal_received:
            command_id = uuid4().hex
            try:
                self._write({
                    "type": "run.cancel",
                    "request_id": context.request_id,
                    "command_id": command_id,
                    "generation": self._process_generation,
                })
            except (BrokenPipeError, OSError, PiRuntimeUnavailable):
                pass
            else:
                cancel_deadline = time.monotonic() + _CANCEL_DRAIN_SECONDS
                cancellation_settled = False
                while time.monotonic() < cancel_deadline:
                    remaining = cancel_deadline - time.monotonic()
                    try:
                        event = event_channel.get(timeout=min(0.05, max(0.001, remaining)))
                    except Empty:
                        process = self._process
                        if process is not None and process.poll() is not None:
                            break
                        continue
                    if event is None:
                        break
                    event_type = str(event.get("type", ""))
                    if event_type in {"run.cancelled", "run.completed", "run.failed"}:
                        cancellation_settled = True
                        break
                    if event_type == "run.cancel_rejected" and str(event.get("command_id", "")) == command_id:
                        cancellation_settled = True
                        break
                if not cancellation_settled:
                    self._record_dispatch_audit(
                        "cancel_drain_timeout",
                        request_id=context.request_id,
                        command_id=command_id,
                        event_type="run.cancel",
                        generation=context.generation,
                    )

        dispatcher.unregister_request(context.request_id)
        if self._active_request_id == context.request_id:
            self._active_request_id = ""
            self._active_session_id = ""
        with self._interaction_lock:
            self._approval_tokens.pop(context.request_id, None)
            stale_interactions = [
                interaction_id
                for interaction_id, (pending_request_id, _kind) in self._pending_interaction_kinds.items()
                if pending_request_id == context.request_id
            ]
            for interaction_id in stale_interactions:
                self._pending_interaction_kinds.pop(interaction_id, None)

    def _run_request(
        self,
        start_message: dict[str, Any],
        *,
        api_key: str,
        timeout_seconds: float,
    ) -> Iterator[dict[str, Any]]:
        if not self._run_lock.acquire(blocking=False):
            raise RuntimeError("This Pi client already has an active run")
        request_id = str(start_message["request_id"])
        session_id = str(start_message["session_id"])
        task_mode = str(start_message.get("task_mode", "general") or "general").strip().lower()
        self._lifecycle_lock.acquire()
        dispatcher: _ProtocolDispatcher | None = None
        event_channel: Queue[dict[str, Any] | None] | None = None
        context: _RunContext | None = None
        request_started = False
        try:
            process = self._ensure_process(api_key=api_key)
            dispatcher = self._ensure_dispatcher()
            event_channel = dispatcher.register_request(request_id)
            dispatcher.start()
            self._request_generation += 1
            context = _RunContext(
                request_id=request_id,
                session_id=session_id,
                generation=self._request_generation,
                supports_images="image" in set(
                    dict(start_message.get("model_runtime", {}) or {}).get("input_modalities", [])
                ),
            )
            with self._contexts_lock:
                self._run_contexts[request_id] = context
            self._active_request_id = request_id
            self._active_session_id = session_id
            self._cancel_requested = context.cancel_event
            self._tool_history = context.history
            request_started = True
            self._write(start_message)
        except BaseException:
            try:
                if context is not None and dispatcher is not None and event_channel is not None:
                    self._teardown_run_context(
                        context,
                        dispatcher,
                        event_channel,
                        terminal_received=False,
                        request_started=request_started,
                    )
                elif dispatcher is not None and event_channel is not None:
                    dispatcher.unregister_request(request_id)
            except BaseException as cleanup_error:  # noqa: BLE001 - preserve the setup failure
                self._record_dispatch_audit(
                    "run_setup_cleanup_failed",
                    request_id=request_id,
                    event_type=type(cleanup_error).__name__,
                )
            finally:
                self._lifecycle_lock.release()
                self._run_lock.release()
            raise
        assert dispatcher is not None
        assert event_channel is not None
        assert context is not None
        deadline = time.monotonic() + timeout_seconds
        terminal_received = False
        task_contract = dict(start_message.get("task_contract", {}) or {})

        try:
            while True:
                while True:
                    try:
                        completion = context.completions.get_nowait()
                    except Empty:
                        break
                    public_event = self._finish_tool_completion(
                        context,
                        completion,
                        task_mode=task_mode,
                    )
                    if public_event is not None:
                        yield public_event

                remaining = deadline - time.monotonic()
                if remaining <= 0:
                    raise TimeoutError("Pi Agent exceeded the request timeout")
                if context.pending_terminal is not None and not context.has_pending_tools():
                    event = context.pending_terminal
                    context.pending_terminal = None
                else:
                    try:
                        routed = event_channel.get(timeout=min(0.05, remaining))
                    except Empty:
                        if process.poll() is not None:
                            detail = "\n".join(self._errors[-8:])
                            raise RuntimeError(f"Pi Agent exited unexpectedly{': ' + detail if detail else ''}")
                        continue
                    if routed is None:
                        detail = "\n".join(self._errors[-8:])
                        raise RuntimeError(f"Pi Agent closed its output stream{': ' + detail if detail else ''}")
                    event = routed

                event_type = str(event.get("type", ""))
                event_request_id = str(event.get("request_id", ""))
                if event_type == "protocol.error":
                    raise RuntimeError(str(event.get("error", "Pi Agent protocol error")))
                if event_type not in {"skill.call", "tool.call"} and event_request_id and event_request_id != request_id:
                    continue
                if event_type == "skill.call":
                    call_id = str(event.get("call_id", ""))
                    name = str(event.get("name", ""))
                    arguments = dict(event.get("arguments", {}) or {})
                    try:
                        if not event_request_id or event_request_id != request_id:
                            raise PermissionError("Skill call belongs to another request")
                        if self._cancel_requested.is_set():
                            raise InterruptedError("Pi Skill loading was cancelled")
                        result = self._execute_skill_tool(request_id, name, arguments)
                        response = {
                            "type": "skill.result",
                            "request_id": request_id,
                            "call_id": call_id,
                            "ok": True,
                            "result": result,
                        }
                        if name in {"load_skill", "restore_skill"}:
                            public = {
                                key: value
                                for key, value in result.items()
                                if key != "content"
                            }
                            yield {"type": "skill.loaded", "name": str(result.get("skill_id", "")), "value": public}
                        else:
                            yield {
                                "type": "skill.searched",
                                "name": "search_skills",
                                "value": {
                                    "query": str(result.get("query", "")),
                                    "count": int(result.get("count", 0) or 0),
                                    "limit": int(result.get("limit", 0) or 0),
                                    "skill_ids": [
                                        str(item.get("id", ""))
                                        for item in list(result.get("skills", []) or [])
                                        if isinstance(item, dict)
                                    ],
                                },
                            }
                    except Exception as error:  # noqa: BLE001 - returned to Pi as bounded tool output
                        public_error = str(_redact_tool_value(str(error)))[:500]
                        response = {
                            "type": "skill.result",
                            "request_id": request_id,
                            "call_id": call_id,
                            "ok": False,
                            "error": f"{type(error).__name__}: {public_error}",
                        }
                        yield {"type": "skill.failed", "name": name, "error": public_error}
                    self._write(response)
                    continue
                if event_type == "tool.call":
                    call_id = str(event.get("call_id", ""))
                    name = str(event.get("name", ""))
                    arguments = dict(event.get("arguments", {}) or {})
                    try:
                        # Tool calls are executable protocol messages, so an
                        # omitted request id is not accepted as a legacy event.
                        # The Python host independently rechecks every bridge
                        # call before entering the dispatcher.
                        authorize_tool_call(
                            tool_name=name,
                            contract=task_contract,
                            descriptor=builtin_capability_descriptor(name),
                            request_id=event_request_id,
                            active_request_id=request_id,
                            approval_token=self._approval_tokens.get(request_id),
                            call_count=context.authorized_call_count,
                        )
                        context.authorized_call_count += 1
                        if context.cancel_event.is_set():
                            raise InterruptedError("Pi tool execution was cancelled before it started")
                        descriptor = builtin_capability_descriptor(name)
                        parallel_safe = (
                            name in _PARALLEL_SAFE_TOOL_NAMES
                            and isinstance(descriptor, dict)
                            and str(descriptor.get("risk_level", "")) == "read_only"
                            and descriptor.get("idempotent") is not False
                        )
                        self._enqueue_tool(context, _PendingToolCall(
                            call_id=call_id,
                            name=name,
                            arguments=arguments,
                            parallel_safe=parallel_safe,
                        ))
                        continue
                    except Exception as error:  # noqa: BLE001 - error is returned to the model as tool output
                        public_error = str(_redact_tool_value(str(error)))[:500]
                        with context.history_lock:
                            context.history.append({
                                "name": name,
                                "status": "failed",
                                "error": public_error[:200],
                            })
                        response = {
                            "type": "tool.result",
                            "request_id": request_id,
                            "call_id": call_id,
                            "ok": False,
                            "error": f"{type(error).__name__}: {public_error}",
                        }
                        yield {"type": "tool.failed", "name": name, "error": public_error}
                    self._write(response)
                    continue
                if event_type == "session.ready":
                    durable_file = str(event.get("session_file", ""))
                    if durable_file:
                        registry = self._load_session_registry()
                        registry[session_id] = durable_file
                        self._save_session_registry(registry)
                    yield {
                        "type": "session",
                        "session_id": session_id,
                        "session_file": durable_file,
                        "resumed": bool(event.get("resumed", False)),
                        "request_id": request_id,
                    }
                elif event_type == "message.delta":
                    yield {"type": "delta", "content": str(event.get("delta", ""))}
                elif event_type == "status.update":
                    yield {
                        "type": "status",
                        "request_id": event_request_id or request_id,
                        "status": str(event.get("status", "")),
                        "name": str(event.get("name", "")),
                        "attempt": event.get("attempt"),
                        "error": str(event.get("error", "")),
                        "duration_ms": event.get("duration_ms"),
                        "details": dict(event.get("details", {}) or {}),
                    }
                elif event_type == "agent.queue_updated":
                    steering = [str(item) for item in list(event.get("steering", []) or [])]
                    follow_up = [str(item) for item in list(event.get("follow_up", []) or [])]
                    yield {
                        "type": "queue",
                        "steering": steering,
                        "follow_up": follow_up,
                        "pending_count": int(
                            event.get("pending_count", len(steering) + len(follow_up)) or 0
                        ),
                    }
                elif event_type in {
                    "agent.started",
                    "agent.turn_started",
                    "agent.message_started",
                    "agent.message_completed",
                    "agent.turn_completed",
                    "agent.completed",
                    "agent.settled",
                }:
                    yield {
                        "type": "lifecycle",
                        "event": event_type.removeprefix("agent."),
                        "role": str(event.get("role", "")),
                        "tool_result_count": int(event.get("tool_result_count", 0) or 0),
                        "will_retry": bool(event.get("will_retry", False)),
                    }
                elif event_type in {
                    "run.steer_ack",
                    "run.steer_rejected",
                    "run.follow_up_ack",
                    "run.follow_up_rejected",
                    "run.cancel_ack",
                    "run.cancel_rejected",
                }:
                    action = event_type.split(".", 1)[1].rsplit("_", 1)[0]
                    accepted = event_type.endswith("_ack")
                    yield {
                        "type": "control",
                        "command_id": str(event.get("command_id", "")),
                        "action": action,
                        "status": "accepted" if accepted else "rejected",
                        "error": str(event.get("error", "")),
                        "queued": int(event.get("queued", 0) or 0),
                    }
                elif event_type == "interaction.requested":
                    interaction_id = str(event.get("interaction_id", ""))
                    interaction_kind = str(event.get("interaction_kind", ""))
                    if interaction_id:
                        with self._interaction_lock:
                            self._pending_interaction_kinds[interaction_id] = (request_id, interaction_kind)
                    yield {
                        "type": "interaction",
                        "request_id": request_id,
                        "session_id": session_id,
                        "interaction_id": interaction_id,
                        "interaction_kind": interaction_kind,
                        "payload": dict(event.get("payload", {}) or {}),
                    }
                elif event_type == "run.completed":
                    if context.has_pending_tools():
                        context.pending_terminal = dict(event)
                        continue
                    terminal_received = True
                    yield {
                        "type": "done",
                        "stats": dict(event.get("stats", {}) or {}),
                        "control": dict(event.get("control", {}) or {}),
                        "truncated": False,
                    }
                    return
                elif event_type == "run.cancelled":
                    terminal_received = True
                    if context.terminal_tool_completed:
                        yield {
                            "type": "done",
                            "stats": {},
                            "truncated": False,
                            "terminal_tool": context.terminal_tool_completed,
                        }
                    else:
                        yield {"type": "cancelled", "request_id": request_id, "session_id": session_id}
                    return
                elif event_type in {"session.compaction_started", "session.compaction_completed"}:
                    payload: dict[str, Any] = {
                        "type": "compaction",
                        "status": "started" if event_type.endswith("started") else "completed",
                        "reason": str(event.get("reason", "")),
                        "aborted": bool(event.get("aborted", False)),
                        "error": str(event.get("error", "")),
                    }
                    if event_type.endswith("completed"):
                        result = dict(event.get("result", {}) or {})
                        if "tokensBefore" in result:
                            payload["tokens_before"] = int(result["tokensBefore"])
                        if "estimatedTokensAfter" in result:
                            payload["tokens_after"] = int(result["estimatedTokensAfter"])
                        if "summary" in result:
                            payload["summary"] = str(result["summary"])[:500]
                    yield payload
                elif event_type == "run.failed":
                    terminal_received = True
                    raise PiAgentRunError(
                        str(event.get("error", "Pi Agent failed")),
                        failure=dict(event.get("failure", {}) or {}),
                    )
        finally:
            try:
                self._teardown_run_context(
                    context,
                    dispatcher,
                    event_channel,
                    terminal_received=terminal_received,
                    request_started=request_started,
                )
            finally:
                self._lifecycle_lock.release()
                self._run_lock.release()

    @property
    def active_request_id(self) -> str:
        return self._active_request_id

    @property
    def transport_records(self) -> list[dict[str, Any]]:
        if self._gateway_adapter is None:
            return []
        return [dict(item) for item in self._gateway_adapter.records]

    def cancel(self, request_id: str | None = None) -> bool:
        """Abort the active SDK operation without killing the durable session."""

        target = request_id or self._active_request_id
        if not target:
            return False
        with self._contexts_lock:
            context = self._run_contexts.get(target)
        if context is None or not context.active:
            return False
        self._mark_context_cancelled(context)
        self._write({
            "type": "run.cancel",
            "request_id": target,
            "command_id": uuid4().hex,
            "generation": self._process_generation,
        })
        return True

    def pause(self, request_id: str | None = None) -> bool:
        """Request cooperative abort-and-resume for the current durable turn.

        The Pi sidecar exposes ``run.cancel`` as its interruption primitive.
        This is not a suspended SDK operation: the current operation aborts,
        while durable session history is retained so the host can explicitly
        start a new or resumed turn.
        """

        return self.cancel(request_id)

    def steer(
        self,
        text: str,
        request_id: str | None = None,
        *,
        images: list[dict[str, Any]] | None = None,
    ) -> bool:
        """Queue steering text into the active Pi tool loop."""

        target = request_id or self._active_request_id
        if not target:
            return False
        with self._contexts_lock:
            context = self._run_contexts.get(target)
        if context is None or not context.active:
            return False
        validated_images = validate_pi_image_blocks(list(images or []))
        if validated_images and not bool(getattr(context, "supports_images", False)):
            raise PiMultimodalUnavailable("The active Pi session does not support image input")
        self._write({
            "type": "run.steer",
            "request_id": target,
            "command_id": uuid4().hex,
            "generation": self._process_generation,
            "text": text,
            "images": validated_images,
        })
        return True

    def follow_up(
        self,
        text: str,
        request_id: str | None = None,
        *,
        images: list[dict[str, Any]] | None = None,
    ) -> bool:
        """Queue a message for the active session after its current turn."""

        target = request_id or self._active_request_id
        if not target or not str(text).strip():
            return False
        with self._contexts_lock:
            context = self._run_contexts.get(target)
        if context is None or not context.active:
            return False
        validated_images = validate_pi_image_blocks(list(images or []))
        if validated_images and not bool(getattr(context, "supports_images", False)):
            raise PiMultimodalUnavailable("The active Pi session does not support image input")
        self._write({
            "type": "run.follow_up",
            "request_id": target,
            "command_id": uuid4().hex,
            "generation": self._process_generation,
            "text": str(text).strip(),
            "images": validated_images,
        })
        return True

    def respond_interaction(
        self,
        interaction_id: str,
        response: dict[str, Any],
        *,
        request_id: str | None = None,
    ) -> bool:
        """Resolve a pending AskUser or plan-approval tool call."""

        target = request_id or self._active_request_id
        if not target or not str(interaction_id).strip():
            return False
        normalized_interaction_id = str(interaction_id).strip()
        with self._interaction_lock:
            pending = self._pending_interaction_kinds.get(normalized_interaction_id)
            if pending and pending[0] == target and pending[1] == "plan":
                token = approval_token_from_response(target, response)
                if token is None:
                    self._approval_tokens.pop(target, None)
                else:
                    self._approval_tokens[target] = token
            if pending and pending[0] == target:
                self._pending_interaction_kinds.pop(normalized_interaction_id, None)
        self._write(
            {
                "type": "interaction.response",
                "request_id": target,
                "interaction_id": normalized_interaction_id,
                "response": dict(response or {}),
            }
        )
        return True

    def _await_command(
        self,
        message: dict[str, Any],
        *,
        terminal_types: set[str],
        timeout_seconds: float,
        session_lock_id: str,
        timeout_message: str,
    ) -> dict[str, Any]:
        if str(message.get("type", "")).startswith("session."):
            message = {
                **message,
                "pi_protocol_version": _PI_PROTOCOL_VERSION,
                "required_features": list(_PI_REQUIRED_FEATURES),
            }
        command_id = str(message.get("command_id", ""))
        if not command_id:
            raise ValueError("Pi session commands require a command_id")
        with self._lifecycle_lock:
            dispatcher = self._ensure_dispatcher()
            message = {**message, "generation": self._process_generation}
            lock = self._session_command_lock(session_lock_id)
            with lock:
                channel = dispatcher.register_command(command_id)
                dispatcher.start()
                try:
                    # Registration must happen before stdin write: a local sidecar
                    # can acknowledge a command in the same scheduling quantum.
                    self._write(message)
                    deadline = time.monotonic() + max(0.01, float(timeout_seconds))
                    while True:
                        remaining = deadline - time.monotonic()
                        if remaining <= 0:
                            raise TimeoutError(timeout_message)
                        try:
                            event = channel.get(timeout=min(0.25, remaining))
                        except Empty:
                            process = self._process
                            if process is not None and process.poll() is not None:
                                raise RuntimeError("Pi Agent exited while awaiting a session command")
                            continue
                        if event is None:
                            raise RuntimeError("Pi Agent closed while awaiting a session command")
                        if str(event.get("type", "")) == "protocol.error":
                            raise RuntimeError(str(event.get("error", "Pi Agent protocol error")))
                        if str(event.get("type", "")) in terminal_types:
                            self._validate_command_correlation(message, event)
                            return dict(event)
                finally:
                    dispatcher.unregister_command(command_id)

    @staticmethod
    def _validate_command_correlation(message: dict[str, Any], event: dict[str, Any]) -> None:
        expected_session = str(message.get("session_id", ""))
        actual_session = str(event.get("session_id", ""))
        if expected_session and actual_session != expected_session:
            raise RuntimeError("Pi session command acknowledgement crossed session boundaries")
        expected_source = str(message.get("source_session_id", ""))
        actual_source = str(event.get("source_session_id", ""))
        if expected_source and actual_source != expected_source:
            raise RuntimeError("Pi session fork acknowledgement crossed source sessions")
        expected_target = str(message.get("target_session_id", ""))
        actual_target = str(event.get("target_session_id", ""))
        if expected_target and actual_target != expected_target:
            raise RuntimeError("Pi session fork acknowledgement crossed target sessions")

    def _await_control_command(
        self,
        message: dict[str, Any],
        *,
        terminal_types: set[str],
        timeout_seconds: float,
        timeout_message: str,
    ) -> dict[str, Any]:
        """Await a non-replacing control without taking lifecycle/session locks."""

        if str(message.get("type", "")).startswith("session."):
            message = {
                **message,
                "pi_protocol_version": _PI_PROTOCOL_VERSION,
                "required_features": list(_PI_REQUIRED_FEATURES),
            }
        command_id = str(message.get("command_id", ""))
        if not command_id:
            raise ValueError("Pi control commands require a command_id")
        dispatcher = self._ensure_dispatcher()
        message = {**message, "generation": self._process_generation}
        channel = dispatcher.register_command(command_id)
        dispatcher.start()
        try:
            self._write(message)
            deadline = time.monotonic() + max(0.01, float(timeout_seconds))
            while True:
                remaining = deadline - time.monotonic()
                if remaining <= 0:
                    raise TimeoutError(timeout_message)
                try:
                    event = channel.get(timeout=min(0.25, remaining))
                except Empty:
                    process = self._process
                    if process is not None and process.poll() is not None:
                        raise RuntimeError("Pi Agent exited while awaiting a control command")
                    continue
                if event is None:
                    raise RuntimeError("Pi Agent closed while awaiting a control command")
                if str(event.get("type", "")) == "protocol.error":
                    raise RuntimeError(str(event.get("error", "Pi Agent protocol error")))
                if str(event.get("type", "")) in terminal_types:
                    self._validate_command_correlation(message, event)
                    return dict(event)
        finally:
            dispatcher.unregister_command(command_id)

    def inspect_queue(self, session_id: str, *, timeout_seconds: float = 30.0) -> dict[str, Any]:
        return self._await_control_command(
            {"type": "session.queue.inspect", "command_id": uuid4().hex, "session_id": str(session_id)},
            terminal_types={"session.queue", "session.queue_failed"},
            timeout_seconds=timeout_seconds,
            timeout_message="Inspecting the Pi session queue exceeded the timeout",
        )

    def clear_queue(self, session_id: str, *, timeout_seconds: float = 30.0) -> dict[str, Any]:
        return self._await_control_command(
            {"type": "session.queue.clear", "command_id": uuid4().hex, "session_id": str(session_id)},
            terminal_types={"session.queue_cleared", "session.queue_failed"},
            timeout_seconds=timeout_seconds,
            timeout_message="Clearing the Pi session queue exceeded the timeout",
        )

    def abort_compaction(self, session_id: str, *, timeout_seconds: float = 30.0) -> dict[str, Any]:
        return self._await_control_command(
            {"type": "session.compact.abort", "command_id": uuid4().hex, "session_id": str(session_id)},
            terminal_types={"session.compact_aborted", "session.compact_abort_failed"},
            timeout_seconds=timeout_seconds,
            timeout_message="Aborting Pi session compaction exceeded the timeout",
        )

    def compact(self, session_id: str, *, instructions: str = "", timeout_seconds: float = 180.0) -> dict[str, Any]:
        """Run Pi's native context compactor for a loaded durable session."""

        if self._active_request_id:
            raise RuntimeError("Cannot manually compact while a run is active")
        if self._process is None or self._process.poll() is not None:
            raise PiRuntimeUnavailable("Load the durable session with stream_chat before compacting it")
        command_id = uuid4().hex
        event = self._await_command(
            {
                "type": "session.compact",
                "command_id": command_id,
                "session_id": session_id,
                "instructions": instructions,
            },
            terminal_types={"session.compact_completed", "session.compact_failed"},
            timeout_seconds=timeout_seconds,
            session_lock_id=session_id,
            timeout_message="Pi session compaction exceeded the timeout",
        )
        if event.get("type") == "session.compact_failed":
            raise RuntimeError(str(event.get("error", "Pi compaction failed")))
        result = dict(event.get("result", {}) or {})
        if isinstance(event.get("stats"), dict):
            result["_session_stats"] = dict(event["stats"])
        return result

    def load_session(
        self,
        session_id: str,
        *,
        provider_kind: str,
        base_url: str,
        api_key: str,
        model_id: str,
        api_surface: str = "chat_completions",
        responses_enabled: bool = False,
        model_runtime: dict[str, Any] | ModelRuntimeDescriptor | None = None,
        thinking_level: str = "medium",
        timeout_seconds: float = 60.0,
    ) -> dict[str, Any]:
        """Load a persisted Pi session into this sidecar without prompting it."""

        if self._active_request_id:
            raise RuntimeError("Cannot load a session while a run is active")
        logical_session_id = str(session_id or "").strip()
        session_file = str(self._load_session_registry().get(logical_session_id, ""))
        if not logical_session_id or not session_file or not Path(session_file).is_file():
            raise PiRuntimeUnavailable("The durable Pi session file is unavailable")
        if isinstance(model_runtime, ModelRuntimeDescriptor):
            descriptor = model_runtime
        elif isinstance(model_runtime, dict):
            descriptor = ModelRuntimeDescriptor.from_payload(model_runtime)
        else:
            descriptor = descriptor_from_model_record(
                provider_id="runtime",
                provider_kind=provider_kind,
                model_id=model_id,
                model_record=None,
                api_surface=api_surface,
            )
        if (
            descriptor.provider_kind != str(provider_kind)
            or descriptor.model_id != str(model_id)
            or descriptor.api_surface != str(api_surface or "chat_completions")
        ):
            raise ValueError(
                "Model runtime descriptor does not match the final provider route"
            )
        with self._lifecycle_lock:
            self._ensure_process(api_key=api_key)
            self._ensure_protocol_negotiated(timeout_seconds=min(8.0, timeout_seconds))
            command_id = uuid4().hex
            event = self._await_command(
                {
                    "type": "session.load",
                    "command_id": command_id,
                    "session_id": logical_session_id,
                    "session_file": session_file,
                    "cwd": str(self.workspace.parent),
                    "agent_dir": str(self.agent_dir),
                    "provider_kind": provider_kind,
                    "base_url": base_url,
                    "model_id": model_id,
                    "api_surface": api_surface,
                    "responses_enabled": responses_enabled,
                    "model_runtime": descriptor.to_dict(),
                    "thinking_level": thinking_level,
                    "mcp_servers": self._enabled_mcp_servers(),
                    "disabled_tools": self._disabled_artifact_tools(),
                },
                terminal_types={"session.loaded", "session.load_failed"},
                timeout_seconds=timeout_seconds,
                session_lock_id=logical_session_id,
                timeout_message="Loading the Pi session exceeded the timeout",
            )
        if event.get("type") == "session.load_failed":
            raise RuntimeError(str(event.get("error", "Pi session load failed")))
        return event

    def close_session(self, session_id: str, *, timeout_seconds: float = 30.0) -> dict[str, Any]:
        """Unload a session while retaining its persisted JSONL file."""

        command_id = uuid4().hex
        if self._process is None or self._process.poll() is not None:
            return {
                "type": "session.closed",
                "command_id": command_id,
                "session_id": str(session_id or ""),
                "generation": self._process_generation,
                "not_loaded": True,
            }
        event = self._await_command(
            {
                "type": "session.close",
                "command_id": command_id,
                "session_id": session_id,
            },
            terminal_types={"session.closed", "session.close_rejected"},
            timeout_seconds=timeout_seconds,
            session_lock_id=session_id,
            timeout_message="Closing the Pi session exceeded the timeout",
        )
        if event.get("type") == "session.close_rejected":
            raise RuntimeError(str(event.get("error", "Pi session close was rejected")))
        return event

    def fork_session(
        self,
        source_session_id: str,
        *,
        target_session_id: str | None = None,
        entry_id: str = "",
        before: bool = False,
        timeout_seconds: float = 90.0,
    ) -> dict[str, Any]:
        """Fork a loaded durable Pi session into a new independently resumable session."""

        if self._active_request_id:
            raise RuntimeError("Cannot fork a session while this client has an active run")
        if self._process is None or self._process.poll() is not None:
            raise PiRuntimeUnavailable("Load the source session before forking it")
        source_id = str(source_session_id or "").strip()
        target_id = str(target_session_id or uuid4().hex).strip()
        if not source_id or not target_id:
            raise ValueError("Source and target session ids are required")
        command_id = uuid4().hex
        event = self._await_command(
            {
                "type": "session.fork",
                "command_id": command_id,
                "generation": self._process_generation,
                "source_session_id": source_id,
                "target_session_id": target_id,
                "entry_id": str(entry_id or "").strip(),
                "before": bool(before),
                "full_history": not bool(str(entry_id or "").strip()),
            },
            terminal_types={"session.forked", "session.fork_failed"},
            timeout_seconds=timeout_seconds,
            session_lock_id=source_id,
            timeout_message="Forking the Pi session exceeded the timeout",
        )
        if event.get("type") == "session.fork_failed":
            raise PiAgentRunError(
                str(event.get("error", "Pi session fork failed")),
                failure=dict(event.get("failure", {}) or {}),
            )
        durable_file = str(event.get("session_file", ""))
        if durable_file:
            self._save_session_registry({target_id: durable_file})
        return event

    def forget_session(self, session_id: str) -> None:
        """Forget a broken session mapping without deleting its audit file.

        A provider can leave a durable session with an incomplete assistant
        turn (for example after a gateway timeout).  The next request should
        be able to rebuild that logical session from ScanSci's durable task
        messages rather than repeatedly resuming the poisoned transcript.
        """

        normalized = str(session_id or "").strip()
        if not normalized:
            return
        path = self._registry_path()
        temporary = path.with_name(f"{path.name}.{uuid4().hex}.tmp")
        with _SESSION_REGISTRY_LOCK:
            registry = self._load_session_registry_unlocked()
            if normalized not in registry:
                return
            registry.pop(normalized, None)
            temporary.write_text(json.dumps(registry, ensure_ascii=False, indent=2), encoding="utf-8")
            os.replace(temporary, path)

    def close(self) -> None:
        """Stop the sidecar; persisted sessions remain recoverable."""

        with self._lifecycle_lock:
            process = self._process
            if process is not None and process.poll() is None:
                try:
                    self._write({"type": "runtime.shutdown"})
                    process.wait(timeout=3)
                except (BrokenPipeError, OSError, subprocess.TimeoutExpired):
                    process.kill()
                    process.wait(timeout=5)
            self._process = None
            self._provider_key_fingerprint = ""
            if self._gateway_adapter is not None:
                self._gateway_adapter.close()
                self._gateway_adapter = None
                self._gateway_adapter_signature = ""

    def __enter__(self) -> "PiAgentClient":
        return self

    def __exit__(self, *_args: object) -> None:
        self.close()

    def _registry_path(self) -> Path:
        return self.agent_dir / "sessions.json"

    def _enabled_mcp_servers(self) -> list[dict[str, Any]]:
        """Return user-enabled MCP connections without provider credentials.

        Saving a record never launches it.  The Pi sidecar connects these
        records only while creating a task session, and write-like tools stay
        hidden unless the user explicitly enabled ``allow_write``.
        """

        settings = load_settings(self.workspace)
        servers = []
        for record in list(settings.get("mcp_servers", []) or []):
            if not isinstance(record, dict) or not record.get("enabled") or record.get("uninstalled"):
                continue
            transport = str(record.get("transport", "stdio") or "stdio")
            if transport == "stdio" and not str(record.get("command", "")).strip():
                continue
            if transport in {"streamable-http", "sse"} and not str(record.get("endpoint", "")).strip():
                continue
            servers.append(
                {
                    key: record.get(key)
                    for key in (
                        "id",
                        "name",
                        "description",
                        "transport",
                        "command",
                        "args",
                        "endpoint",
                        "connector_kind",
                        "allow_write",
                        "tool_effects",
                        "tool_policies",
                        "deferred",
                        "enabled",
                        "uninstalled",
                    )
                }
            )
        return servers

    def _disabled_artifact_tools(self) -> list[str]:
        plugin_tools = {
            "zotero": [
                "zotero_search",
                "zotero_fulltext",
                "zotero_attachment",
                "zotero_export_bibtex",
                "zotero_citations",
            ],
            "documents": ["create_document"],
            "pdf": ["create_pdf"],
            "spreadsheets": ["create_spreadsheet"],
            "presentations": ["create_presentation"],
            "latex": ["compile_latex"],
        }
        settings = load_settings(self.workspace)
        plugins = {str(item.get("id", "")): item for item in list(settings.get("plugins", []) or [])}
        disabled: list[str] = []
        for plugin_id, tool_names in plugin_tools.items():
            plugin = plugins.get(plugin_id)
            if not plugin or not plugin.get("enabled") or plugin.get("uninstalled"):
                disabled.extend(tool_names)
        return disabled

    def _effective_base_url(self, *, base_url: str, api_key: str) -> str:
        if _MANAGED_GATEWAY_HOST not in base_url.lower():
            return base_url
        signature = hashlib.sha256(f"{base_url.rstrip('/')}\0{api_key}".encode("utf-8")).hexdigest()
        if self._gateway_adapter is not None and self._gateway_adapter_signature == signature:
            return self._gateway_adapter.base_url
        if self._gateway_adapter is not None:
            self._gateway_adapter.close()
        self._gateway_adapter = _ManagedGatewayAdapter(upstream_base_url=base_url, api_key=api_key)
        self._gateway_adapter_signature = signature
        return self._gateway_adapter.base_url

    def _load_session_registry(self) -> dict[str, str]:
        with _SESSION_REGISTRY_LOCK:
            try:
                payload = json.loads(self._registry_path().read_text(encoding="utf-8"))
            except (FileNotFoundError, json.JSONDecodeError, OSError):
                return {}
        if not isinstance(payload, dict):
            return {}
        return {str(key): str(value) for key, value in payload.items() if value}

    def _save_session_registry(self, registry: dict[str, str]) -> None:
        path = self._registry_path()
        temporary = path.with_name(f"{path.name}.{uuid4().hex}.tmp")
        with _SESSION_REGISTRY_LOCK:
            latest = self._load_session_registry_unlocked()
            latest.update(registry)
            temporary.write_text(json.dumps(latest, ensure_ascii=False, indent=2), encoding="utf-8")
            os.replace(temporary, path)

    def _load_session_registry_unlocked(self) -> dict[str, str]:
        try:
            payload = json.loads(self._registry_path().read_text(encoding="utf-8"))
        except (FileNotFoundError, json.JSONDecodeError, OSError):
            return {}
        if not isinstance(payload, dict):
            return {}
        return {str(key): str(value) for key, value in payload.items() if value}

    @staticmethod
    def _recorded_document_paths(run: dict[str, Any]) -> list[Path]:
        """Resolve only document paths already registered by a ScanSci task."""

        candidates: list[Path] = []

        def visit_path_collection(value: Any) -> None:
            if isinstance(value, (list, tuple)):
                for nested in value:
                    visit_path_collection(nested)
                return
            if isinstance(value, dict):
                # A path-bearing record is allowed, but arbitrary nested
                # metadata is not.  Import results can contain a full
                # workspace snapshot with unrelated Zotero/Obsidian paths.
                for key in ("path", "file_path", "local_path", "source_path"):
                    visit_path_collection(value.get(key))
                return
            if not isinstance(value, str) or not value.strip():
                return
            try:
                candidate = Path(value).expanduser().resolve()
            except (OSError, RuntimeError):
                return
            if candidate.suffix.lower() in SUPPORTED_INGESTION_SUFFIXES and candidate.is_file():
                candidates.append(candidate)

        def visit_artifact(value: Any) -> None:
            if not isinstance(value, dict):
                return
            visit_path_collection(value.get("file_path"))
            payload = dict(value.get("payload", {}) or {})
            for key in ("files", "file", "file_path", "source_files", "documents"):
                visit_path_collection(payload.get(key))

        # Only explicit artifact/stage file contracts are trusted.  Do not
        # recursively walk the complete task object: importer diagnostics may
        # embed the entire workspace and would contaminate this task with
        # unrelated library documents.
        visit_artifact(run.get("output_artifact"))
        for artifact in list(run.get("artifacts", []) or []):
            visit_artifact(artifact)
        for stage in list(run.get("stages", []) or []):
            if not isinstance(stage, dict):
                continue
            output = dict(stage.get("output", {}) or {})
            for key in ("files", "file", "file_path", "source_files", "documents"):
                visit_path_collection(output.get(key))
        visit_path_collection(dict(run.get("input") or {}).get("source_files"))

        unique: dict[str, Path] = {}
        for candidate in candidates:
            unique.setdefault(os.path.normcase(str(candidate)), candidate)
        size_groups: dict[tuple[str, int], list[Path]] = {}
        for candidate in unique.values():
            try:
                size_groups.setdefault((candidate.suffix.lower(), candidate.stat().st_size), []).append(candidate)
            except OSError:
                continue
        content_unique: list[Path] = []
        for group in size_groups.values():
            if len(group) == 1:
                content_unique.extend(group)
                continue
            hashes: set[str] = set()
            for candidate in group:
                try:
                    digest = hashlib.sha256()
                    with candidate.open("rb") as handle:
                        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
                            digest.update(chunk)
                    fingerprint = digest.hexdigest()
                except OSError:
                    fingerprint = f"path:{os.path.normcase(str(candidate))}"
                if fingerprint in hashes:
                    continue
                hashes.add(fingerprint)
                content_unique.append(candidate)
        document_priority = {
            ".pdf": 0,
            ".docx": 1,
            ".pptx": 2,
            ".md": 3,
            ".markdown": 3,
            ".txt": 4,
        }
        return sorted(
            content_unique,
            key=lambda path: (document_priority.get(path.suffix.lower(), 10), path.name.lower()),
        )

    def _task_run_with_documents(self, requested_run_id: str = "") -> tuple[dict[str, Any], list[Path]]:
        store = ResearchRunStore(self.workspace)
        explicit_run_id = str(requested_run_id or self.active_run_id).strip()
        if explicit_run_id:
            run = store.get_run(explicit_run_id)
            return run, self._recorded_document_paths(run)

        for summary in store.list_runs(limit=50, archived=None):
            run_id = str(summary.get("run_id", "")).strip()
            if not run_id:
                continue
            run = store.get_run(run_id)
            paths = self._recorded_document_paths(run)
            if paths:
                return run, paths
        raise FileNotFoundError("No recent ScanSci task has readable document artifacts")

    def _read_task_documents(self, arguments: dict[str, Any]) -> dict[str, Any]:
        """Extract bounded excerpts from the active or most recent task artifacts."""

        run, paths = self._task_run_with_documents(str(arguments.get("run_id", "")))
        if not paths:
            return {
                "ok": False,
                "run_id": str(run.get("run_id", "")),
                "title": str(run.get("title", "")),
                "workflow_type": str(run.get("workflow_type", "")),
                "document_count": 0,
                "total_recorded": 0,
                "documents": [],
                "failures": [],
                "message": "The selected ScanSci task has no readable document artifacts.",
            }

        max_files = min(24, max(1, int(arguments.get("max_files", 20) or 20)))
        # The managed gateway accepts a smaller request body than the model's
        # advertised context window.  Keep the complete tool result bounded so
        # the follow-up generation (which also includes the conversation and
        # task context) cannot be rejected for exceeding that transport limit.
        per_file_chars = min(6_000, max(1_000, int(arguments.get("per_file_chars", 1_300) or 1_300)))
        selected = paths[:max_files]
        run_id = str(run.get("run_id", "")).strip() or "recent"
        safe_run_id = re.sub(r"[^A-Za-z0-9._-]+", "-", run_id)[:80] or "recent"
        output_dir = self.agent_dir / "task-documents" / safe_run_id
        documents: list[dict[str, Any]] = []
        failures: list[dict[str, str]] = []
        remaining_chars = 28_000

        for path in selected:
            if remaining_chars <= 0:
                break
            try:
                stat = path.stat()
                cache_key = hashlib.sha256(
                    f"{path}\0{stat.st_size}\0{stat.st_mtime_ns}".encode("utf-8")
                ).hexdigest()[:24]
                cache_text_path = output_dir / f"document-{cache_key}.md"
                cache_meta_path = output_dir / f"document-{cache_key}.json"
                if cache_text_path.is_file():
                    text = cache_text_path.read_text(encoding="utf-8")
                    try:
                        extracted = json.loads(cache_meta_path.read_text(encoding="utf-8"))
                    except (FileNotFoundError, OSError, json.JSONDecodeError):
                        extracted = {}
                    if not isinstance(extracted, dict):
                        extracted = {}
                    extracted.setdefault("parser", "cached")
                    extracted.setdefault("page_count", 0)
                    extracted.setdefault("character_count", len(text))
                    extracted.setdefault("warnings", [])
                else:
                    extracted = extract_local_document(path, output_dir=output_dir, parser="auto")
                    text_path = Path(str(extracted.get("text_path", ""))).resolve()
                    text = text_path.read_text(encoding="utf-8") if text_path.is_file() else ""
                    cache_text_path.write_text(text, encoding="utf-8")
                    cache_meta_path.write_text(
                        json.dumps(
                            {
                                "parser": str(extracted.get("parser", "")),
                                "page_count": int(extracted.get("page_count", 0) or 0),
                                "character_count": int(extracted.get("character_count", len(text)) or len(text)),
                                "warnings": list(extracted.get("warnings", []) or []),
                            },
                            ensure_ascii=False,
                        ),
                        encoding="utf-8",
                    )
                excerpt_limit = min(per_file_chars, remaining_chars)
                excerpt = text[:excerpt_limit].strip()
                if len(text) > len(excerpt) and len(excerpt) >= 1_200:
                    # Preserve a little of the conclusion as well as the
                    # opening abstract/introduction.  This is more useful for
                    # cross-paper synthesis than taking only the first page.
                    head_chars = max(700, int(len(excerpt) * 0.62))
                    tail_chars = max(200, len(excerpt) - head_chars)
                    excerpt = (
                        text[:head_chars].strip()
                        + "\n\n[…中间内容已省略…]\n\n"
                        + text[-tail_chars:].strip()
                    ).strip()
                remaining_chars -= len(excerpt)
                documents.append(
                    {
                        "name": path.name,
                        "path": str(path),
                        "parser": str(extracted.get("parser", "")),
                        "page_count": int(extracted.get("page_count", 0) or 0),
                        "character_count": int(extracted.get("character_count", len(text)) or len(text)),
                        "excerpt": excerpt,
                        "excerpt_truncated": len(text) > len(excerpt),
                        "warnings": list(extracted.get("warnings", []) or []),
                    }
                )
            except Exception as error:
                failures.append({"name": path.name, "error": f"{type(error).__name__}: {error}"[:500]})

        return {
            "ok": bool(documents),
            "run_id": run_id,
            "title": str(run.get("title", "")),
            "workflow_type": str(run.get("workflow_type", "")),
            "artifact_type": str(dict(run.get("output_artifact") or {}).get("artifact_type", "")),
            "document_count": len(documents),
            "total_recorded": len(paths),
            "truncated": len(paths) > len(selected) or len(documents) < len(selected) or remaining_chars <= 0,
            "documents": documents,
            "failures": failures,
        }

    def _prepare_tool_result(
        self,
        tool_name: str,
        payload: dict[str, Any],
        *,
        request_id: str = "",
        generation: int = 0,
        call_id: str = "",
    ) -> _PreparedToolResult:
        """Write a request-scoped temporary result without publishing it."""

        safe_name = re.sub(r"[^A-Za-z0-9._-]+", "-", str(tool_name).strip())[:32] or "tool"
        result_dir = self.agent_dir / "tool-results"
        result_dir.mkdir(parents=True, exist_ok=True)
        result_id = f"{safe_name}-{time.strftime('%Y%m%dT%H%M%S')}-{uuid4().hex[:12]}.json"
        result_path = result_dir / result_id
        encoded = json.dumps(_json_safe(payload), ensure_ascii=False, indent=2).encode("utf-8")
        scope_hash = hashlib.sha256(
            f"{request_id}\0{max(0, int(generation))}\0{call_id}".encode("utf-8")
        ).hexdigest()[:16]
        temporary = result_dir / f".{safe_name[:24]}-{scope_hash}-{uuid4().hex[:8]}.tmp"
        try:
            temporary.write_bytes(encoded)
        except BaseException:
            temporary.unlink(missing_ok=True)
            raise
        return _PreparedToolResult(
            temporary_path=temporary,
            final_path=result_path,
            reference=str(result_path.relative_to(self.agent_dir)).replace("\\", "/"),
            encoded_bytes=len(encoded),
        )

    @staticmethod
    def _discard_prepared_tool_result(prepared: _PreparedToolResult) -> None:
        prepared.temporary_path.unlink(missing_ok=True)
        prepared.final_path.unlink(missing_ok=True)

    def _persist_tool_result(self, tool_name: str, payload: dict[str, Any]) -> tuple[str, int]:
        """Keep complete tool output off-model while retaining a durable task record."""

        prepared = self._prepare_tool_result(tool_name, payload)
        try:
            os.replace(prepared.temporary_path, prepared.final_path)
        except BaseException:
            self._discard_prepared_tool_result(prepared)
            raise
        return prepared.reference, prepared.encoded_bytes

    def _selected_writable_notebook_id(self, requested: str = "") -> str:
        """Resolve the notebook a Pi write action is allowed to update."""

        notebook_id = str(requested or "").strip()
        if notebook_id:
            if self.notebook_ids and notebook_id not in self.notebook_ids:
                raise PermissionError("The requested notebook is not selected in this conversation")
        elif self.notebook_ids:
            notebook_id = self.notebook_ids[0]
        if not notebook_id:
            default_notebook = initialize_notebook(
                self.workspace,
                notebook_id="pi-research-downloads",
                title="研究下载",
                description="由 ScanSci Agent 下载并建立全文索引的文献。",
                root_path=self.workspace.parent / "research-downloads",
                metadata={
                    "library_kind": "files",
                    "created_by": "pi-agent",
                    "purpose": "download-and-index",
                },
            )
            notebook_id = str(default_notebook["notebook_id"])
            self.notebook_ids.append(notebook_id)
        summary = load_workspace_summary(self.workspace, notebook_id=notebook_id)
        if not list(summary.get("notebooks", []) or []):
            raise FileNotFoundError(f"Notebook does not exist: {notebook_id}")
        return notebook_id

    def _download_and_index(self, arguments: dict[str, Any]) -> dict[str, Any]:
        """Run the durable download -> full-text index contract as one Pi tool."""

        raw_identifiers = arguments.get("identifiers")
        if isinstance(raw_identifiers, str):
            identifiers = [raw_identifiers]
        elif isinstance(raw_identifiers, list):
            identifiers = [str(value).strip() for value in raw_identifiers if str(value).strip()]
        else:
            identifier = str(arguments.get("identifier", "")).strip()
            identifiers = [identifier] if identifier else []
        identifiers = list(dict.fromkeys(identifiers))[:20]
        if not identifiers:
            raise ValueError("download_and_index requires at least one DOI or arXiv identifier")

        notebook_id = self._selected_writable_notebook_id(str(arguments.get("notebook_id", "")))
        strategy = str(arguments.get("strategy", "oa_first") or "oa_first")
        timeout = min(600.0, max(30.0, float(arguments.get("timeout_seconds", 180) or 180)))
        workflow_type = "paper_download" if len(identifiers) == 1 else "paper_download_batch"
        store = ResearchRunStore(self.workspace)
        run = store.create_run(
            notebook_id=notebook_id,
            workflow_type=workflow_type,
            title=f"Download and index {len(identifiers)} paper{'s' if len(identifiers) != 1 else ''}",
            input_payload={
                "identifiers": identifiers,
                "strategy": strategy,
                "requested_by": "pi-agent",
            },
            stages=[
                StageSpec("download", "Download full text", "tool", "paper_download"),
                StageSpec("index", "Index full text", "tool", "import_library_files"),
                StageSpec("deliver", "Verify task output", "delivery"),
            ],
            metadata={
                "runtime": "scansci-pi-agent.v1",
                "workflow_orchestration": "durable-high-level-tool",
                "evidence_policy": "strict",
            },
        )
        run_id = str(run["run_id"])
        # Bind every follow-up read/summarize/check action to this task even
        # when download or indexing fails.  Without this, a model could
        # accidentally summarize documents from an older successful task.
        self.active_run_id = run_id
        current_stage = "download"
        try:
            store.begin_run(run_id)
            store.start_stage(run_id, "download")
            # Use the staged batch downloader even for one identifier. It
            # supplies cooperative cancellation and commits only validated
            # PDFs, while its own bounded retry policy avoids a second layer
            # of duplicate network attempts in the Pi bridge.
            downloaded = download_papers(
                identifiers,
                workspace=self.workspace,
                strategy=strategy,
                timeout=timeout,
                cancel_check=self._current_cancel_event().is_set,
            )
            reported_files = [
                str(value)
                for value in list(downloaded.get("files", []) or [])
                if str(value).strip()
            ]
            files = list(
                dict.fromkeys(
                    str(Path(value).resolve())
                    for value in reported_files
                    if Path(value).expanduser().is_file()
                )
            )
            missing_files = [value for value in reported_files if not Path(value).expanduser().is_file()]
            download_failures = [
                {
                    "identifier": str(item.get("identifier", "")),
                    "error": str(item.get("error", "Download failed"))[:500],
                }
                for item in list(downloaded.get("items", []) or [])
                if isinstance(item, dict) and str(item.get("status", "")) == "failed"
            ]
            raw_failures = downloaded.get("failures", [])
            if isinstance(raw_failures, list):
                download_failures.extend(
                    dict(item) if isinstance(item, dict) else {"error": str(item)[:500]}
                    for item in raw_failures
                )
            download_failures.extend(
                {"file": value, "error": "Downloader reported a file that does not exist"}
                for value in missing_files
            )
            successful_identifiers = [
                str(item.get("identifier", "")).strip()
                for item in list(downloaded.get("items", []) or [])
                if isinstance(item, dict)
                and str(item.get("status", "")) == "completed"
                and any(
                    Path(value).expanduser().is_file()
                    for value in list(item.get("files", []) or [])
                    if str(value).strip()
                )
            ]
            if not successful_identifiers and len(identifiers) == 1 and files:
                successful_identifiers = [identifiers[0]]
            elif (
                not list(downloaded.get("items", []) or [])
                and len(files) >= len(identifiers)
            ):
                successful_identifiers = list(identifiers)
            successful_identifiers = list(dict.fromkeys(successful_identifiers))
            requested_count_satisfied = len(successful_identifiers) == len(identifiers)
            downloaded = {
                **downloaded,
                "files": files,
                "missing_files": missing_files,
                "failure_records": download_failures,
                "requested_identifiers": identifiers,
                "successful_identifiers": successful_identifiers,
                "requested_count": len(identifiers),
                "successful_count": len(successful_identifiers),
                "failed_count": max(0, len(identifiers) - len(successful_identifiers)),
                "requested_count_satisfied": requested_count_satisfied,
            }
            store.complete_stage(
                run_id,
                "download",
                summary=f"Downloaded {len(successful_identifiers)}/{len(identifiers)} requested paper(s)",
                output=downloaded,
            )

            current_stage = "index"
            store.start_stage(run_id, "index")
            if not files:
                raise RuntimeError("The download completed without a readable document file")
            imported = import_library_files(
                self.workspace,
                self.root_evidence_db,
                notebook_id=notebook_id,
                file_paths=files,
            )
            import_failures = len(list(imported.get("skipped_files", []) or []))
            imported_count = max(
                int(imported.get("added_files", 0) or 0),
                max(0, len(files) - import_failures),
            )
            goal_complete = requested_count_satisfied and imported_count >= len(successful_identifiers)
            result = {
                **downloaded,
                "ok": goal_complete,
                "partial": not goal_complete,
                "goal_complete": goal_complete,
                "imported": imported,
                "indexed_file_count": imported_count,
                "evidence_status": "indexed_fulltext",
                "notebook_id": notebook_id,
            }
            store.complete_stage(
                run_id,
                "index",
                summary=f"Indexed {imported_count} file(s)",
                output=result,
            )

            current_stage = "deliver"
            store.start_stage(run_id, "deliver")
            artifact = store.create_artifact(
                run_id,
                artifact_type="downloaded_papers" if len(identifiers) > 1 else "downloaded_paper",
                title=str(run["title"]),
                summary=(
                    f"Downloaded and indexed {len(successful_identifiers)}/{len(identifiers)} "
                    "requested paper(s)"
                ),
                payload=result,
                file_path=files[0] if len(files) == 1 else "",
            )
            if goal_complete:
                store.complete_stage(
                    run_id,
                    "deliver",
                    summary="Download and full-text indexing verified",
                    output={"artifact_id": artifact["artifact_id"], "evidence_status": "indexed_fulltext"},
                )
                completed = store.complete_run(run_id, output_artifact_id=str(artifact["artifact_id"]))
            else:
                store.fail_stage(
                    run_id,
                    "deliver",
                    RuntimeError(
                        f"Downloaded and indexed {len(successful_identifiers)}/{len(identifiers)} "
                        "requested paper(s)"
                    ),
                    output_artifact_id=str(artifact["artifact_id"]),
                )
                completed = store.get_run(run_id)
            return {
                "ok": goal_complete,
                "run_id": run_id,
                "status": str(completed.get("status", "")),
                "notebook_id": notebook_id,
                "downloaded_file_count": len(files),
                "requested_count": len(identifiers),
                "successful_count": len(successful_identifiers),
                "failed_count": max(0, len(identifiers) - len(successful_identifiers)),
                "requested_count_satisfied": requested_count_satisfied,
                "partial": not goal_complete,
                "goal_complete": goal_complete,
                "evidence_status": "indexed_fulltext",
                "files": files,
                "failures": download_failures[:20],
                "artifact": {
                    "artifact_id": str(artifact.get("artifact_id", "")),
                    "artifact_type": str(artifact.get("artifact_type", "")),
                    "summary": str(artifact.get("summary", "")),
                },
            }
        except _BatchCancelled as error:
            try:
                store.request_cancel(run_id)
                store.mark_cancelled(run_id, summary="Download cancelled; no staged file was committed")
            except Exception:
                pass
            raise InterruptedError("Download and indexing cancelled") from error
        except Exception as error:
            try:
                store.fail_stage(run_id, current_stage, error)
            except Exception:
                pass
            raise

    @staticmethod
    def _document_map(text: str, *, focus: str = "", maximum_chars: int = 2_400) -> dict[str, Any]:
        """Build a compact, section-aware map from a complete extracted document."""

        normalized = str(text or "").replace("\r\n", "\n").replace("\r", "\n")
        raw_sections: list[str] = []
        for section in re.split(r"\n{2,}", normalized):
            compact = " ".join(section.split())
            if len(compact) <= 2_000:
                raw_sections.append(compact)
                continue
            raw_sections.extend(
                " ".join(value.split())
                for value in re.split(r"(?<=[.!?。！？])\s+", compact)
            )
        paragraphs = [value for value in raw_sections if len(value) >= 50]
        if not paragraphs:
            paragraphs = [" ".join(normalized.split())] if normalized.strip() else []
        focus_tokens = {
            token.casefold()
            for token in re.findall(r"[\w\u4e00-\u9fff-]{3,}", str(focus or ""))
            if token.strip()
        }
        aspects = {
            "research_question": (
                "abstract", "introduction", "objective", "purpose", "aim", "question",
                "摘要", "引言", "目的", "问题",
            ),
            "methods": (
                "method", "materials", "dataset", "data", "experiment", "model", "sample",
                "方法", "材料", "数据", "实验", "模型", "样本",
            ),
            "findings": (
                "result", "finding", "conclusion", "showed", "demonstrate", "significant",
                "结果", "发现", "结论", "表明", "显著",
            ),
            "limitations": (
                "limitation", "caveat", "uncertainty", "future work", "however",
                "局限", "不足", "不确定", "未来研究", "然而",
            ),
        }
        mapped: dict[str, str] = {}
        per_aspect_limit = max(300, maximum_chars // max(1, len(aspects)))
        for aspect, keywords in aspects.items():
            candidates: list[tuple[float, int, str]] = []
            for index, paragraph in enumerate(paragraphs):
                lowered = paragraph.casefold()
                keyword_hits = sum(1 for keyword in keywords if keyword in lowered)
                focus_hits = sum(1 for token in focus_tokens if token in lowered)
                heading_bonus = 2 if len(paragraph) < 180 and keyword_hits else 0
                position_bonus = 0.4 if aspect == "research_question" and index < max(3, len(paragraphs) // 8) else 0
                score = keyword_hits * 3 + focus_hits * 2 + heading_bonus + position_bonus
                if score > 0:
                    candidates.append((score, index, paragraph))
            candidates.sort(key=lambda item: (-item[0], item[1]))
            excerpts: list[str] = []
            used = 0
            for _score, _index, paragraph in candidates:
                excerpt = paragraph[: min(900, per_aspect_limit - used)].strip()
                if not excerpt:
                    continue
                excerpts.append(excerpt)
                used += len(excerpt)
                if used >= per_aspect_limit or len(excerpts) >= 2:
                    break
            mapped[aspect] = "\n".join(excerpts)

        if not mapped["research_question"] and paragraphs:
            mapped["research_question"] = paragraphs[0][:per_aspect_limit]
        if not mapped["findings"] and paragraphs:
            mapped["findings"] = paragraphs[-1][:per_aspect_limit]
        return {
            **mapped,
            "character_count": len(normalized),
            "paragraph_count": len(paragraphs),
            "mapped_aspects": [key for key, value in mapped.items() if value],
        }

    def _summarize_task_documents(self, arguments: dict[str, Any]) -> dict[str, Any]:
        """Map full task documents into a bounded cross-document synthesis input."""

        run, paths = self._task_run_with_documents(str(arguments.get("run_id", "")))
        maximum_files = min(24, max(1, int(arguments.get("max_files", 12) or 12)))
        focus = str(arguments.get("focus", "") or "").strip()
        selected = paths[:maximum_files]
        run_id = str(run.get("run_id", "") or "recent")
        safe_run_id = re.sub(r"[^A-Za-z0-9._-]+", "-", run_id)[:80] or "recent"
        output_dir = self.agent_dir / "task-document-maps" / safe_run_id
        output_dir.mkdir(parents=True, exist_ok=True)
        documents: list[dict[str, Any]] = []
        failures: list[dict[str, str]] = []
        remaining_chars = 34_000
        for path in selected:
            if remaining_chars <= 1_000:
                break
            try:
                stat = path.stat()
                cache_key = hashlib.sha256(
                    f"{path}\0{stat.st_size}\0{stat.st_mtime_ns}".encode("utf-8")
                ).hexdigest()[:24]
                cache_path = output_dir / f"document-{cache_key}.md"
                if cache_path.is_file():
                    text = cache_path.read_text(encoding="utf-8")
                    parser = "cached"
                else:
                    extracted = extract_local_document(path, output_dir=output_dir, parser="auto")
                    extracted_path = Path(str(extracted.get("text_path", ""))).resolve()
                    text = extracted_path.read_text(encoding="utf-8") if extracted_path.is_file() else ""
                    cache_path.write_text(text, encoding="utf-8")
                    parser = str(extracted.get("parser", ""))
                mapped = self._document_map(
                    text,
                    focus=focus,
                    maximum_chars=min(2_400, remaining_chars),
                )
                encoded_size = len(json.dumps(mapped, ensure_ascii=False))
                remaining_chars -= encoded_size
                documents.append(
                    {
                        "name": path.name,
                        "path": str(path),
                        "parser": parser,
                        **mapped,
                    }
                )
            except Exception as error:
                failures.append({"name": path.name, "error": f"{type(error).__name__}: {error}"[:500]})
        return {
            "ok": bool(documents),
            "run_id": run_id,
            "title": str(run.get("title", "")),
            "workflow_type": str(run.get("workflow_type", "")),
            "focus": focus,
            "document_count": len(documents),
            "total_recorded": len(paths),
            "coverage": round(len(documents) / max(1, len(paths)), 4),
            "truncated": len(documents) < len(paths),
            "documents": documents,
            "failures": failures,
            "synthesis_instruction": (
                "Compare the mapped research questions, methods, findings, and limitations across all documents. "
                "Name documents that could not be read and do not treat missing aspects as negative findings."
            ),
        }

    def _check_task_completion(self, arguments: dict[str, Any]) -> dict[str, Any]:
        """Evaluate persisted task state instead of trusting a model completion claim."""

        store = ResearchRunStore(self.workspace)
        requested_run_id = str(arguments.get("run_id", "") or self.active_run_id).strip()
        if requested_run_id:
            run = store.get_run(requested_run_id)
        else:
            recent = store.list_runs(limit=1, archived=None)
            if not recent:
                raise FileNotFoundError("No ScanSci task is available for completion checking")
            run = store.get_run(str(recent[0]["run_id"]))
        stages = list(run.get("stages", []) or [])
        artifact = dict(run.get("output_artifact") or {})
        payload = dict(artifact.get("payload") or {})
        workflow_type = str(run.get("workflow_type", ""))
        files = self._recorded_document_paths(run)
        checks = {
            "run_completed": str(run.get("status", "")) == "completed",
            "stages_completed": bool(stages) and all(
                str(stage.get("status", "")) in {"completed", "skipped"} for stage in stages
            ),
            "artifact_created": bool(artifact.get("artifact_id")),
            "registered_documents": len(files),
        }
        if workflow_type in {"paper_download", "paper_download_batch", "paper_search_download"}:
            checks["fulltext_indexed"] = str(payload.get("evidence_status", "")) == "indexed_fulltext"
            checks["downloaded_files"] = len(files)
            run_input = dict(run.get("input", {}) or {})
            raw_requested_identifiers = (
                payload.get("requested_identifiers")
                or run_input.get("identifiers")
                or ([run_input.get("identifier")] if run_input.get("identifier") else [])
            )
            requested_identifiers = [
                str(value).strip()
                for value in list(raw_requested_identifiers or [])
                if str(value).strip()
            ]
            requested_count = int(payload.get("requested_count", len(requested_identifiers)) or len(requested_identifiers))
            successful_count = int(payload.get("successful_count", len(files)) or 0)
            download_failures = list(
                payload.get("failure_records")
                or payload.get("failures")
                or []
            )
            checks["requested_documents"] = requested_count
            checks["successful_documents"] = successful_count
            checks["requested_count_satisfied"] = bool(
                payload.get(
                    "requested_count_satisfied",
                    requested_count > 0 and successful_count >= requested_count,
                )
            )
            checks["download_failures"] = len(download_failures)
        blockers = [
            key
            for key, value in checks.items()
            if (isinstance(value, bool) and not value)
            or (key == "downloaded_files" and int(value) <= 0)
            or (key == "download_failures" and int(value) > 0)
        ]
        return {
            "ok": not blockers,
            "complete": not blockers,
            "run_id": str(run.get("run_id", "")),
            "workflow_type": workflow_type,
            "status": str(run.get("status", "")),
            "progress": float(run.get("progress", 0.0) or 0.0),
            "checks": checks,
            "blockers": blockers,
            "error": dict(run.get("error", {}) or {}),
            "next_action": "deliver" if not blockers else "inspect the failed or incomplete check and retry that stage",
        }

    def _build_verified_answer(self, arguments: dict[str, Any]) -> dict[str, Any]:
        """Build independently verified answers across every selected evidence store."""

        self._require_evidence_store()
        limit = _bounded_limit(arguments.get("result_limit"), default=12)
        question = str(arguments.get("question", "")).strip()
        available = [path for path in self.evidence_dbs if path.is_file()]
        results: list[dict[str, Any]] = []
        failures: list[dict[str, str]] = []
        for library_index, evidence_db in enumerate(available):
            try:
                payload = answer_question(
                    evidence_db,
                    question,
                    limit=limit,
                    max_quotes=min(8, limit),
                    adequacy_profile="manual",
                    agentic_profile="custom",
                    query_variants=1,
                    max_followup_queries=1,
                    embedding_provider=self.embedding_provider,
                    reranker=self.reranker,
                    filters=self._knowledge_filters(evidence_db, question),
                )
                compact = _compact_verified_answer_for_model(payload)
                compact["library_index"] = library_index
                results.append(compact)
            except Exception as error:
                failures.append(
                    {
                        "library_index": str(library_index),
                        "error": f"{type(error).__name__}: {error}"[:500],
                    }
                )
        if not results:
            detail = failures[0]["error"] if failures else "No selected evidence store could be queried"
            raise RuntimeError(detail)
        scope_summary = self._knowledge_scope_summary(question)
        if len(results) == 1 and not failures:
            if scope_summary is not None:
                results[0]["knowledge_scope"] = scope_summary
            return results[0]

        citations: list[dict[str, Any]] = []
        text_sections: list[str] = []
        for item in results:
            library_index = int(item.get("library_index", 0) or 0)
            reader = dict(item.get("reader_answer", {}) or {})
            text = str(reader.get("text", "")).strip()
            if text:
                text_sections.append(f"[Knowledge base {library_index + 1}]\n{text}")
            for citation in list(reader.get("citations", []) or []):
                citations.append({**dict(citation), "library_index": library_index})
        adequacies = [dict(item.get("adequacy", {}) or {}) for item in results]
        verifications = [dict(item.get("citation_verification", {}) or {}) for item in results]
        return {
            "question": question,
            "federated": True,
            "library_count": len(available),
            "successful_library_count": len(results),
            "reader_answer": {
                "text": "\n\n".join(text_sections)[:12_000],
                "citation_count": len(citations),
                "citations": citations[:24],
            },
            "citation_verification": {
                "passed": all(bool(item.get("passed", False)) for item in verifications),
                "claim_count": sum(int(item.get("claim_count", 0) or 0) for item in verifications),
                "supported_claim_count": sum(
                    int(item.get("supported_claim_count", 0) or 0) for item in verifications
                ),
                "missing_quote_ids": [
                    value
                    for item in verifications
                    for value in list(item.get("missing_quote_ids", []) or [])
                ][:16],
            },
            "answer": {
                "insufficient_evidence": all(
                    bool(dict(item.get("answer", {}) or {}).get("insufficient_evidence", False))
                    for item in results
                ),
                "limitations": [
                    limitation
                    for item in results
                    for limitation in list(dict(item.get("answer", {}) or {}).get("limitations", []) or [])
                ][:12],
            },
            "adequacy": {
                "is_sufficient": any(bool(item.get("is_sufficient", False)) for item in adequacies),
                "quote_count": sum(int(item.get("quote_count", 0) or 0) for item in adequacies),
                "document_count": sum(int(item.get("document_count", 0) or 0) for item in adequacies),
                "followup_reason": "; ".join(
                    str(item.get("followup_reason", "")).strip()
                    for item in adequacies
                    if str(item.get("followup_reason", "")).strip()
                )[:1000],
            },
            "libraries": results,
            "failures": failures,
            **({"knowledge_scope": scope_summary} if scope_summary is not None else {}),
        }

    def _execute_tool(self, name: str, arguments: dict[str, Any]) -> dict[str, Any]:
        if name in {
            "delegate_scientific_agents",
            "list_scientific_agents",
            "collect_scientific_agents",
            "cancel_scientific_agents",
        }:
            if not self.active_run_id or self._scientific_agent_control is None:
                raise PermissionError("Scientific agent controls require a host-bound active durable run")
            sanitized = {
                str(key): value
                for key, value in dict(arguments or {}).items()
                if str(key) not in {"parent_run_id", "run_id"}
            }
            return self._scientific_agent_control(name, self.active_run_id, sanitized)
        if name in {"create_document", "create_pdf", "create_spreadsheet", "create_presentation", "compile_latex"}:
            plugin_by_tool = {
                "create_document": "documents",
                "create_pdf": "pdf",
                "create_spreadsheet": "spreadsheets",
                "create_presentation": "presentations",
                "compile_latex": "latex",
            }
            plugin_id = plugin_by_tool[name]
            plugins = list(load_settings(self.workspace).get("plugins", []) or [])
            plugin = next((item for item in plugins if str(item.get("id", "")) == plugin_id), None)
            if not plugin or not plugin.get("enabled") or plugin.get("uninstalled"):
                raise PermissionError(f"Built-in plugin is disabled: {plugin_id}")
            return execute_artifact_tool(name, arguments, workspace=self.workspace)
        if name in {
            "zotero_search",
            "zotero_fulltext",
            "zotero_attachment",
            "zotero_export_bibtex",
            "zotero_citations",
        }:
            plugins = list(load_settings(self.workspace).get("plugins", []) or [])
            plugin = next((item for item in plugins if str(item.get("id", "")) == "zotero"), None)
            if not plugin or not plugin.get("enabled") or plugin.get("uninstalled"):
                raise PermissionError("Built-in plugin is disabled: zotero")
        if name == "inspect_workspace":
            return load_workspace_summary(self.workspace, notebook_id=str(arguments.get("notebook_id", "")))
        if name == "inspect_available_tools":
            snapshot = capability_snapshot(workspace=self.workspace, evidence_db=self.evidence_db)
            snapshot["pi_high_level_tools"] = [
                {
                    "id": "agent_reach",
                    "status": "ready",
                    "description": "Read and search public internet channels through built-in zero-install Agent Reach routes.",
                },
                {
                    "id": "browser_access",
                    # A probe is a side-effect-free localhost health check; do
                    # not advertise a ready tool that would fail on first use.
                    "status": (
                        "ready"
                        if browser_access_status(timeout=1.5).get("ready")
                        else "requires-setup"
                    ),
                    "description": "Read rendered public pages through the bundled web-access CDP bridge when public readers are insufficient.",
                },
                {
                    "id": "download_and_index",
                    "status": "ready" if self.notebook_ids else "needs-selected-library",
                    "description": "Download DOI/arXiv documents, register a durable task, and index readable full text.",
                },
                {
                    "id": "summarize_documents",
                    "status": "ready",
                    "description": "Map complete task documents by question, methods, findings, and limitations.",
                },
                {
                    "id": "check_task_completion",
                    "status": "ready",
                    "description": "Verify persisted stages, artifacts, documents, and full-text indexing.",
                },
                {
                    "id": "build_verified_answer",
                    "status": "ready" if any(path.is_file() for path in self.evidence_dbs) else "needs-data",
                    "description": f"Build verified answers across {len(self.evidence_dbs)} selected evidence store(s).",
                },
            ]
            snapshot["selected_notebook_ids"] = list(self.notebook_ids)
            snapshot["selected_evidence_store_count"] = len(
                [path for path in self.evidence_dbs if path.is_file()]
            )
            return snapshot
        if name == "read_task_documents":
            return self._read_task_documents(arguments)
        if name == "download_and_index":
            return self._download_and_index(arguments)
        if name == "summarize_documents":
            return self._summarize_task_documents(arguments)
        if name == "check_task_completion":
            return self._check_task_completion(arguments)
        if name in {"search_local_evidence", "kb_search", "zotero_search"}:
            limit = _bounded_limit(arguments.get("result_limit"), default=8)
            hits: list[dict[str, Any]] = []
            selected_notebooks = self._selected_notebooks()
            plugins = list(load_settings(self.workspace).get("plugins", []) or [])
            zotero_plugin = next((item for item in plugins if str(item.get("id", "")) == "zotero"), None)
            zotero_enabled = bool(zotero_plugin and zotero_plugin.get("enabled") and not zotero_plugin.get("uninstalled"))
            zotero_notebooks = [
                notebook
                for notebook in selected_notebooks
                if zotero_enabled
                if str(dict(notebook.get("metadata", {}) or {}).get("library_kind", "")) == "zotero"
                or bool(dict(notebook.get("metadata", {}) or {}).get("zotero"))
            ]
            search_evidence = name == "search_local_evidence" and (
                not selected_notebooks or any(list(notebook.get("sources", []) or []) for notebook in selected_notebooks)
            )
            query_text = str(arguments.get("query", ""))
            if search_evidence:
                for library_index, evidence_db in enumerate(self.evidence_dbs):
                    if not evidence_db.is_file():
                        continue
                    library_hits = search_evidence_store(
                        evidence_db,
                        query_text,
                        limit=limit,
                        context_mode="sentence",
                        embedding_provider=self.embedding_provider,
                        reranker=self.reranker,
                        filters=self._knowledge_filters(evidence_db, query_text),
                    )
                    for hit in library_hits:
                        hits.append({**hit, "library_index": library_index})
            hits.sort(key=lambda hit: float(hit.get("score", 0.0) or 0.0), reverse=True)
            hits = hits[:limit]
            zotero_results = []
            if zotero_enabled and (name in {"kb_search", "zotero_search"} or zotero_notebooks):
                collection_key = str(arguments.get("collection_key", ""))
                zotero_scope = self._knowledge_scope_summary(query_text) or self.knowledge_scope
                zotero_results.append(
                    filter_zotero_result(
                        search_zotero_library(
                            query_text,
                            limit=limit,
                            collection_key=collection_key,
                            include_fulltext=bool(arguments.get("include_fulltext", True)),
                        ),
                        zotero_scope,
                    )
                )
            payload = {
                "query": str(arguments.get("query", "")),
                "count": len(hits) + sum(int(result.get("count", 0) or 0) for result in zotero_results),
                "library_count": len([path for path in self.evidence_dbs if path.is_file()]),
                "hits": [
                    {**_compact_evidence_hit(hit), "library_index": int(hit.get("library_index", 0) or 0)}
                    for hit in hits
                ],
                "zotero": zotero_results,
            }
            scope_summary = self._knowledge_scope_summary(query_text)
            if scope_summary is not None:
                payload["knowledge_scope"] = scope_summary
            if not payload["count"] and not zotero_results:
                self._require_evidence_store()
            return payload
        if name == "zotero_status":
            status = dict(zotero_status())
            plugins = list(load_settings(self.workspace).get("plugins", []) or [])
            plugin = next((item for item in plugins if str(item.get("id", "")) == "zotero"), None)
            status["plugin_enabled"] = bool(plugin and plugin.get("enabled") and not plugin.get("uninstalled"))
            return status
        if name == "zotero_fulltext":
            return zotero_fulltext(
                str(arguments.get("attachment_key", "")),
                max_chars=min(80_000, max(1_000, int(arguments.get("max_chars", 40_000) or 40_000))),
            )
        if name == "zotero_attachment":
            return zotero_file_url(str(arguments.get("attachment_key", "")))
        if name == "zotero_export_bibtex":
            return zotero_export_bibtex(
                item_key=str(arguments.get("item_key", "")),
                limit=_bounded_limit(arguments.get("result_limit"), default=50),
            )
        if name == "zotero_citations":
            return zotero_formatted_citations(
                style=str(arguments.get("style", "apa")),
                limit=_bounded_limit(arguments.get("result_limit"), default=30),
            )
        if name in {"obsidian_status", "obsidian_search", "obsidian_read", "obsidian_backlinks"}:
            vaults = self._selected_obsidian_vaults()
            if not vaults:
                raise FileNotFoundError("No linked Obsidian vault is selected for this conversation")
            if name == "obsidian_status":
                return {"ok": True, "count": len(vaults), "vaults": [obsidian_status(vault) for vault in vaults]}
            if name == "obsidian_search":
                results = [
                    search_obsidian_vault(
                        vault,
                        str(arguments.get("query", "")),
                        limit=_bounded_limit(arguments.get("result_limit"), default=10),
                    )
                    for vault in vaults
                ]
                return {
                    "ok": True,
                    "count": sum(int(result.get("count", 0) or 0) for result in results),
                    "vaults": results,
                }
            vault = vaults[0]
            note_path = str(arguments.get("note_path", ""))
            if name == "obsidian_read":
                return read_obsidian_note(
                    vault,
                    note_path,
                    max_chars=min(100_000, max(1_000, int(arguments.get("max_chars", 40_000) or 40_000))),
                )
            return obsidian_backlinks(
                vault,
                note_path,
                limit=_bounded_limit(arguments.get("result_limit"), default=30),
            )
        if name == "build_verified_answer":
            return self._build_verified_answer(arguments)
        if name == "verify_doi":
            return verify_doi_metadata(
                str(arguments.get("doi", "")),
                expected_title=str(arguments.get("expected_title", "")),
            )
        if name == "discover_papers":
            providers = arguments.get("providers")
            provider_names = [str(value) for value in providers] if isinstance(providers, list) else None
            plan = plan_academic_search(
                str(arguments.get("query", "")),
                explicit_providers=provider_names,
            )
            year_from = arguments.get("year_from")
            try:
                normalized_year = int(year_from) if year_from not in {None, ""} else plan.get("year_from")
            except (TypeError, ValueError):
                normalized_year = plan.get("year_from")
            requested_limit = _bounded_limit(
                arguments.get("result_limit", arguments.get("limit")),
                default=8,
            )
            result = search_academic_papers(
                str(plan["topic"]),
                query_variants=list(plan.get("query_variants", []) or []),
                required_terms=list(plan.get("required_terms", []) or []),
                limit=requested_limit,
                per_source=_bounded_limit(arguments.get("per_source"), default=min(8, requested_limit)),
                provider_names=list(plan.get("providers", []) or []),
                year_from=normalized_year,
                embedding_provider=self.embedding_provider,
                reranker=self.reranker,
            )
            result["search_plan"] = plan
            if isinstance(getattr(self._worker_context, "run_context", None), _RunContext):
                # The owning event thread publishes this complete provider payload
                # together with its bounded model view.  A worker must not create
                # a durable file that cancellation can strand before tool.result.
                return result
            result_reference, result_bytes = self._persist_tool_result("discover-papers", result)
            return _compact_academic_search_result_for_model(
                result,
                requested_limit=requested_limit,
                result_reference=result_reference,
                result_bytes=result_bytes,
            )
        if name == "search_web":
            return search_public_web(
                str(arguments.get("query", "")),
                limit=_bounded_limit(arguments.get("result_limit"), default=8),
            )
        if name == "agent_reach":
            return run_agent_reach(
                str(arguments.get("operation", "")),
                target=str(arguments.get("target", "")),
                query=str(arguments.get("query", "")),
                channel=str(arguments.get("channel", "auto")),
                limit=_bounded_limit(arguments.get("limit"), default=8),
                timeout=min(60.0, max(3.0, float(arguments.get("timeout_seconds", 30) or 30))),
            )
        if name == "browser_access":
            operation = str(arguments.get("operation", "")).strip().lower()
            timeout = min(60.0, max(5.0, float(arguments.get("timeout_seconds", 30) or 30)))
            if operation == "status":
                return browser_access_status(timeout=min(3.0, timeout))
            if operation == "read":
                return browser_access_read(
                    str(arguments.get("target", "")),
                    timeout=timeout,
                )
            raise ValueError("browser_access operation must be status or read")
        if name == "search_journal":
            return search_journals(
                str(arguments.get("query", "")),
                limit=_bounded_limit(arguments.get("result_limit"), default=8),
            )
        if name == "audit_references":
            mode = "references" if str(arguments.get("mode", "references")) == "references" else "full"
            return analyze_references(str(arguments.get("text", "")), mode=mode)
        if name == "build_presentation_outline":
            summary = load_workspace_summary(self.workspace, notebook_id=str(arguments.get("notebook_id", "")))
            notebooks = list(summary.get("notebooks", []) or [])
            if not notebooks:
                raise FileNotFoundError("The current ScanSci workspace has no usable notebook")
            return build_ppt_outline(
                dict(notebooks[0]),
                topic=str(arguments.get("topic", "")),
                template_id=str(arguments.get("template_id", "")),
            )
        if name == "edit_section":
            return _edit_section(
                str(arguments.get("file_path", "")),
                str(arguments.get("old_string", "")),
                str(arguments.get("new_string", "")),
                workspace=self.workspace,
            )
        if name == "edit_slide":
            return _edit_slide_text(
                str(arguments.get("pptx_path", "")),
                int(arguments.get("slide_index", 0)),
                str(arguments.get("old_string", "")),
                str(arguments.get("new_string", "")),
                workspace=self.workspace,
            )
        if name == "self_assess":
            context = getattr(self._worker_context, "run_context", None)
            if isinstance(context, _RunContext):
                with context.history_lock:
                    return _build_self_assessment(list(context.history))
            with self._tool_history_lock:
                return _build_self_assessment(list(self._tool_history))
        raise ValueError(f"Unsupported ScanSci Pi tool: {name}")

    def _selected_notebooks(self) -> list[dict[str, Any]]:
        if not self.notebook_ids:
            return []
        selected: list[dict[str, Any]] = []
        for notebook_id in self.notebook_ids:
            summary = load_workspace_summary(self.workspace, notebook_id=notebook_id)
            selected.extend(dict(notebook) for notebook in list(summary.get("notebooks", []) or []))
        return selected

    def _knowledge_scope_resolution(
        self,
        evidence_db: str | Path,
        question: str = "",
    ) -> dict[str, Any]:
        """Resolve explicit or automatic tag scope for one selected index."""

        requested = dict(self.knowledge_scope or {})
        requested_type = str(requested.get("type", "")).strip()
        if requested_type not in {"", "zotero-tag"}:
            return {"active": False, "status": "inactive"}
        database = Path(evidence_db).resolve()
        cache_question = "" if requested_type == "zotero-tag" else str(question or "").strip().casefold()
        cache_key = f"{database}|{cache_question}"
        if cache_key in self._knowledge_scope_cache:
            return dict(self._knowledge_scope_cache[cache_key])
        selected = self._selected_notebooks()
        requested_id = str(requested.get("notebook_id", "")).strip()
        notebook = next(
            (
                item
                for item in selected
                if requested_id and str(item.get("notebook_id", "")).strip() == requested_id
            ),
            None,
        )
        if notebook is None:
            notebook = next(
                (
                    item
                    for item in selected
                    if str(dict(item.get("metadata", {}) or {}).get("library_kind", "")).strip().lower() == "zotero"
                ),
                None,
            )
        if notebook is None:
            return {"active": False, "status": "unavailable"}
        notebook_databases = {
            str(Path(str(dict(source).get("evidence_db_path", ""))).resolve())
            for source in list(notebook.get("sources", []) or [])
            if str(dict(source).get("evidence_db_path", "")).strip()
        }
        if notebook_databases and str(database) not in notebook_databases:
            result = {"active": False, "status": "other-library"}
            self._knowledge_scope_cache[cache_key] = result
            return dict(result)
        if not notebook_databases and len(self.evidence_dbs) > 1 and database != self.evidence_dbs[0].resolve():
            result = {"active": False, "status": "other-library"}
            self._knowledge_scope_cache[cache_key] = result
            return dict(result)
        library_kind = str(dict(notebook.get("metadata", {}) or {}).get("library_kind", "")).strip().lower()
        if requested_type == "zotero-tag":
            # Keep the old explicit scope format readable for old callers, but
            # the UI no longer creates it. Normal Zotero retrieval uses the
            # tag sidecar as a soft hybrid-ranking signal instead of a filter.
            result = resolve_zotero_tag_scope(notebook, database, requested)
        elif library_kind == "zotero":
            result = {"active": False, "status": "automatic-ranking"}
        else:
            result = {"active": False, "status": "no-query" if not str(question or "").strip() else "not-zotero"}
        result["notebook_id"] = str(notebook.get("notebook_id", ""))
        self._knowledge_scope_cache[cache_key] = dict(result)
        return dict(result)

    def _knowledge_filters(self, evidence_db: str | Path, question: str = "") -> dict[str, Any]:
        resolution = self._knowledge_scope_resolution(evidence_db, question)
        if not bool(resolution.get("active")):
            return {}
        return {"doc_ids": list(resolution.get("doc_ids", []) or [])}

    def _knowledge_scope_summary(self, question: str = "") -> dict[str, Any] | None:
        for evidence_db in self.evidence_dbs:
            resolution = self._knowledge_scope_resolution(evidence_db, question)
            if resolution.get("type") == "zotero-tag" and resolution.get("status") == "applied":
                return resolution
        return None

    def _selected_obsidian_vaults(self) -> list[str]:
        vaults: list[str] = []
        for notebook in self._selected_notebooks():
            metadata = dict(notebook.get("metadata", {}) or {})
            if str(metadata.get("library_kind", "")) != "obsidian":
                continue
            root = str(metadata.get("library_root", "") or notebook.get("root_path", "")).strip()
            if root and root not in vaults:
                vaults.append(root)
        return vaults

    def _require_evidence_store(self) -> None:
        if not any(path.is_file() for path in self.evidence_dbs):
            raise FileNotFoundError(f"Evidence store does not exist: {self.evidence_db}")


def _bounded_limit(value: Any, *, default: int) -> int:
    try:
        return max(1, min(20, int(value if value is not None else default)))
    except (TypeError, ValueError):
        return default


def _compact_academic_search_result_for_model(
    payload: dict[str, Any],
    *,
    requested_limit: int,
    result_reference: str,
    result_bytes: int,
) -> dict[str, Any]:
    """Return a bounded discovery view while preserving download identifiers."""

    source_items = [dict(item) for item in list(payload.get("items", []) or []) if isinstance(item, dict)]
    detailed_items: list[dict[str, Any]] = []
    download_identifiers: list[str] = []
    seen_identifiers: set[str] = set()
    for item in source_items:
        doi = str(item.get("doi", "") or "").strip()
        arxiv_id = str(item.get("arxiv_id", "") or "").strip()
        identifier = str(item.get("identifier", "") or doi or arxiv_id).strip()
        if identifier and identifier.casefold() not in seen_identifiers:
            seen_identifiers.add(identifier.casefold())
            download_identifiers.append(identifier)
        if len(detailed_items) >= 8:
            continue
        authors_value = item.get("authors")
        authors = (
            [str(value).strip() for value in authors_value if str(value).strip()]
            if isinstance(authors_value, list)
            else [str(authors_value).strip()] if str(authors_value or "").strip() else []
        )
        sources_value = item.get("sources")
        sources = (
            [str(value).strip() for value in sources_value if str(value).strip()]
            if isinstance(sources_value, list)
            else [str(sources_value).strip()] if str(sources_value or "").strip() else []
        )
        compact = {
            "title": " ".join(str(item.get("title", "")).split())[:240],
            "identifier": identifier,
            "doi": doi,
            "arxiv_id": arxiv_id,
            "year": item.get("year"),
            "venue": " ".join(str(item.get("venue", "")).split())[:160],
            "authors": authors[:6],
            "url": str(item.get("url", "") or "")[:500],
            "oa_url": str(item.get("oa_url", "") or "")[:500],
            "citation_count": int(item.get("citation_count", 0) or 0),
            "sources": sources[:6],
            "score": round(float(item.get("score", 0.0) or 0.0), 6),
            "abstract_excerpt": " ".join(str(item.get("abstract", "")).split())[:600],
            "discovery_only": True,
        }
        detailed_items.append({key: value for key, value in compact.items() if value not in ("", None, [])})

    provider_errors = {
        str(name): str(error)[:240]
        for name, error in list(dict(payload.get("provider_errors", {}) or {}).items())[:5]
    }
    total_count = int(payload.get("count", len(source_items)) or len(source_items))
    return {
        "query": str(payload.get("query", "")),
        "query_variants": list(payload.get("query_variants", []) or [])[:3],
        "search_plan": {
            key: value
            for key, value in dict(payload.get("search_plan", {}) or {}).items()
            if key in {"topic", "normalized_topic", "domain", "providers", "query_variants", "required_terms", "planner"}
        },
        "items": detailed_items,
        "count": len(detailed_items),
        "total_count": total_count,
        "returned_item_count": len(detailed_items),
        "requested_result_limit": requested_limit,
        "candidate_count": int(payload.get("candidate_count", total_count) or total_count),
        "deduplicated_count": int(payload.get("deduplicated_count", total_count) or total_count),
        "providers_requested": list(payload.get("providers_requested", []) or [])[:12],
        "providers_succeeded": list(payload.get("providers_succeeded", []) or [])[:12],
        "provider_counts": dict(payload.get("provider_counts", {}) or {}),
        "provider_errors": provider_errors,
        "quality_gate": {
            key: value
            for key, value in dict(payload.get("quality_gate", {}) or {}).items()
            if key in {"status", "accepted_count", "rejected_count", "candidate_count", "reason"}
        },
        "year_from": payload.get("year_from"),
        "latency_ms": payload.get("latency_ms"),
        "evidence_status": str(payload.get("evidence_status", "discovery_leads")),
        "evidence_notice": str(
            payload.get(
                "evidence_notice",
                "Bibliographic results are discovery leads until lawful full text is indexed and verified.",
            )
        )[:500],
        "download_identifiers": download_identifiers[:12],
        "next_action": (
            "Pass download_identifiers to download_and_index. "
            "Then call summarize_documents and check_task_completion."
        ),
        "full_result_reference": result_reference,
        "full_result_bytes": result_bytes,
        "full_result_note": "The complete provider payload is retained by ScanSci and is not injected into model context.",
    }


def _compact_evidence_hit(hit: dict[str, Any]) -> dict[str, Any]:
    return {
        "evidence_id": str(hit.get("evidence_id", "")),
        "doc_id": str(hit.get("doc_id", "")),
        "paper": str(hit.get("paper", "")),
        "doi": str(hit.get("doi", "")),
        "section": str(hit.get("section", "")),
        "html_anchor": str(hit.get("html_anchor", "")),
        "text": " ".join(str(hit.get("text", "")).split())[:1600],
        "score": round(float(hit.get("score", 0.0) or 0.0), 6),
    }


def _compact_verified_answer_for_model(payload: dict[str, Any]) -> dict[str, Any]:
    reader = dict(payload.get("reader_answer", {}) or {})
    citations = []
    for item in list(reader.get("citations", []) or [])[:8]:
        if not isinstance(item, dict):
            continue
        citations.append(
            {
                "citation_id": str(item.get("citation_id", "")),
                "evidence_id": str(item.get("evidence_id", "")),
                "doc_id": str(item.get("doc_id", "")),
                "paper": str(item.get("paper", "")),
                "doi": str(item.get("doi", "")),
                "section": str(item.get("section", "")),
                "html_anchor": str(item.get("html_anchor", "")),
                "exact_quote": " ".join(str(item.get("exact_quote", "")).split())[:700],
                "source_href": str(item.get("source_href", "")),
            }
        )
    answer = dict(payload.get("answer", {}) or {})
    adequacy = dict(payload.get("adequacy", {}) or {})
    verification = dict(payload.get("citation_verification", {}) or {})
    return {
        "question": str(payload.get("question", "")),
        "reader_answer": {
            "text": str(reader.get("text", ""))[:3000],
            "citation_count": int(reader.get("citation_count", len(citations)) or 0),
            "citations": citations,
        },
        "citation_verification": {
            "passed": bool(verification.get("passed", False)),
            "claim_count": int(verification.get("claim_count", 0) or 0),
            "supported_claim_count": int(verification.get("supported_claim_count", 0) or 0),
            "missing_quote_ids": list(verification.get("missing_quote_ids", []) or [])[:8],
        },
        "answer": {
            "insufficient_evidence": bool(answer.get("insufficient_evidence", False)),
            "limitations": list(answer.get("limitations", []) or [])[:6],
        },
        "adequacy": {
            "is_sufficient": bool(adequacy.get("is_sufficient", False)),
            "quote_count": int(adequacy.get("quote_count", 0) or 0),
            "document_count": int(adequacy.get("document_count", 0) or 0),
            "followup_reason": str(adequacy.get("followup_reason", ""))[:500],
        },
    }


def _edit_section(file_path: str, old_string: str, new_string: str, *, workspace: Path) -> dict[str, Any]:
    """Replace the first occurrence of ``old_string`` with ``new_string`` in a
    text-based artifact (docx/txt/md). ZCode-style exact-string editing.
    Returns ``{"ok": true, "matched": true, ...}`` on success.
    """
    if not file_path or not old_string:
        return {"ok": False, "error": "file_path and old_string are required", "matched": False}
    path = Path(file_path)
    if not path.is_absolute():
        path = (workspace.parent / file_path).resolve()
    if not path.exists():
        return {"ok": False, "error": f"文件不存在：{path}", "matched": False}
    try:
        original = path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": f"无法读取文件：{exc}", "matched": False}
    if old_string not in original:
        return {"ok": False, "error": "old_string not found in file", "matched": False}
    updated = original.replace(old_string, new_string, 1)
    checkpoint_id = ""
    try:
        checkpoint = CheckpointStore(workspace.parent).begin(turn_id="edit_section", label=path.name)
        CheckpointStore(workspace.parent).capture(checkpoint.checkpoint_id, path)
        checkpoint_id = checkpoint.checkpoint_id
    except CheckpointError:
        # Preserve the legacy edit behaviour for artifacts outside the local
        # workspace; in-workspace edits are checkpointed automatically.
        checkpoint_id = ""
    path.write_text(updated, encoding="utf-8")
    return {"ok": True, "matched": True, "file_path": str(path), **({"checkpoint_id": checkpoint_id} if checkpoint_id else {})}


def _edit_slide_text(
    pptx_path: str,
    slide_index: int,
    old_string: str,
    new_string: str,
    *,
    workspace: Path | None = None,
) -> dict[str, Any]:
    """Replace text in a specific slide of an existing PPTX file.

    Opens the file, finds the slide (1-based index), walks all text runs in all
    shapes, and replaces the first matching occurrence. Saves in-place.
    """
    import pptx

    if not pptx_path or not old_string:
        return {"ok": False, "error": "pptx_path and old_string are required", "matched": False}
    path = Path(pptx_path)
    if not path.exists():
        return {"ok": False, "error": f"文件不存在：{path}", "matched": False}
    if slide_index < 1:
        return {"ok": False, "error": "slide_index must be >= 1", "matched": False}
    try:
        presentation = pptx.Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": f"无法打开 PPTX：{exc}", "matched": False}
    slides = presentation.slides
    if slide_index > len(slides):
        return {"ok": False, "error": f"幻灯片索引 {slide_index} 超出范围（共 {len(slides)} 页）", "matched": False}
    slide = slides[slide_index - 1]
    matched = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if old_string in run.text:
                    run.text = run.text.replace(old_string, new_string, 1)
                    matched = True
                    break
            if matched:
                break
        if matched:
            break
    if not matched:
        return {"ok": False, "error": f"old_string not found on slide {slide_index}", "matched": False}
    checkpoint_id = ""
    if workspace is not None:
        try:
            checkpoint = CheckpointStore(workspace.parent).begin(turn_id="edit_slide", label=path.name)
            CheckpointStore(workspace.parent).capture(checkpoint.checkpoint_id, path)
            checkpoint_id = checkpoint.checkpoint_id
        except CheckpointError:
            checkpoint_id = ""
    try:
        presentation.save(str(path))
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "error": f"无法保存 PPTX：{exc}", "matched": True}
    return {
        "ok": True,
        "matched": True,
        "slide_index": slide_index,
        "pptx_path": str(path),
        **({"checkpoint_id": checkpoint_id} if checkpoint_id else {}),
    }


# -- L1: search auto-retry ---------------------------------------------------

_SEARCH_TOOLS: set[str] = {"discover_papers", "search_local_evidence", "kb_search", "search_web", "search_journal"}


def _search_result_is_empty(result: dict[str, Any]) -> bool:
    """Detect an empty or near-empty search result across all known schemas."""
    for key in ("items", "hits", "results"):
        value = result.get(key)
        if isinstance(value, list) and len(value) == 0:
            return True
    for key in ("count", "total"):
        value = result.get(key)
        if isinstance(value, int) and value == 0:
            return True
    return False


def _widen_search_query(arguments: dict[str, Any]) -> dict[str, Any]:
    """Broaden a search query that returned zero results.

    Strategies: drop quotes (exact-phrase matching is narrow), shorten the
    query to its first 2 words, double the result limit.
    """
    args = dict(arguments)
    for key in ("query", "question", "q"):
        value = str(args.get(key, "")).strip()
        if value:
            no_quotes = value.replace('"', "").replace("'", "")
            if no_quotes != value:
                args[key] = no_quotes
            elif len(value.split()) > 4:
                args[key] = " ".join(value.split()[:4])
            break
    for key in ("limit", "max_results"):
        current = int(args.get(key, 10) or 10)
        args[key] = min(current * 2, 100)
    for key in ("year_from",):
        if key in args:
            args[key] = max(1900, int(args.get(key, 2000)) - 10)
    return args


def _retry_empty_search(
    name: str,
    arguments: dict[str, Any],
    first_result: dict[str, Any],
    *,
    executor: Any = None,
) -> dict[str, Any] | None:
    """If a search returned zero results, broaden and retry once.

    Returns the retry result on success, None if no retry was needed.
    """
    if name not in _SEARCH_TOOLS:
        return None
    if not isinstance(first_result, dict):
        return None
    if not _search_result_is_empty(first_result):
        return None
    if executor is None:
        return None
    broadened = _widen_search_query(arguments)
    if broadened == arguments:
        return None  # nothing to change
    try:
        return executor(name, broadened)
    except Exception:  # noqa: BLE001 — retry failed, keep original result
        return None


# -- L4: self-assess meta-tool -----------------------------------------------

_RESULT_KEYS = ("count", "total", "items", "hits", "results", "ok", "passed", "status", "matched", "message")


def _summarize_tool_result(name: str, result: Any) -> str:
    if not isinstance(result, dict):
        return str(result)[:120]
    parts: list[str] = []
    for key in _RESULT_KEYS:
        value = result.get(key)
        if isinstance(value, int):
            parts.append(f"{key}={value}")
        elif isinstance(value, list):
            parts.append(f"{key}=[{len(value)}]")
        elif isinstance(value, bool):
            parts.append(f"{key}={value}")
        elif isinstance(value, str) and len(value) < 80:
            parts.append(f"{key}={value}")
    return f"{name}: " + ", ".join(parts[:3]) if parts else f"{name}: done"


def _compact_workspace_result_for_model(result: dict[str, Any]) -> dict[str, Any]:
    """Return counts and notebook identities, never the complete source inventory."""

    notebooks: list[dict[str, Any]] = []
    for item in list(result.get("notebooks", []) or [])[:12]:
        if not isinstance(item, dict):
            continue
        counts = dict(item.get("counts", {}) or {})
        metadata = dict(item.get("metadata", {}) or {})
        notebooks.append(
            {
                "notebook_id": str(item.get("notebook_id", "")),
                "title": str(item.get("title", "")),
                "description": str(item.get("description", ""))[:400],
                "library_kind": str(metadata.get("library_kind", "")),
                "counts": {
                    "sources": int(counts.get("sources", 0) or 0),
                    "notes": int(counts.get("notes", 0) or 0),
                    "layers": int(counts.get("layers", 0) or 0),
                },
            }
        )
    return {
        "status": str(result.get("status", "ready") or "ready"),
        "workspace": str(result.get("workspace", "")),
        "counts": dict(result.get("counts", {}) or {}),
        "notebooks": notebooks,
        "_scansci_compacted": True,
        "_notice": "The full source inventory is intentionally kept out of model context.",
    }


def _compact_json_value(
    value: Any,
    *,
    depth: int = 0,
    max_depth: int = 5,
    max_items: int = 10,
    max_keys: int = 32,
    max_string: int = 900,
) -> Any:
    if depth >= max_depth:
        if isinstance(value, dict):
            return {"_omitted_keys": len(value)}
        if isinstance(value, (list, tuple, set)):
            return {"_omitted_items": len(value)}
        return str(value)[:max_string]
    if isinstance(value, dict):
        items = list(value.items())
        compacted = {
            str(key): _compact_json_value(
                item,
                depth=depth + 1,
                max_depth=max_depth,
                max_items=max_items,
                max_keys=max_keys,
                max_string=max_string,
            )
            for key, item in items[:max_keys]
        }
        if len(items) > max_keys:
            compacted["_omitted_keys"] = len(items) - max_keys
        return compacted
    if isinstance(value, (list, tuple, set)):
        items = list(value)
        compacted = [
            _compact_json_value(
                item,
                depth=depth + 1,
                max_depth=max_depth,
                max_items=max_items,
                max_keys=max_keys,
                max_string=max_string,
            )
            for item in items[:max_items]
        ]
        if len(items) > max_items:
            compacted.append({"_omitted_items": len(items) - max_items})
        return compacted
    if isinstance(value, str):
        if len(value) <= max_string:
            return value
        return value[:max_string] + f"… [omitted {len(value) - max_string} chars]"
    if value is None or isinstance(value, (bool, int, float)):
        return value
    return str(value)[:max_string]


def _redact_tool_value(value: Any, *, key: str = "", depth: int = 0) -> Any:
    """Remove credentials from tool results before model or disk persistence."""

    if re.search(
        r"(?:^|[_-])(?:api[_-]?key|access[_-]?token|refresh[_-]?token|password|secret|authorization)(?:$|[_-])",
        key,
        re.IGNORECASE,
    ):
        return "[REDACTED]"
    if depth > 12:
        return "[TRUNCATED]"
    if isinstance(value, str):
        provider_key = str(os.getenv("SCANSCIPI_PROVIDER_KEY", "") or "")
        text = value.replace(provider_key, "[REDACTED]") if provider_key else value
        text = re.sub(r"\bBearer\s+[A-Za-z0-9._~+/=-]+", "Bearer [REDACTED]", text, flags=re.IGNORECASE)
        text = re.sub(r"\bsk-[A-Za-z0-9_-]{12,}\b", "[REDACTED]", text)
        return re.sub(
            r"([?&](?:api[_-]?key|access[_-]?token|token|password|secret)=)[^&\s]+",
            r"\1[REDACTED]",
            text,
            flags=re.IGNORECASE,
        )
    if isinstance(value, list):
        return [_redact_tool_value(item, depth=depth + 1) for item in value]
    if isinstance(value, dict):
        return {
            str(nested_key): _redact_tool_value(item, key=str(nested_key), depth=depth + 1)
            for nested_key, item in value.items()
        }
    return value


def _bounded_tool_result_for_model(
    name: str,
    result: Any,
    *,
    max_bytes: int = _MAX_MODEL_TOOL_RESULT_BYTES,
) -> tuple[dict[str, Any], dict[str, int | bool]]:
    """Bound every native tool result before it can enter a paid model context."""

    safe = _redact_tool_value(_json_safe(result))
    original_bytes = len(json.dumps(safe, ensure_ascii=False).encode("utf-8"))
    if name == "inspect_workspace" and isinstance(safe, dict):
        safe = _compact_workspace_result_for_model(safe)
    if not isinstance(safe, dict):
        safe = {"result": safe}
    model_bytes = len(json.dumps(safe, ensure_ascii=False).encode("utf-8"))
    intentionally_compacted = name == "inspect_workspace" and model_bytes != original_bytes
    if model_bytes <= max_bytes:
        return safe, {
            "original_bytes": original_bytes,
            "model_bytes": model_bytes,
            "truncated": intentionally_compacted,
            "persist_full": False,
        }

    compacted = _compact_json_value(safe)
    if not isinstance(compacted, dict):
        compacted = {"result": compacted}
    compacted.update(
        {
            "_scansci_truncated": True,
            "_original_bytes": original_bytes,
            "_notice": "Tool output exceeded the model-context budget; a bounded preview is shown.",
        }
    )
    encoded = json.dumps(compacted, ensure_ascii=False).encode("utf-8")
    if len(encoded) > max_bytes:
        preview_budget = max(800, max_bytes - 1_000)
        compacted = {
            "_scansci_truncated": True,
            "_original_bytes": original_bytes,
            "_notice": "Tool output exceeded the model-context budget; use focused follow-up tools.",
            "summary": _summarize_tool_result(name, safe),
            "preview": encoded[:preview_budget].decode("utf-8", errors="ignore"),
        }
        encoded = json.dumps(compacted, ensure_ascii=False).encode("utf-8")
    return compacted, {
        "original_bytes": original_bytes,
        "model_bytes": len(encoded),
        "truncated": True,
        "persist_full": True,
    }


def _build_self_assessment(history: list[dict[str, Any]]) -> dict[str, Any]:
    if not history:
        return {"steps": 0, "summary": "尚未执行任何工具", "suggestion": "分析请求，选择最合适的工具开始执行"}
    steps: list[dict[str, Any]] = []
    failed = 0
    for entry in history:
        steps.append({
            "tool": entry.get("name", "?"),
            "status": entry.get("status", "?"),
            "summary": str(entry.get("result_summary", entry.get("error", "")))[:200],
        })
        if entry.get("status") == "failed":
            failed += 1
    suggestion = "继续执行"
    if failed > 0:
        suggestion = f"已有 {failed} 个工具失败，考虑尝试不同的方法或参数"
    if len(history) >= 5:
        suggestion += "。已执行多步，请检查结果是否足够交付"
    return {
        "steps_taken": len(history),
        "steps": steps,
        "failed": failed,
        "suggestion": suggestion,
    }


def _json_safe(value: Any) -> Any:
    return json.loads(json.dumps(value, ensure_ascii=False, default=str))
