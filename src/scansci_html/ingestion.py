"""Unified, local-first document ingestion for chat, libraries, and slides.

The ingestion boundary deliberately hides parser-specific objects from the
rest of ScanSci.  MarkItDown is the fast default, Docling is an optional
enhanced parser for difficult documents, and small built-in readers keep the
desktop useful when optional packages are unavailable.
"""

from __future__ import annotations

import base64
import csv
from datetime import datetime, timezone
import json
from pathlib import Path
import re
import shutil
from typing import Any, Iterable
from uuid import uuid4

from bs4 import BeautifulSoup
from docx import Document
from pypdf import PdfReader
from pptx import Presentation


INGESTION_SCHEMA_VERSION = 1
MAX_INGESTION_FILES = 8
MAX_INGESTION_FILE_BYTES = 50 * 1024 * 1024
MAX_INGESTION_TOTAL_BYTES = 120 * 1024 * 1024
MAX_EXTRACTED_CHARS = 240_000
SUPPORTED_INGESTION_SUFFIXES = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".xls",
    ".csv",
    ".json",
    ".xml",
    ".html",
    ".htm",
    ".md",
    ".markdown",
    ".txt",
    ".rtf",
    ".epub",
    ".zip",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".gif",
    ".wav",
    ".mp3",
    ".m4a",
    ".flac",
}


def ingest_sources(
    workspace: str | Path,
    values: object,
    *,
    parser: str = "auto",
    max_files: int = MAX_INGESTION_FILES,
    max_total_bytes: int = MAX_INGESTION_TOTAL_BYTES,
    max_file_bytes: int = MAX_INGESTION_FILE_BYTES,
) -> dict[str, Any]:
    """Persist and extract one-off attachments without requiring a library."""

    normalized_parser = str(parser or "auto").strip().lower()
    if normalized_parser not in {"auto", "fast", "enhanced"}:
        raise ValueError("parser must be auto, fast, or enhanced")
    job = persist_ingestion_sources(
        workspace,
        values,
        max_files=max_files,
        max_total_bytes=max_total_bytes,
        max_file_bytes=max_file_bytes,
    )
    manifest_path = _manifest_path(workspace, str(job["job_id"]))
    records: list[dict[str, Any]] = []
    try:
        for source in list(job["sources"]):
            records.append(_extract_source(source, parser=normalized_parser))
    except Exception as error:
        job["status"] = "failed"
        job["error"] = str(error)
        job["updated_at"] = _utc_now()
        _write_manifest(manifest_path, job)
        raise

    job["status"] = "completed"
    job["sources"] = records
    job["updated_at"] = _utc_now()
    job["summary"] = {
        "files": len(records),
        "characters": sum(int(item.get("character_count", 0)) for item in records),
        "pages": sum(int(item.get("page_count", 0)) for item in records),
        "parsers": sorted({str(item.get("parser", "")) for item in records if item.get("parser")}),
    }
    _write_manifest(manifest_path, job)
    return public_ingestion_job(job)


def extract_local_document(
    path: str | Path,
    *,
    output_dir: str | Path,
    parser: str = "auto",
) -> dict[str, Any]:
    """Extract a referenced local file without copying the original document."""

    source_path = Path(path).expanduser().resolve()
    if not source_path.is_file():
        raise FileNotFoundError(f"文件不存在：{source_path}")
    suffix = source_path.suffix.lower()
    if suffix not in SUPPORTED_INGESTION_SUFFIXES:
        raise ValueError(f"暂不支持 {suffix or '无扩展名'} 文件")
    normalized_parser = str(parser or "auto").strip().lower()
    if normalized_parser not in {"auto", "fast", "enhanced"}:
        raise ValueError("parser must be auto, fast, or enhanced")
    target = Path(output_dir).expanduser().resolve()
    target.mkdir(parents=True, exist_ok=True)
    source = {
        "source_id": f"source-{uuid4().hex[:12]}",
        "name": source_path.name,
        "path": str(source_path),
        "suffix": suffix,
        "size": source_path.stat().st_size,
        "status": "persisted",
    }
    return _extract_source(source, parser=normalized_parser, output_dir=target)


def persist_ingestion_sources(
    workspace: str | Path,
    values: object,
    *,
    max_files: int = MAX_INGESTION_FILES,
    max_total_bytes: int = MAX_INGESTION_TOTAL_BYTES,
    max_file_bytes: int = MAX_INGESTION_FILE_BYTES,
) -> dict[str, Any]:
    """Copy local paths or data URLs into a private, durable ingestion job."""

    supplied = list(values or []) if isinstance(values, list) else []
    if not supplied:
        raise ValueError("请至少添加一个文件")
    allowed_files = max(1, min(_safe_positive_int(max_files, MAX_INGESTION_FILES), 2_000))
    allowed_total = max(
        MAX_INGESTION_FILE_BYTES,
        min(_safe_positive_int(max_total_bytes, MAX_INGESTION_TOTAL_BYTES), 2 * 1024 * 1024 * 1024),
    )
    if len(supplied) > allowed_files:
        raise ValueError(f"一次最多添加 {allowed_files} 个文件")

    job_id = f"ing-{uuid4().hex}"
    root = _ingestion_root(workspace) / job_id
    root.mkdir(parents=True, exist_ok=False)
    total_bytes = 0
    sources: list[dict[str, Any]] = []
    for index, raw in enumerate(supplied, start=1):
        if not isinstance(raw, dict):
            raise ValueError("附件记录格式无效")
        supplied_name = _safe_name(str(raw.get("name", "")))
        source_path = str(raw.get("path", "")).strip()
        data_url = str(raw.get("data_url", "")).strip()
        if source_path:
            candidate = Path(source_path).expanduser().resolve(strict=True)
            if not candidate.is_file():
                raise ValueError("选择的附件不是文件")
            supplied_name = supplied_name or _safe_name(candidate.name)
            suffix = candidate.suffix.lower()
            size = candidate.stat().st_size
            target = root / f"source-{index:02d}{suffix}"
            _validate_source(supplied_name, suffix, size, total_bytes, max_total_bytes=allowed_total, max_file_bytes=max_file_bytes)
            shutil.copy2(candidate, target)
        elif data_url:
            suffix = Path(supplied_name).suffix.lower()
            payload = _decode_data_url(data_url)
            size = len(payload)
            _validate_source(supplied_name, suffix, size, total_bytes, max_total_bytes=allowed_total, max_file_bytes=max_file_bytes)
            target = root / f"source-{index:02d}{suffix}"
            target.write_bytes(payload)
        else:
            raise ValueError("附件必须包含本地路径或上传数据")
        total_bytes += size
        sources.append(
            {
                "source_id": f"source-{index:02d}",
                "name": supplied_name,
                "suffix": suffix,
                "size": size,
                "path": str(target),
                "status": "persisted",
            }
        )

    now = _utc_now()
    job = {
        "schema_version": INGESTION_SCHEMA_VERSION,
        "job_id": job_id,
        "status": "processing",
        "created_at": now,
        "updated_at": now,
        "sources": sources,
        "summary": {"files": len(sources), "characters": 0, "pages": 0, "parsers": []},
    }
    _write_manifest(root / "manifest.json", job)
    return job


def load_ingestion_job(workspace: str | Path, job_id: str) -> dict[str, Any]:
    path = _manifest_path(workspace, job_id)
    if not path.is_file():
        raise FileNotFoundError(f"Ingestion job does not exist: {job_id}")
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as error:
        raise RuntimeError("无法读取附件处理记录") from error
    if not isinstance(payload, dict):
        raise RuntimeError("附件处理记录格式无效")
    return payload


def public_ingestion_job(job: dict[str, Any]) -> dict[str, Any]:
    """Return a browser-safe job without private filesystem paths or full text."""

    job_id = str(job.get("job_id", ""))
    sources = []
    for raw in list(job.get("sources", []) or []):
        item = dict(raw)
        source_id = str(item.get("source_id", ""))
        sources.append(
            {
                "source_id": source_id,
                "name": str(item.get("name", "")),
                "suffix": str(item.get("suffix", "")),
                "size": int(item.get("size", 0) or 0),
                "status": str(item.get("status", "")),
                "parser": str(item.get("parser", "")),
                "page_count": int(item.get("page_count", 0) or 0),
                "character_count": int(item.get("character_count", 0) or 0),
                "warnings": list(item.get("warnings", []) or []),
                "preview": str(item.get("preview", "")),
                "file_url": f"/api/ingestions/{job_id}/sources/{source_id}/file",
                "text_url": f"/api/ingestions/{job_id}/sources/{source_id}/text",
            }
        )
    return {
        "schema_version": int(job.get("schema_version", INGESTION_SCHEMA_VERSION)),
        "job_id": job_id,
        "status": str(job.get("status", "")),
        "created_at": str(job.get("created_at", "")),
        "updated_at": str(job.get("updated_at", "")),
        "summary": dict(job.get("summary", {}) or {}),
        "sources": sources,
        **({"error": str(job.get("error", ""))} if job.get("error") else {}),
    }


def ingestion_source_path(workspace: str | Path, job_id: str, source_id: str) -> Path:
    job = load_ingestion_job(workspace, job_id)
    source = _job_source(job, source_id)
    path = Path(str(source.get("path", ""))).resolve()
    root = (_ingestion_root(workspace) / _safe_job_id(job_id)).resolve()
    if root not in path.parents or not path.is_file():
        raise FileNotFoundError("Attachment file is unavailable")
    return path


def ingestion_source_text(workspace: str | Path, job_id: str, source_id: str) -> str:
    job = load_ingestion_job(workspace, job_id)
    source = _job_source(job, source_id)
    text_path = Path(str(source.get("text_path", ""))).resolve()
    root = (_ingestion_root(workspace) / _safe_job_id(job_id)).resolve()
    if root not in text_path.parents or not text_path.is_file():
        raise FileNotFoundError("Extracted attachment text is unavailable")
    return text_path.read_text(encoding="utf-8")


def ingestion_context(workspace: str | Path, job_id: str, *, limit: int = 80_000) -> str:
    """Build a bounded, source-labelled context for one chat request."""

    job = load_ingestion_job(workspace, job_id)
    remaining = max(1_000, int(limit))
    pieces: list[str] = []
    for source in list(job.get("sources", []) or []):
        if remaining <= 0:
            break
        text_path = Path(str(source.get("text_path", "")))
        text = text_path.read_text(encoding="utf-8") if text_path.is_file() else ""
        if not text:
            continue
        excerpt = text[:remaining]
        pieces.append(f"[附件：{source.get('name', source.get('source_id', 'document'))}]\n{excerpt}")
        remaining -= len(excerpt)
    return "\n\n".join(pieces)


def _extract_source(
    source: dict[str, Any],
    *,
    parser: str,
    output_dir: str | Path | None = None,
) -> dict[str, Any]:
    path = Path(str(source["path"])).resolve()
    suffix = path.suffix.lower()
    warnings: list[str] = []
    text = ""
    selected_parser = ""

    if parser == "enhanced":
        text = _docling_text(path)
        if text:
            selected_parser = "docling"
        else:
            warnings.append("增强解析器不可用，已回退到快速解析")

    # For PDFs the fast parser below preserves page boundaries and is both
    # substantially faster and more useful for evidence citations than doing a
    # full MarkItDown conversion first.  Enhanced mode may still opt into
    # Docling, with the same page-aware fallback when it is unavailable.
    if not text and suffix != ".pdf":
        text = _markitdown_text(path)
        if text:
            selected_parser = "markitdown"

    page_sections: list[dict[str, Any]] = []
    page_count = 1
    if suffix == ".pdf":
        page_sections = _pdf_pages(path)
        page_count = len(page_sections)
        paged_text = "\n\n".join(
            f"# 第 {item['page']} 页\n\n{item['text']}" for item in page_sections if item.get("text")
        )
        if paged_text and parser != "enhanced":
            text = paged_text
            selected_parser = "pypdf"
    if not text:
        text, selected_parser = _builtin_text(path)

    cleaned_text = _clean_text(text)
    was_truncated = len(cleaned_text) > MAX_EXTRACTED_CHARS
    text = cleaned_text[:MAX_EXTRACTED_CHARS]
    if not text:
        raise ValueError(f"未能从《{source.get('name', path.name)}》提取可用内容")
    if was_truncated:
        warnings.append("文档较长，当前会话仅载入前部内容")
    extraction_root = Path(output_dir).resolve() if output_dir is not None else path.parent
    extraction_root.mkdir(parents=True, exist_ok=True)
    text_path = extraction_root / f"{source['source_id']}.md"
    text_path.write_text(text, encoding="utf-8")
    pages_path = extraction_root / f"{source['source_id']}.pages.json"
    pages_path.write_text(json.dumps(page_sections, ensure_ascii=False, indent=2), encoding="utf-8")
    return {
        **source,
        "status": "completed",
        "parser": selected_parser or "builtin",
        "page_count": page_count,
        "character_count": len(text),
        "text_path": str(text_path),
        "pages_path": str(pages_path),
        "preview": _clean_text(text)[:360],
        "warnings": warnings,
    }


def _markitdown_text(path: Path) -> str:
    try:
        from markitdown import MarkItDown

        result = MarkItDown(enable_plugins=False).convert(str(path))
        return str(getattr(result, "text_content", "") or "").strip()
    except Exception:  # optional parser boundary
        return ""


def _docling_text(path: Path) -> str:
    try:
        from docling.document_converter import DocumentConverter

        result = DocumentConverter().convert(str(path))
        document = getattr(result, "document", None)
        exporter = getattr(document, "export_to_markdown", None)
        return str(exporter() if callable(exporter) else "").strip()
    except Exception:  # optional, heavyweight parser boundary
        return ""


def _pdf_pages(path: Path) -> list[dict[str, Any]]:
    # PyMuPDF is substantially faster for ordinary text PDFs.  Keep pypdf as
    # the compatibility fallback for installations without fitz or files it
    # cannot open, while preserving the page-aware evidence structure.
    try:
        import fitz

        document = fitz.open(str(path))
        try:
            if document.needs_pass and not document.authenticate(""):
                raise ValueError("鍙楀瘑鐮佷繚鎶ょ殑 PDF 鏆備笉鏀寔")
            return [
                {"page": index, "text": _clean_text(page.get_text("text") or "")}
                for index, page in enumerate(document, start=1)
            ]
        finally:
            document.close()
    except ValueError:
        raise
    except Exception:
        pass
    reader = PdfReader(str(path))
    if reader.is_encrypted and not reader.decrypt(""):
        raise ValueError("受密码保护的 PDF 暂不支持")
    pages: list[dict[str, Any]] = []
    for index, page in enumerate(reader.pages, start=1):
        pages.append({"page": index, "text": _clean_text(page.extract_text() or "")})
    return pages


def _builtin_text(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        document = Document(str(path))
        return "\n\n".join(p.text for p in document.paragraphs if p.text.strip()), "python-docx"
    if suffix == ".pptx":
        presentation = Presentation(str(path))
        slides = []
        for index, slide in enumerate(presentation.slides, start=1):
            texts = [str(shape.text).strip() for shape in slide.shapes if hasattr(shape, "text") and str(shape.text).strip()]
            slides.append(f"## 幻灯片 {index}\n\n" + "\n\n".join(texts))
        return "\n\n".join(slides), "python-pptx"
    if suffix in {".html", ".htm"}:
        html = path.read_text(encoding="utf-8", errors="replace")
        return BeautifulSoup(html, "html.parser").get_text("\n", strip=True), "beautifulsoup"
    if suffix == ".csv":
        with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
            rows = list(csv.reader(handle))
        return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows), "csv"
    if suffix == ".json":
        value = json.loads(path.read_text(encoding="utf-8-sig"))
        return json.dumps(value, ensure_ascii=False, indent=2), "json"
    if suffix == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=True)
        sections = []
        try:
            for sheet in workbook.worksheets:
                lines = []
                for row in sheet.iter_rows(values_only=True):
                    values = ["" if value is None else str(value).strip() for value in row]
                    if any(values):
                        lines.append(" | ".join(values))
                    if len(lines) >= 2_000:
                        break
                sections.append(f"## {sheet.title}\n\n" + "\n".join(lines))
        finally:
            workbook.close()
        return "\n\n".join(sections), "openpyxl"
    if suffix == ".xls":
        import xlrd

        workbook = xlrd.open_workbook(path, on_demand=True)
        sections = []
        try:
            for sheet in workbook.sheets():
                lines = []
                for row_index in range(min(sheet.nrows, 2_000)):
                    values = [str(sheet.cell_value(row_index, column)).strip() for column in range(sheet.ncols)]
                    if any(values):
                        lines.append(" | ".join(values))
                sections.append(f"## {sheet.name}\n\n" + "\n".join(lines))
        finally:
            workbook.release_resources()
        return "\n\n".join(sections), "xlrd"
    if suffix in {".xml", ".md", ".markdown", ".txt", ".rtf"}:
        return path.read_text(encoding="utf-8-sig", errors="replace"), "text"
    return "", ""


def _validate_source(
    name: str,
    suffix: str,
    size: int,
    total_bytes: int,
    *,
    max_total_bytes: int = MAX_INGESTION_TOTAL_BYTES,
    max_file_bytes: int = MAX_INGESTION_FILE_BYTES,
) -> None:
    if not name:
        raise ValueError("附件缺少文件名")
    if suffix not in SUPPORTED_INGESTION_SUFFIXES:
        raise ValueError(f"暂不支持 {suffix or '无扩展名'} 文件")
    if size <= 0:
        raise ValueError("不能上传空文件")
    allowed_file_bytes = max(MAX_INGESTION_FILE_BYTES, min(int(max_file_bytes), 512 * 1024 * 1024))
    if size > allowed_file_bytes:
        raise ValueError(f"单个附件不能超过 {allowed_file_bytes // (1024 * 1024)} MB")
    if total_bytes + size > max_total_bytes:
        raise ValueError(f"本次附件总大小不能超过 {max_total_bytes // (1024 * 1024)} MB")


def _safe_positive_int(value: object, fallback: int) -> int:
    try:
        parsed = int(value)
    except (TypeError, ValueError):
        return fallback
    return parsed if parsed > 0 else fallback


def _decode_data_url(value: str) -> bytes:
    header, separator, encoded = value.partition(",")
    if not separator or ";base64" not in header.lower():
        raise ValueError("上传数据格式无效")
    try:
        return base64.b64decode(encoded, validate=True)
    except ValueError as error:
        raise ValueError("上传文件无法解码") from error


def _ingestion_root(workspace: str | Path) -> Path:
    root = Path(workspace).resolve().parent / ".scansci-ingestions"
    root.mkdir(parents=True, exist_ok=True)
    return root


def _manifest_path(workspace: str | Path, job_id: str) -> Path:
    return _ingestion_root(workspace) / _safe_job_id(job_id) / "manifest.json"


def _safe_job_id(value: str) -> str:
    normalized = str(value or "").strip()
    if not re.fullmatch(r"ing-[a-f0-9]{32}", normalized):
        raise FileNotFoundError("Ingestion job does not exist")
    return normalized


def _job_source(job: dict[str, Any], source_id: str) -> dict[str, Any]:
    normalized = str(source_id or "").strip()
    source = next((dict(item) for item in list(job.get("sources", []) or []) if str(item.get("source_id", "")) == normalized), None)
    if source is None:
        raise FileNotFoundError("Attachment source does not exist")
    return source


def _write_manifest(path: Path, payload: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_suffix(".tmp")
    temporary.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    temporary.replace(path)


def _safe_name(value: str) -> str:
    return Path(str(value or "")).name.strip()[:180]


def _clean_text(value: str) -> str:
    text = str(value or "").replace("\x00", "")
    text = re.sub(r"(?<=[\u3400-\u4dbf\u4e00-\u9fff\uf900-\ufaff])\s+(?=[\u3400-\u4dbf\u4e00-\u9fff\uf900-\ufaff])", "", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _utc_now() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")
