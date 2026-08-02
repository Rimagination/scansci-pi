from __future__ import annotations

import json
from pathlib import Path
import time
import base64
from io import BytesIO
import zipfile

import fitz
from pptx import Presentation
from pptx.util import Inches

from scansci_html.slide_studio import (
    _balanced_cover_lines,
    _display_source_name,
    _ground_model_adaptation_claim,
    _normalise_model_outline,
    _slide_excerpt,
    _source_key_sentences,
    create_source_slide_deck,
    persist_slide_sources,
    save_browser_rendered_deck,
)
from scansci_html.research_runs import StageSpec
from scansci_html.webapp import NotebookWebApp


def _make_pdf(path: Path) -> Path:
    document = fitz.open()
    page = document.new_page()
    page.insert_text(
        (72, 90),
        "Research question: how can evidence-linked analysis improve research decisions?\n"
        "The study compares transparent review steps with unstructured synthesis.\n"
        "Results show that traceable source links reduce unsupported claims.\n"
        "The conclusion is limited to the supplied evaluation setting.",
        fontsize=13,
    )
    document.save(path)
    document.close()
    return path


def _classic_easyslides_fixture(tmp_path: Path) -> Path:
    root = tmp_path / "easyslides"
    layouts = root / "templates" / "layouts"
    template = layouts / "defense_leftnav"
    template.mkdir(parents=True)
    (layouts / "layouts_index.json").write_text(
        json.dumps(
            {
                "defense_leftnav": {
                    "summary": "Classic academic defense template.",
                    "keywords": ["academic", "defense"],
                }
            }
        ),
        encoding="utf-8",
    )
    (template / "design_spec.md").write_text(
        """---
template_id: defense_leftnav
display_name_zh: Defense Left Navigation
category: defense
primary_color: "#8B0012"
canvas_format: ppt169
replication_mode: classic
---
""",
        encoding="utf-8",
    )
    roles = ("cover", "toc", "chapter", "content", "ending")
    svg = '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1280 720"><rect width="1280" height="720" fill="#fff"/></svg>'
    for index, role in enumerate(roles, start=1):
        (template / f"{index:02d}_{role}.svg").write_text(svg, encoding="utf-8")
    (template / "layouts.json").write_text(
        json.dumps(
            {
                "template_id": "defense_leftnav",
                "replication_mode": "classic",
                "shells": [{"role": role} for role in roles],
                "slot_models": {role: [{"id": f"{role}-default"}] for role in roles},
                "text_fit_policy": {"overflow": "shrink"},
            }
        ),
        encoding="utf-8",
    )
    return root


def test_pdf_source_creates_editable_pptx_without_a_library(tmp_path: Path):
    workspace = tmp_path / "workspace.sqlite"
    source = _make_pdf(tmp_path / "study.pdf")

    persisted = persist_slide_sources(workspace, [{"name": source.name, "path": str(source)}])
    output = create_source_slide_deck(workspace=workspace, sources=persisted, topic="Evidence-linked scientific slides")

    deck_path = Path(output["pptx_path"])
    presentation = Presentation(deck_path)
    assert deck_path.is_file()
    assert deck_path.parent == tmp_path / "presentations"
    assert len(presentation.slides) >= 5
    assert len(presentation.slides) == len(output["outline"]["slides"]) + 1 == output["slide_count"]
    assert presentation.slide_width > presentation.slide_height
    assert output["source_count"] == 1
    assert output["planning"]["mode"] == "source-grounded"
    assert output["slide_plan"]["schema"] == "scansci.slide-plan.v1"
    assert Path(output["slide_plan_path"]).is_file()
    assert output["enhanced_renderer"] == "pptxgenjs-browser"
    assert any("Research question" in shape.text for shape in presentation.slides[1].shapes if hasattr(shape, "text"))


def test_selected_classic_template_routes_to_the_native_easyslides_renderer(tmp_path: Path, monkeypatch):
    workspace = tmp_path / "workspace.sqlite"
    slides_root = _classic_easyslides_fixture(tmp_path)
    source = _make_pdf(tmp_path / "study.pdf")
    persisted = persist_slide_sources(workspace, [{"name": source.name, "path": str(source)}])
    calls: list[dict[str, object]] = []

    def fake_classic_renderer(**kwargs):
        calls.append(kwargs)
        deck_path = Path(kwargs["deck_path"])
        deck_path.parent.mkdir(parents=True, exist_ok=True)
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(0.5)).text = "研究背景"
        presentation.save(deck_path)
        plan_path = deck_path.with_suffix(".classic-plan.json")
        plan = {"slides": [{"role": "cover"}]}
        plan_path.write_text(json.dumps(plan), encoding="utf-8")
        return {
            "project_path": str(deck_path.parent / "projects" / "native-classic"),
            "plan": plan,
            "plan_path": str(plan_path),
            "quality_gate": {"status": "pass", "blocking_count": 0},
            "classic_text_fit": {"status": "pass"},
            "visual_measure": {"status": "pass"},
            "visual_review": {"status": "needs_review"},
            "plugin": {"version": "fixture"},
            "render_manifest": {"renderer": "easyslides-classic-shells"},
        }

    monkeypatch.setattr("scansci_html.slide_studio.render_classic_deck", fake_classic_renderer)

    output = create_source_slide_deck(
        workspace=workspace,
        sources=persisted,
        topic="Template-aware scientific briefing",
        template_id="defense_leftnav",
        slides_root=slides_root,
    )
    presentation = Presentation(output["pptx_path"])
    content_shapes = presentation.slides[0].shapes

    assert output["template_id"] == "defense_leftnav"
    assert output["template"]["primary_color"] == "#8B0012"
    assert output["renderer"] == "easyslides-classic-native"
    assert output["quality_gate"]["status"] == "pass"
    assert calls and calls[0]["template"]["id"] == "defense_leftnav"
    assert any("研究背景" in shape.text for shape in content_shapes if hasattr(shape, "text"))


def test_browser_rendered_pptx_is_validated_and_saved(tmp_path: Path):
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("[Content_Types].xml", "<Types/>")
        archive.writestr("ppt/presentation.xml", "<presentation/>")

    result = save_browser_rendered_deck(
        tmp_path / "workspace.sqlite",
        file_name="Evidence review.pptx",
        base64_data=base64.b64encode(buffer.getvalue()).decode("ascii"),
    )

    assert result["renderer"] == "pptxgenjs-browser"
    assert Path(result["file_path"]).read_bytes().startswith(b"PK")
    assert result["download_url"].startswith("/api/presentations/")


def test_web_workflow_accepts_pdf_and_serves_pptx_without_a_notebook(tmp_path: Path, monkeypatch):
    workspace = tmp_path / "workspace.sqlite"
    app = NotebookWebApp(workspace=workspace, evidence_db=tmp_path / "evidence.sqlite")
    source = _make_pdf(tmp_path / "source.pdf")
    monkeypatch.setattr(app.research_agent, "_slides_chat_client", lambda: None)

    created = json.loads(
        app.dispatch(
            "POST",
            "/api/runs",
            json.dumps(
                {
                    "workflow_type": "pdf_to_ppt",
                    "topic": "PDF to slide deck",
                    "source_files": [{"name": source.name, "path": str(source)}],
                }
            ).encode("utf-8"),
        ).body.decode("utf-8")
    )
    completed = created
    for _ in range(160):
        completed = json.loads(app.dispatch("GET", f"/api/runs/{created['run_id']}").body.decode("utf-8"))
        if completed["status"] in {"completed", "failed"}:
            break
        time.sleep(0.025)

    assert completed["status"] == "completed"
    assert completed["notebook_id"] == ""
    assert completed["output_artifact"]["artifact_type"] == "presentation_deck"
    assert Path(completed["output_artifact"]["file_path"]).is_file()
    download = app.dispatch("GET", f"/api/runs/{created['run_id']}/download")
    assert download.status == 200
    assert download.content_type.endswith("presentationml.presentation")
    assert download.body.startswith(b"PK")
    preview = app.dispatch("GET", f"/api/runs/{created['run_id']}/preview")
    assert preview.status == 200
    assert preview.content_type.startswith("image/svg+xml")
    assert b"ScanSci Presentation Studio" in preview.body


def test_presentation_ui_uses_stable_rendering_and_native_save_bridge():
    app_js = (Path(__file__).parents[1] / "src" / "scansci_html" / "web" / "app.js").read_text(encoding="utf-8")

    assert 'pdf_to_ppt: "none"' in app_js
    assert "const shouldFollow = distanceFromBottom < 72;" in app_js
    assert "continueTaskConversation" in app_js
    assert "}/messages`" in app_js
    assert "isTaskFollowUp" in app_js
    assert "isPptRecreate" not in app_js
    assert "originalSlideSources" not in app_js
    assert "taskConversationMarkup" in app_js
    assert "const isTaskConversation = inputId === \"chatQuestionInput\"" in app_js
    assert "const isTaskFollowUp = isTaskConversation" in app_js
    assert "renderAssistantContent(message.content)" in app_js
    assert '<a href=\"$2\" target=\"_blank\" rel=\"noopener noreferrer\">$1</a>' in app_js
    assert 'input.dispatchEvent(new Event("input", { bubbles: true }));' in app_js
    assert 'mode === "knowledge" && !selectedKnowledge.length && !isTaskFollowUp' in app_js
    # Public academic/deep research no longer requires a selected local
    # notebook.  Only the two evidence-bound shortcuts retain that guard.
    assert '["novelty", "idea"].includes(mode) && !state.notebook && !isTaskFollowUp' in app_js
    assert 'skills: extractSkillMentions(content)' in app_js
    assert 'document.addEventListener("keyup", (event) =>' in app_js
    assert 'renderSkillSuggestions(event.target);' in app_js
    assert '!["ArrowDown", "ArrowUp", "Enter", "Escape"].includes(event.key)' in app_js
    assert "partialDownloadArtifactMarkup" in app_js
    assert "继续下载并交付" in app_js
    assert "taskOriginalPrompt" not in app_js
    assert "原始任务" not in app_js
    assert "任务续聊" not in app_js
    assert 'aria-label="任务消息"' in app_js
    assert 'run.error?.code === "app_restarted"' in app_js
    assert '`/api/runs/${encodeURIComponent(id)}/resume`' in app_js
    assert '["queued", "planning", "running", "verifying"].includes(String(run.status || ""))' in app_js
    assert 'window.localStorage.setItem("scansci.active.task", run.run_id)' in app_js
    assert "任务上下文 · 续聊" not in app_js


def test_presentation_follow_up_route_targets_the_existing_run(tmp_path: Path, monkeypatch):
    app_js = (Path(__file__).parents[1] / "src" / "scansci_html" / "web" / "app.js").read_text(encoding="utf-8")
    app = NotebookWebApp(workspace=tmp_path / "workspace.sqlite", evidence_db=tmp_path / "evidence.sqlite")
    run = app.research_agent.store.create_run(
        notebook_id="",
        workflow_type="pdf_to_ppt",
        title="Existing deck",
        input_payload={"question": "Create a deck"},
        stages=[StageSpec("deliver", "Deliver", "delivery")],
    )
    captured: dict[str, object] = {}

    def continue_run(run_id: str, payload: dict[str, object]) -> dict[str, object]:
        captured.update({"run_id": run_id, "payload": payload})
        return {"run": app.research_agent.store.get_run(run_id), "message": {"role": "assistant", "content": "Kept together."}}

    monkeypatch.setattr(app.research_agent, "continue_run_conversation", continue_run)
    response = app.dispatch(
        "POST",
        f"/api/runs/{run['run_id']}/messages",
        json.dumps({"content": "What does slide 3 mean?"}).encode("utf-8"),
    )

    assert response.status == 200
    assert captured == {"run_id": run["run_id"], "payload": {"content": "What does slide 3 mean?"}}
    assert "data-action=\"save-presentation\"" in app_js
    assert "重新排版导出" in app_js


def test_slide_planner_preserves_a_long_research_question():
    question = "Can a physics-informed, data-efficient model accurately predict the long-term degradation trajectory of retired batteries using only partial field-accessible signals without historical records?"
    candidate = {
        "title": "Long-title deck",
        "central_question": question,
        "story": "Question to evidence.",
        "slides": [
            {"title": "Question", "takeaway": "Frame it.", "layout": "comparison", "bullets": ["Evidence."], "source_pages": [1]},
            {"title": "Evidence", "takeaway": "Show it.", "layout": "process", "bullets": ["Evidence."], "source_pages": [1]},
            {"title": "Boundary", "takeaway": "Limit it.", "layout": "branches", "bullets": ["Evidence."], "source_pages": [1]},
        ],
    }
    fallback = {"title": "Fallback", "central_question": "Fallback", "story": "Fallback", "slides": []}

    outline = _normalise_model_outline(candidate, fallback=fallback, sources=[{"page_count": 1}])

    assert outline["central_question"] == question
    assert [slide["layout"] for slide in outline["slides"]] == ["comparison", "process", "branches"]


def test_slide_planner_rebalances_a_repetitive_model_layout():
    titles = [
        "背景：Transformer 架构",
        "BERT：双向编码",
        "GPT-3：自回归生成",
        "训练范式对比",
        "实证结果",
        "规模与效率",
        "三条路线的总结与边界",
    ]
    candidate = {
        "title": "Transformer 到 BERT 与 GPT-3",
        "central_question": "三条路线如何比较？",
        "story": "从基础到分支。",
        "slides": [
            {
                "title": title,
                "takeaway": "材料支持的结论。",
                "layout": "comparison",
                "bullets": ["第一项", "第二项", "第三项", "第四项"],
                "source_pages": [1],
            }
            for title in titles
        ],
    }

    outline = _normalise_model_outline(
        candidate,
        fallback={"title": "Fallback", "central_question": "Fallback", "story": "Fallback", "slides": []},
        sources=[{"page_count": 3}],
    )

    assert [slide["layout"] for slide in outline["slides"]] == [
        "process",
        "cards",
        "process",
        "comparison",
        "cards",
        "process",
        "branches",
    ]


def test_cover_line_balancing_avoids_a_single_word_orphan():
    title = "Transformer 到 BERT 与 GPT-3：架构、训练范式与实证边界"

    lines = _balanced_cover_lines(title).splitlines()

    assert len(lines) == 2
    assert min(len(line) for line in lines) >= 12


def test_slide_claims_do_not_conflate_output_layers_or_gpt3_context_learning():
    corpus = (
        "During fine-tuning, all parameters are fine-tuned. "
        "We evaluate GPT-3 without any gradient updates or fine-tuning."
    ).casefold()

    assert _ground_model_adaptation_claim("BERT 仅需微调一层输出层。", corpus) == (
        "BERT 在下游任务中新增任务输出层，并联合微调全部预训练参数。"
    )
    assert _ground_model_adaptation_claim("GPT-3 在部分任务仍需微调才能达到最优。", corpus) == (
        "GPT-3 论文的零样本、单样本与少样本评估通过上下文示例适配，不进行梯度更新或微调。"
    )


def test_source_names_are_presentable_on_the_cover():
    assert _display_source_name("Attention_Is_All_You_Need.pdf") == "Attention Is All You Need"


def test_offline_slide_fallback_filters_front_matter_and_author_contributions():
    source = {
        "name": "Attention_Is_All_You_Need.pdf",
        "text": (
            "Provided proper attribution is provided, Google hereby grants permission. "
            "Ashish Vaswani Google Brain author@example.com. "
            "The Transformer uses multi-head attention in its encoder and decoder. "
            "Our model achieves 28.4 BLEU on the WMT 2014 translation task."
        ),
    }

    sentences = _source_key_sentences(source, topic="Transformer architecture")

    assert any("multi-head attention" in item for item in sentences)
    assert all("@" not in item and "proper attribution" not in item for item in sentences)


def test_offline_slide_fallback_filters_install_and_preview_commands():
    source = {
        "name": "README.md",
        "text": (
            "# ScanSci\n"
            "ScanSci is an evidence-first research workspace for traceable scientific analysis.\n"
            "Browser preview always serves this checkout at http://127.0.0.1:8781.\n"
            "Run npm install and python -m scansci_html to start the application.\n"
            "Windows 构建会将 Pi sidecar、Node.js runtime 与本地检索运行时放入应用目录。\n"
            "It organizes source evidence, review boundaries, and editable research deliverables."
        ),
    }

    sentences = _source_key_sentences(source, topic="ScanSci research workflow")

    assert any("evidence-first" in item for item in sentences)
    assert all(
        "127.0.0.1" not in item and "npm install" not in item and "Windows 构建" not in item
        for item in sentences
    )


def test_offline_slide_fallback_collapses_a_markdown_heading_repeated_in_extracted_text():
    source = {
        "name": "ScanSci README.md",
        "text": "# ScanSci ScanSci 是一款 evidence-first 科研工作台，支持可追溯证据与可编辑交付物。",
    }

    sentences = _source_key_sentences(source, topic="ScanSci research workflow")

    assert sentences[0].startswith("ScanSci 是一款")
    assert "ScanSci ScanSci" not in sentences[0]


def test_slide_excerpt_marks_a_visual_truncation_instead_of_returning_a_half_sentence():
    value = "This deliberately long sentence has no early punctuation but still needs a visible and honest truncation marker"

    excerpt = _slide_excerpt(value, limit=42)

    assert excerpt.endswith("…")
    assert len(excerpt) <= 43
